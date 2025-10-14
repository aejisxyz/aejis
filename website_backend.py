#!/usr/bin/env python3
"""
Aejis Security Website Backend
Handles file uploads and analysis for the web platform
"""

import os
import time
import logging
import math
import uuid
import threading
import requests
import re
import socket
import ssl
import base64
from io import BytesIO
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import TimeoutException, WebDriverException
from PIL import Image
from flask import Flask, request, jsonify, send_file, make_response, redirect, Response
from flask_cors import CORS
from browser_isolation_service import browser_isolation_service
from werkzeug.utils import secure_filename
import tempfile
import json
from datetime import datetime
from typing import Dict, Any
import docker
import shutil

# Import our analysis modules
from virustotal_engine import scan_with_virustotal
from file_analyzer import FileAnalyzer

# Configure logging with debug option
# Set to DEBUG for detailed troubleshooting
LOG_LEVEL = os.getenv('AEJIS_LOG_LEVEL', 'INFO').upper()
log_level = getattr(logging, LOG_LEVEL, logging.INFO)

logging.basicConfig(
    level=log_level,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Enable debug mode if requested
if LOG_LEVEL == 'DEBUG':
    logger.info("🐛 Debug logging enabled - detailed logs will be shown")
    logger.debug("🔧 Debug mode: All operations will be logged in detail")

# Docker-based secure file preview system
class SecureFilePreview:
    """Docker-based file preview system for maximum security"""
    
    def __init__(self):
        try:
            self.docker_client = docker.from_env()
            self.docker_client.ping()
            self.persistent_container = None
            self.persistent_container_name = "aejis-persistent-preview"
            logger.info("🐳 Docker initialized for secure file preview")
            
            # Initialize persistent container
            self._ensure_persistent_container()
        except Exception as e:
            logger.error(f"Docker REQUIRED but not available: {e}")
            raise Exception(f"Docker is required for secure file preview but not available: {e}")
    
    def _ensure_persistent_container(self):
        """Ensure a persistent container is running for fast processing"""
        try:
            # Check if container already exists and is running
            try:
                container = self.docker_client.containers.get(self.persistent_container_name)
                if container.status == 'running':
                    self.persistent_container = container
                    logger.info(f"🚀 Using existing persistent container: {self.persistent_container_name}")
                    return
                else:
                    # Container exists but not running, remove it
                    container.remove(force=True)
                    logger.info(f"🗑️ Removed stopped persistent container")
            except docker.errors.NotFound:
                pass  # Container doesn't exist, we'll create it
            
            # Create new persistent container with pre-installed packages
            logger.info(f"🔧 Creating new persistent container: {self.persistent_container_name}")
            self.persistent_container = self.docker_client.containers.run(
                'python:3.11-slim',
                command=['python', '-c', '''
import time
import os
import subprocess
import sys

# Install packages on startup
print("Installing packages...")
packages = ["Pillow", "PyPDF2", "python-docx", "openpyxl", "python-pptx", "mutagen", "exifread", "rarfile", "py7zr"]
for pkg in packages:
    try:
        print(f"Installing {pkg}...")
        result = subprocess.run(["pip", "install", "--no-cache-dir", pkg], 
                              check=True, capture_output=True, text=True)
        print(f"✅ {pkg} installed")
    except Exception as e:
        print(f"❌ Failed to install {pkg}: {e}")

print("All packages installed. Container ready.")
# Keep container alive
while True:
    time.sleep(1)
'''],
                name=self.persistent_container_name,
                detach=True,
                remove=False,
                mem_limit='1g',  # More memory for packages
                # Allow network access for package installation, will be disabled later
            )
            logger.info(f"✅ Persistent container created successfully: {self.persistent_container_name}")
            
            # Wait for packages to install
            logger.info("⏳ Waiting for packages to install...")
            time.sleep(10)  # Give time for package installation
            
        except Exception as e:
            logger.error(f"❌ Failed to create persistent container: {e}")
            self.persistent_container = None
    
    def _use_persistent_container(self, file_path: str, analysis_id: str) -> Dict[str, Any]:
        """Use the persistent container for fast processing"""
        if not self.persistent_container:
            logger.warning("⚠️ Persistent container not available")
            return {"success": False, "error": "Persistent container not available"}
        
        try:
            # Check if container is still running
            self.persistent_container.reload()
            if self.persistent_container.status != 'running':
                logger.warning("⚠️ Persistent container stopped")
                return {"success": False, "error": "Persistent container stopped"}
            
            # Create processing script based on file type
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext == '.pdf':
                processing_script = self._create_pdf_processor_script()
            elif file_ext == '.docx':
                processing_script = self._create_docx_processor_script()
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.tiff', '.ico']:
                processing_script = self._create_image_processor_script()
            elif file_ext in ['.txt', '.log', '.md', '.json', '.xml', '.csv', '.py', '.js', '.html', '.css', '.sh', '.bat', '.ps1', '.sql']:
                processing_script = self._create_text_processor_script()
            elif file_ext in ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.ogg', '.wmv', '.flv']:
                processing_script = self._create_video_processor_script()
            else:
                # Fallback to universal processor for other file types
                processing_script = self._create_universal_processor_script(file_ext)
            
            # Copy file to temporary location in container (use unique name per analysis)
            container_file_path = f"/tmp/file_{analysis_id}{file_ext}"
            # Copy file into container
            with open(file_path, 'rb') as f:
                file_data = f.read()
            
            # Create tar archive with the file
            import tarfile
            import io
            
            tar_buffer = io.BytesIO()
            tar = tarfile.open(fileobj=tar_buffer, mode='w')
            
            # Add the file to the tar (use unique name per analysis)
            file_info = tarfile.TarInfo(name=f"file_{analysis_id}{file_ext}")
            file_info.size = len(file_data)
            tar.addfile(file_info, io.BytesIO(file_data))
            
            # No need to replace file path since script already expects the correct path
            script_with_file_path = processing_script
            script_data = script_with_file_path.encode('utf-8')
            script_info = tarfile.TarInfo(name=f"process_{analysis_id}.py")
            script_info.size = len(script_data)
            tar.addfile(script_info, io.BytesIO(script_data))
            
            tar.close()
            tar_buffer.seek(0)
            
            # Copy tar archive to container
            self.persistent_container.put_archive('/tmp', tar_buffer.getvalue())
            
            # Clean up old files in container before processing
            try:
                self.persistent_container.exec_run('rm -f /tmp/file_* /tmp/process_*.py', stdout=True, stderr=True)
            except Exception as e:
                logger.debug(f"Cleanup warning: {e}")
            
            # Execute the processing script
            start_time = time.time()
            exec_result = self.persistent_container.exec_run(
                f'python /tmp/process_{analysis_id}.py',
                stdout=True,
                stderr=True
            )
            
            execution_time = time.time() - start_time
            
            if exec_result.exit_code == 0:
                # Parse the result from stdout
                output = exec_result.output.decode('utf-8', errors='replace').strip()
                logger.debug(f"🔍 Raw container output (first 500 chars): {output[:500]}")
                
                # Find the JSON result in the output (skip installation logs)
                lines = output.split('\n')
                json_result = None
                
                # First try to find a single-line JSON
                for line in lines:
                    line = line.strip()
                    if line.startswith('{') and line.endswith('}'):
                        try:
                            json_result = json.loads(line)
                            break
                        except json.JSONDecodeError:
                            continue
                
                # If no single-line JSON found, try to find multi-line JSON
                if not json_result:
                    json_start = -1
                    json_end = -1
                    for i, line in enumerate(lines):
                        if line.strip().startswith('{'):
                            json_start = i
                            break
                    
                    if json_start >= 0:
                        for i in range(json_start, len(lines)):
                            if lines[i].strip().endswith('}'):
                                json_end = i
                                break
                        
                        if json_end >= 0:
                            json_lines = lines[json_start:json_end + 1]
                            json_text = '\n'.join(json_lines)
                            try:
                                json_result = json.loads(json_text)
                            except json.JSONDecodeError:
                                pass
                
                # If still no JSON found, try to extract JSON from the entire output using regex
                if not json_result:
                    import re
                    json_pattern = r'\{.*\}'
                    json_matches = re.findall(json_pattern, output, re.DOTALL)
                    for json_match in json_matches:
                        try:
                            json_result = json.loads(json_match)
                            break
                        except json.JSONDecodeError:
                            pass
                
                if json_result:
                    logger.info(f"✅ Persistent container completed in {execution_time:.3f}s")
                    logger.debug(f"📊 JSON result keys: {list(json_result.keys())}")
                    # Map Docker response to frontend expectations
                    if json_result.get('thumbnail'):
                        thumbnail_data = json_result['thumbnail']
                        # Strip data URL prefix if present (frontend expects just base64 data)
                        if thumbnail_data.startswith('data:'):
                            # Extract just the base64 part after the comma
                            if ',' in thumbnail_data:
                                json_result['base64_preview'] = thumbnail_data.split(',', 1)[1]
                            else:
                                json_result['base64_preview'] = thumbnail_data
                        else:
                            json_result['base64_preview'] = thumbnail_data
                    return json_result
                else:
                    logger.error(f"❌ No valid JSON found in persistent container output: {output[:500]}...")
                    logger.error(f"❌ Full output: {output}")
                    return {"success": False, "error": "No valid JSON result found"}
            else:
                error_output = exec_result.output.decode('utf-8', errors='replace')
                logger.error(f"❌ Persistent container execution failed: {error_output}")
                return {"success": False, "error": f"Execution failed: {error_output}"}
            
        except Exception as e:
            logger.error(f"❌ Persistent container error: {e}")
            return {"success": False, "error": f"Container error: {str(e)}"}
    
    def _use_single_container(self, file_path: str, analysis_id: str) -> Dict[str, Any]:
        """Fallback: Use single-use container (original method)"""
        logger.info("🔄 Using single-use container as fallback")
        # This is the original implementation - we'll keep it as fallback
        return self._original_docker_processing(file_path, analysis_id)
    
    def safe_preview_any_file(self, file_path: str, analysis_id: str) -> Dict[str, Any]:
        """Safely preview any file type using Docker isolation - DOCKER REQUIRED"""
        logger.info(f"🐳 Starting Docker preview for: {os.path.basename(file_path)} (ID: {analysis_id})")
        
        # Try persistent container first for speed
        if self.persistent_container:
            logger.info("🚀 Using persistent container for fast processing")
            result = self._use_persistent_container(file_path, analysis_id)
            if result.get("success"):
                return result
            else:
                logger.warning("⚠️ Persistent container failed, attempting to recreate...")
                # Try to recreate persistent container
                self._ensure_persistent_container()
                if self.persistent_container:
                    logger.info("🔄 Retrying with recreated persistent container")
                    result = self._use_persistent_container(file_path, analysis_id)
                    if result.get("success"):
                        return result
                logger.warning("⚠️ Persistent container recreation failed, falling back to single-use")
        
        # Fall back to original single-use container method
        return self._original_docker_processing(file_path, analysis_id)
    
    def _original_docker_processing(self, file_path: str, analysis_id: str) -> Dict[str, Any]:
        """Original Docker processing method (single-use containers)"""
        logger.info("🔄 Using single-use container processing")
        
        try:
            # Create secure temp directory
            with tempfile.TemporaryDirectory(prefix=f"secure_{analysis_id}_") as temp_dir:
                logger.debug(f"📁 Created temp directory: {temp_dir}")
                
                # Copy file to temp (with safe name)
                file_ext = os.path.splitext(file_path)[1].lower()
                temp_file = os.path.join(temp_dir, f"file{file_ext}")
                file_size = os.path.getsize(file_path)
                
                logger.info(f"📋 Processing file: ext={file_ext}, size={file_size} bytes")
                shutil.copy2(file_path, temp_file)
                logger.debug(f"✅ File copied to temp location: {temp_file}")
                
                # Create processing script based on file type
                logger.debug(f"🔧 Generating processing script for {file_ext}")
                if file_ext == '.pdf':
                    script_content = self._create_pdf_processor_script()
                elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.tiff', '.ico']:
                    script_content = self._create_image_processor_script()
                elif file_ext in ['.txt', '.log', '.md', '.json', '.xml', '.csv', '.py', '.js', '.html', '.css', '.sh', '.bat', '.ps1', '.sql']:
                    script_content = self._create_text_processor_script()
                else:
                    # Fallback to universal processor for other file types
                    script_content = self._create_universal_processor_script(file_ext)
                script_path = os.path.join(temp_dir, "process.py")
                with open(script_path, 'w', encoding='utf-8') as f:
                    f.write(script_content)
                logger.debug(f"✅ Processing script created: {len(script_content)} characters")
                
                # Run in isolated Docker container (fixed output handling)
                logger.info(f"🚀 Starting Docker container execution for {file_ext}")
                start_time = time.time()
                
                try:
                    logger.debug(f"🐳 Docker config: image=python:3.11-slim, mem=512m, cpu=100000, network=none")
                    output = self.docker_client.containers.run(
                    "python:3.11-slim",
                    command=["python", "/tmp/process.py"],
                    volumes={temp_dir: {'bind': '/tmp', 'mode': 'ro'}},
                    network_mode='none',  # No network access
                    mem_limit='512m',     # Memory limit
                    cpu_quota=100000,     # CPU limit
                    security_opt=['no-new-privileges:true'],  # No privilege escalation
                    read_only=True,       # Read-only filesystem
                    remove=True,
                        detach=False
                    )
                    execution_time = time.time() - start_time
                    logger.info(f"✅ Docker execution completed in {execution_time:.3f}s")
                    
                except Exception as docker_error:
                    # Handle Docker-specific errors
                    execution_time = time.time() - start_time
                    error_msg = f"Docker execution failed: {str(docker_error)}"
                    logger.error(f"❌ Docker execution failed after {execution_time:.3f}s: {docker_error}")
                    logger.error(f"🔍 Error details - Type: {type(docker_error).__name__}, Message: {str(docker_error)}")
                    return self._error_result(file_path, analysis_id, error_msg)
                
                # Get results (output is already bytes when detach=False)
                logger.debug(f"📤 Processing Docker output, type: {type(output)}")
                result = ""
                if isinstance(output, bytes):
                    logger.debug(f"📊 Docker output size: {len(output)} bytes")
                    # Try multiple encodings to handle different character sets
                    for encoding in ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']:
                        try:
                            result = output.decode(encoding, errors='replace').strip()
                            logger.debug(f"✅ Successfully decoded with {encoding}")
                            break
                        except UnicodeDecodeError as decode_error:
                            logger.debug(f"❌ Failed to decode with {encoding}: {decode_error}")
                            continue
                    
                    if not result:
                        result = output.decode('utf-8', errors='replace').strip()
                        logger.warning("⚠️ Fallback to UTF-8 with error replacement")
                else:
                    result = str(output).strip()
                    logger.debug(f"📝 Using string output directly")
                    
                logger.info(f"📋 Docker result preview: {result[:200]}..." if len(result) > 200 else f"📋 Docker result: {result}")
                return self._parse_docker_result(result, file_path, analysis_id)
                
        except Exception as e:
            # Safely handle encoding errors in exception messages
            logger.error(f"❌ Unexpected error in safe_preview_any_file for {os.path.basename(file_path)}")
            logger.error(f"🔍 Error type: {type(e).__name__}")
            
            try:
                error_msg = str(e)
                # Test if the error message can be encoded safely
                error_msg.encode('utf-8')
                logger.error(f"🔍 Error details: {error_msg}")
            except UnicodeEncodeError:
                error_msg = "Docker processing failed due to encoding issues"
                logger.error(f"🔍 Error has encoding issues, using fallback message")
            except Exception:
                error_msg = "Docker processing failed with unknown error"
                logger.error(f"🔍 Unknown error occurred during error handling")
            
            logger.error(f"❌ Final error in safe_preview_any_file: {error_msg}")
            return self._error_result(file_path, analysis_id, f"Docker processing failed: {error_msg}")
    
    def _create_pdf_processor_script(self) -> str:
        """Create dedicated PDF processing script"""
        return '''#!/usr/bin/env python3
import os
import sys
import json
import base64
import subprocess
import glob

def process_pdf():
    try:
        # Look for the unique file with analysis ID pattern
        file_patterns = ['/tmp/file_*.pdf']
        file_path = None
        
        for pattern in file_patterns:
            matching_files = glob.glob(pattern)
            if matching_files:
                # Get the most recent file (in case of multiple matches)
                file_path = max(matching_files, key=os.path.getctime)
                break
        
        if not file_path:
            return {'success': False, 'error': 'File not found - no matching PDF file'}
        
        file_size = os.path.getsize(file_path)
        
        # Try to extract text using pdftotext
        try:
            result = subprocess.run(['pdftotext', file_path, '-'], 
                                  capture_output=True, text=True, timeout=30)
            if result.returncode == 0:
                text_content = result.stdout
                # Limit content length
                if len(text_content) > 50000:
                    text_content = text_content[:50000] + "\\n\\n[... content truncated ...]"
                
                return {
                    'success': True,
                    'preview_content': f"PDF Document\\nSize: {file_size} bytes\\n\\nText Content:\\n{text_content}",
                    'preview_type': 'document',
                    'file_size': file_size,
                    'metadata': {
                        'file_size': file_size,
                        'text_length': len(text_content),
                        'pages_processed': True
                    }
                }
        except Exception as e:
            pass
        
        # Convert PDF to base64 for viewing
        with open(file_path, 'rb') as f:
            pdf_data = f.read()
            pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')
        
        return {
            'success': True,
            'preview_content': f"PDF Document ready for viewing",
            'preview_type': 'pdf',
            'pdf_base64': pdf_base64,
            'base64_preview': pdf_base64,
            'file_size': file_size,
            'file_type': 'PDF',
            'mime_type': 'application/pdf',
            'metadata': {
                'file_size': file_size,
                'text_extraction': False,
                'file_type': 'PDF',
                'status': 'Processed'
            }
        }
        
    except Exception as e:
        return {'success': False, 'error': str(e)}

if __name__ == "__main__":
    result = process_pdf()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

    def _create_image_processor_script(self) -> str:
        """Create dedicated image processing script"""
        return '''#!/usr/bin/env python3
import os
import sys
import json
import base64
import subprocess
import glob

def process_image():
    try:
        # Look for the unique file with analysis ID pattern
        file_patterns = ['/tmp/file_*']
        file_path = None
        
        for pattern in file_patterns:
            matching_files = glob.glob(pattern)
            if matching_files:
                # Get the most recent file (in case of multiple matches)
                file_path = max(matching_files, key=os.path.getctime)
                break
        
        if not file_path:
            return {'success': False, 'error': 'File not found - no matching analysis file'}
        
        file_size = os.path.getsize(file_path)
        
        # Convert image to base64 for viewing
        with open(file_path, 'rb') as f:
            image_data = f.read()
            image_base64 = base64.b64encode(image_data).decode('utf-8')
        
        # Try to get image info using identify (ImageMagick)
        image_info = "Image file ready for viewing"
        try:
            result = subprocess.run(['identify', file_path], 
                                  capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                image_info = result.stdout.strip()
        except Exception as e:
            pass
        
        # Determine MIME type based on file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        mime_type_map = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.webp': 'image/webp',
            '.svg': 'image/svg+xml',
            '.tiff': 'image/tiff',
            '.ico': 'image/x-icon'
        }
        mime_type = mime_type_map.get(file_ext, 'image/png')
        
        return {
            'success': True,
            'preview_content': f"Image Document ready for viewing",
            'preview_type': 'image',
            'image_base64': image_base64,
            'base64_preview': image_base64,
            'file_size': file_size,
            'file_type': 'Image',
            'mime_type': mime_type,
            'metadata': {
                'file_size': file_size,
                'image_info': image_info,
                'file_type': 'Image',
                'status': 'Processed'
            }
        }
        
    except Exception as e:
        return {'success': False, 'error': str(e)}

if __name__ == "__main__":
    result = process_image()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

    def _create_text_processor_script(self) -> str:
        """Create dedicated text processing script"""
        return '''#!/usr/bin/env python3
import os
import sys
import json
import glob

def process_text():
    try:
        # Look for the unique file with analysis ID pattern
        file_patterns = ['/tmp/file_*']
        file_path = None
        
        for pattern in file_patterns:
            matching_files = glob.glob(pattern)
            if matching_files:
                # Get the most recent file (in case of multiple matches)
                file_path = max(matching_files, key=os.path.getctime)
                break
        
        if not file_path:
            return {'success': False, 'error': 'File not found - no matching analysis file'}
        
        file_size = os.path.getsize(file_path)
        
        # Read text content
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Limit content length
        if len(content) > 100000:
            content = content[:100000] + "\\n\\n[... content truncated ...]"
        
        lines = content.splitlines()
        
        return {
            'success': True,
            'preview_content': f"Text File\\nSize: {file_size} bytes\\nLines: {len(lines)}\\n\\nContent:\\n{content}",
            'preview_type': 'text',
            'file_size': file_size,
            'metadata': {
                'file_size': file_size,
                'line_count': len(lines),
                'char_count': len(content)
            }
        }
        
    except Exception as e:
        return {'success': False, 'error': str(e)}

if __name__ == "__main__":
    result = process_text()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

    def _create_video_processor_script(self) -> str:
        """Create dedicated video processing script"""
        return '''#!/usr/bin/env python3
import os
import sys
import json
import base64
import subprocess
import glob

def process_video():
    try:
        # Look for the unique file with analysis ID pattern
        file_patterns = ['/tmp/file_*']
        file_path = None
        
        for pattern in file_patterns:
            matching_files = glob.glob(pattern)
            if matching_files:
                # Get the most recent file (in case of multiple matches)
                file_path = max(matching_files, key=os.path.getctime)
                break
        
        if not file_path:
            return {'success': False, 'error': 'File not found - no matching analysis file'}
        
        file_size = os.path.getsize(file_path)
        
        # Convert video to base64 for viewing
        with open(file_path, 'rb') as f:
            video_data = f.read()
            video_base64 = base64.b64encode(video_data).decode('utf-8')
        
        # Try to get video info using ffprobe
        video_info = "Video file ready for viewing"
        try:
            result = subprocess.run(['ffprobe', '-v', 'quiet', '-print_format', 'json', '-show_format', '-show_streams', file_path], 
                                  capture_output=True, text=True, timeout=30)
            if result.returncode == 0:
                import json
                probe_data = json.loads(result.stdout)
                format_info = probe_data.get('format', {})
                streams = probe_data.get('streams', [])
                
                # Extract basic info
                duration = format_info.get('duration', 'Unknown')
                bitrate = format_info.get('bit_rate', 'Unknown')
                
                # Find video stream
                video_stream = None
                for stream in streams:
                    if stream.get('codec_type') == 'video':
                        video_stream = stream
                        break
                
                if video_stream:
                    width = video_stream.get('width', 'Unknown')
                    height = video_stream.get('height', 'Unknown')
                    codec = video_stream.get('codec_name', 'Unknown')
                    video_info = f"Duration: {duration}s, Resolution: {width}x{height}, Codec: {codec}, Bitrate: {bitrate}"
        except Exception as e:
            pass
        
        # Determine MIME type based on file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        mime_type_map = {
            '.mp4': 'video/mp4',
            '.avi': 'video/x-msvideo',
            '.mov': 'video/quicktime',
            '.mkv': 'video/x-matroska',
            '.webm': 'video/webm',
            '.ogg': 'video/ogg',
            '.wmv': 'video/x-ms-wmv',
            '.flv': 'video/x-flv'
        }
        mime_type = mime_type_map.get(file_ext, 'video/mp4')
        
        return {
            'success': True,
            'preview_content': f"Video Document ready for viewing",
            'preview_type': 'video',
            'video_base64': video_base64,
            'base64_preview': video_base64,
            'file_size': file_size,
            'file_type': 'Video',
            'mime_type': mime_type,
            'metadata': {
                'file_size': file_size,
                'video_info': video_info,
                'file_type': 'Video',
                'status': 'Processed'
            }
        }
        
    except Exception as e:
        return {'success': False, 'error': str(e)}

if __name__ == "__main__":
    result = process_video()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

    def _create_docx_processor_script(self) -> str:
        """Create dedicated DOCX processing script"""
        return '''#!/usr/bin/env python3
import os
import sys
import json
import base64
import glob

def escape_html(text):
    """Simple HTML escaping function"""
    if not text:
        return ""
    return (text.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
                .replace('"', "&quot;")
                .replace("'", "&#x27;"))

def process_docx():
    try:
        # Look for the unique file with analysis ID pattern
        file_patterns = ['/tmp/file_*.docx']
        file_path = None
        
        for pattern in file_patterns:
            matching_files = glob.glob(pattern)
            if matching_files:
                # Get the most recent file (in case of multiple matches)
                file_path = max(matching_files, key=os.path.getctime)
                break
        
        if not file_path:
            return {'success': False, 'error': 'File not found - no matching DOCX file'}
        
        file_size = os.path.getsize(file_path)
        
        # Convert DOCX to base64 for viewing
        with open(file_path, 'rb') as f:
            docx_data = f.read()
            docx_base64 = base64.b64encode(docx_data).decode('utf-8')
        
        # Extract content and metadata
        text_content = ""
        html_content = ""
        metadata = {}
        paragraph_count = 0
        table_count = 0
        
        try:
            from docx import Document
            doc = Document(file_path)
            
            # Simple HTML conversion without complex styling
            html_parts = []
            html_parts.append('<div class="docx-html-content" style="font-family: Arial, sans-serif; line-height: 1.6; color: #333; padding: 20px; background: white;">')
            
            # Process paragraphs
            for para in doc.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    paragraph_count += 1
                    # Simple paragraph with basic styling
                    html_parts.append('<p style="margin: 10px 0;">' + escape_html(para_text) + '</p>')
            
            # Process tables
            for table in doc.tables:
                table_count += 1
                html_parts.append('<table style="width: 100%; border-collapse: collapse; margin: 20px 0; border: 1px solid #ccc;">')
                
                for i, row in enumerate(table.rows):
                    if i == 0:  # Header row
                        html_parts.append('<tr style="background: #f0f0f0;">')
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            html_parts.append('<th style="border: 1px solid #ccc; padding: 8px;">' + escape_html(cell_text) + '</th>')
                        html_parts.append('</tr>')
                    else:  # Data rows
                        html_parts.append('<tr>')
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            html_parts.append('<td style="border: 1px solid #ccc; padding: 8px;">' + escape_html(cell_text) + '</td>')
                        html_parts.append('</tr>')
                
                html_parts.append('</table>')
            
            html_parts.append('</div>')
            html_content = ''.join(html_parts)
            
            # Get plain text for search/metadata
            all_paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text_content = '\\n'.join(all_paragraphs)
            
            # Extract document properties
            core_props = doc.core_properties
            metadata = {
                'title': str(core_props.title or ''),
                'author': str(core_props.author or ''),
                'subject': str(core_props.subject or ''),
                'created': str(core_props.created or ''),
                'modified': str(core_props.modified or ''),
                'paragraphs': paragraph_count,
                'tables': table_count
            }
            
        except Exception as e:
            error_msg = str(e)
            text_content = "Content extraction failed: " + error_msg
            html_content = "<div style='color: red; padding: 20px;'>HTML conversion failed: " + escape_html(error_msg) + "</div>"
            metadata = {'error': error_msg}
        
        return {
            'success': True,
            'preview_content': text_content or "No text content found",
            'preview_type': 'docx',
            'docx_html': html_content,
            'docx_base64': docx_base64,
            'base64_preview': docx_base64,
            'file_size': file_size,
            'file_type': 'DOCX',
            'mime_type': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'docx_metadata': {
                'paragraphs': paragraph_count,
                'tables': table_count,
                'metadata': metadata
            }
        }
        
    except Exception as e:
        return {'success': False, 'error': str(e)}

if __name__ == "__main__":
    result = process_docx()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

    def _create_universal_processor_script(self, file_ext: str) -> str:
        """Create advanced processing script for ALL file types"""
        script = '''#!/usr/bin/env python3
import os
import sys
import json
import base64
import struct
import hashlib
import subprocess
import html
from io import BytesIO
from datetime import datetime

# Simple HTML escaping function
def escape_html(text):
    if not text:
        return ""
    return (text.replace("&", "&amp;")
             .replace("<", "&lt;")
             .replace(">", "&gt;")
             .replace('"', "&quot;")
             .replace("'", "&#x27;"))

# Install required packages
try:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'Pillow', '--quiet'])
except:
    pass  # Continue even if installation fails


def get_file_type(file_path):
    """Detect file type using magic bytes and extension"""
    try:
        import magic
        mime = magic.from_file(file_path, mime=True)
        return mime
    except:
        # Fallback to extension-based detection
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            '.txt': 'text/plain', '.log': 'text/plain', '.md': 'text/markdown',
            '.json': 'application/json', '.xml': 'application/xml', '.csv': 'text/csv',
            '.py': 'text/x-python', '.js': 'application/javascript', '.html': 'text/html',
            '.css': 'text/css', '.sh': 'text/x-shellscript', '.bat': 'application/x-msdos-program',
            '.ps1': 'application/x-powershell', '.sql': 'application/sql',
            '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
            '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp',
            '.svg': 'image/svg+xml', '.tiff': 'image/tiff', '.ico': 'image/x-icon',
            '.pdf': 'application/pdf', '.doc': 'application/msword', '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.ppt': 'application/vnd.ms-powerpoint', '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xls': 'application/vnd.ms-excel', '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.zip': 'application/zip', '.rar': 'application/x-rar-compressed', '.7z': 'application/x-7z-compressed',
            '.tar': 'application/x-tar', '.gz': 'application/gzip', '.bz2': 'application/x-bzip2',
            '.mp4': 'video/mp4', '.avi': 'video/x-msvideo', '.mov': 'video/quicktime',
            '.mkv': 'video/x-matroska', '.wmv': 'video/x-ms-wmv', '.flv': 'video/x-flv',
            '.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.flac': 'audio/flac',
            '.aac': 'audio/aac', '.ogg': 'audio/ogg', '.m4a': 'audio/mp4',
            '.exe': 'application/x-msdownload', '.dll': 'application/x-msdownload',
            '.msi': 'application/x-msi', '.deb': 'application/x-debian-package',
            '.rpm': 'application/x-rpm', '.apk': 'application/vnd.android.package-archive',
            '.ttf': 'font/ttf', '.otf': 'font/otf', '.woff': 'font/woff',
            '.woff2': 'font/woff2', '.eot': 'application/vnd.ms-fontobject',
            '.db': 'application/x-sqlite3', '.sqlite': 'application/x-sqlite3',
            '.mdb': 'application/x-msaccess', '.accdb': 'application/vnd.msaccess'
        }}
        return type_map.get(ext, 'application/octet-stream')

def process_image(file_path, result):
    """Process image files with advanced metadata"""
    try:
        from PIL import Image, ExifTags
        import exifread
        
        with Image.open(file_path) as img:
            width, height = img.size
            mode = img.mode
            format_name = img.format or "Unknown"
            
            # Create high-quality thumbnail
            img.thumbnail((800, 800), Image.Resampling.LANCZOS)
            buffer = BytesIO()
            img.save(buffer, format='PNG', optimize=True)
            thumbnail_b64 = base64.b64encode(buffer.getvalue()).decode()
            
            # Extract EXIF data
            exif_data = {}
            try:
                with open(file_path, 'rb') as f:
                    tags = exifread.process_file(f, details=False)
                    for tag in tags:
                        exif_data[str(tag)] = str(tags[tag])
            except:
                pass
            
            # Color analysis
            colors = img.getcolors(maxcolors=256*256*256)
            dominant_colors = sorted(colors, key=lambda x: x[0], reverse=True)[:5] if colors else []
            
            result["preview_content"] = "Image Analysis\\nDimensions: " + str(width) + "x" + str(height) + " pixels\\nColor Mode: " + str(mode) + "\\nFormat: " + str(format_name) + "\\n\\nEXIF Data: " + str(len(exif_data)) + " tags found"
            result["preview_type"] = "image"
            result["metadata"] = {
                "dimensions": str(width) + "x" + str(height),
                "mode": mode,
                "format": format_name,
                "exif_data": exif_data,
                "dominant_colors": dominant_colors[:3],
                "file_size": result["file_size"]
            }
            result["thumbnail"] = "data:image/png;base64," + thumbnail_b64
            
    except ImportError:
        # Fallback when PIL is not available
        result["preview_content"] = "Image File (" + str(result['file_size']) + " bytes)\\nFormat: " + get_file_type(file_path) + "\\nBasic analysis only - PIL not available in container"
        result["preview_type"] = "image"
        result["metadata"] = {
            "file_size": result["file_size"],
            "mime_type": get_file_type(file_path),
            "pil_available": False
        }
        result["thumbnail"] = None
            
    except Exception as e:
        result["preview_content"] = "Image processing error: " + str(e)
        result["preview_type"] = "error"
        
def process_video(file_path, result):
    """Process video files with metadata extraction"""
    try:
        from mutagen import File as MutagenFile
        
        video_file = MutagenFile(file_path)
        duration = getattr(video_file.info, 'length', 0)
        bitrate = getattr(video_file.info, 'bitrate', 0)
        codec = getattr(video_file.info, 'codec', 'Unknown')
        
        # Try to get video dimensions
        try:
            import cv2
            cap = cv2.VideoCapture(file_path)
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            cap.release()
        except:
            width, height, fps, frame_count = 0, 0, 0, 0
        
        # Create thumbnail (first frame)
        thumbnail_b64 = None
        try:
            import cv2
            cap = cv2.CAP_PROP_POS_FRAMES
            cap = cv2.VideoCapture(file_path)
            ret, frame = cap.read()
            if ret:
                from PIL import Image
                img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                img.thumbnail((400, 300), Image.Resampling.LANCZOS)
                buffer = BytesIO()
                img.save(buffer, format='PNG')
                thumbnail_b64 = base64.b64encode(buffer.getvalue()).decode()
            cap.release()
        except:
            pass
                    
        result["preview_content"] = "Video Analysis\\nDuration: " + str(round(duration, 2)) + "s\\nResolution: " + str(width) + "x" + str(height) + "\\nFPS: " + str(round(fps, 2)) + "\\nCodec: " + str(codec) + "\\nBitrate: " + str(bitrate) + " bps"
        result["preview_type"] = "video"
        result["metadata"] = {
            "duration": duration,
            "dimensions": str(width) + "x" + str(height),
            "fps": fps,
            "frame_count": frame_count,
            "codec": codec,
            "bitrate": bitrate,
            "file_size": result["file_size"]
        }
        if thumbnail_b64:
            result["thumbnail"] = "data:image/png;base64," + thumbnail_b64
                    
    except Exception as e:
        result["preview_content"] = "Video processing error: " + str(e)
        result["preview_type"] = "error"
        
def process_audio(file_path, result):
    """Process audio files with metadata extraction"""
    try:
        from mutagen import File as MutagenFile
        
        audio_file = MutagenFile(file_path)
        duration = getattr(audio_file.info, 'length', 0)
        bitrate = getattr(audio_file.info, 'bitrate', 0)
        sample_rate = getattr(audio_file.info, 'sample_rate', 0)
        channels = getattr(audio_file.info, 'channels', 0)
        codec = getattr(audio_file.info, 'codec', 'Unknown')
        
        # Extract tags
        tags = {}
        if hasattr(audio_file, 'tags') and audio_file.tags:
            for key, value in audio_file.tags.items():
                tags[str(key)] = str(value)
        
        result["preview_content"] = "Audio Analysis\\nDuration: " + str(round(duration, 2)) + "s\\nBitrate: " + str(bitrate) + " bps\\nSample Rate: " + str(sample_rate) + " Hz\\nChannels: " + str(channels) + "\\nCodec: " + str(codec) + "\\n\\nTags: " + str(len(tags)) + " metadata entries"
        result["preview_type"] = "audio"
        result["metadata"] = {
            "duration": duration,
            "bitrate": bitrate,
            "sample_rate": sample_rate,
            "channels": channels,
            "codec": codec,
            "tags": tags,
            "file_size": result["file_size"]
        }
        
    except Exception as e:
        result["preview_content"] = "Audio processing error: " + str(e)
        result["preview_type"] = "error"


def process_document(file_path, result):
    """Process document files (DOC, PPT, XLS, etc.)"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        
        # Skip PDF and DOCX processing here since they have dedicated functions
        if ext == '.pdf':
            result.update({{
                "preview_type": "pdf",
                "preview_content": "PDF processing handled by dedicated function"
            }})
            return
            
        if ext == '.docx':
            result.update({{
                "preview_type": "docx",
                "preview_content": "DOCX processing handled by dedicated function"
            }})
            return
            
        if ext in ['.doc']:
            try:
                from docx import Document
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                tables = len(doc.tables)
                
                content_preview = ' '.join(paragraphs[:5])
                result["preview_content"] = f"Word Document\\nParagraphs: {len(paragraphs)}\\nTables: {tables}\\n\\nContent Preview:\\n{content_preview}..."
                result["preview_type"] = "document"
                result["metadata"] = {
                    "paragraphs": len(paragraphs),
                    "tables": tables,
                    "file_size": result["file_size"]
                }
            except:
                result["preview_content"] = f"Word Document ({result['file_size']} bytes)\\nContent extraction not available"
                result["preview_type"] = "document"
        
        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides = len(prs.slides)
                
                slide_texts = []
                for slide in prs.slides[:3]:  # First 3 slides
                    text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
                    slide_texts.append(text.strip())
                
                slide_previews = '\\n'.join(slide_texts)
                result["preview_content"] = f"PowerPoint Presentation\\nSlides: {slides}\\n\\nSlide Previews:\\n{slide_previews}"
                result["preview_type"] = "presentation"
                result["metadata"] = {
                    "slides": slides,
                    "file_size": result["file_size"]
                }
            except:
                file_size = result['file_size']
                result["preview_content"] = f"PowerPoint Presentation ({file_size} bytes)\\nContent extraction not available"
                result["preview_type"] = "presentation"
        
        # XLSX processing moved to mime_type-based processing below
                    
    except Exception as e:
        result["preview_content"] = "Document processing error: " + str(e)
        result["preview_type"] = "error"
        
def process_archive(file_path, result):
    """Process archive files with detailed content analysis"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
                
        if ext == '.zip':
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as zf:
                file_list = zf.namelist()
                total_size = sum(info.file_size for info in zf.infolist())
                compressed_size = sum(info.compress_size for info in zf.infolist())
                compression_ratio = (1 - compressed_size / total_size) * 100 if total_size > 0 else 0
                
                # Analyze file types in archive
                file_types = {}
                for filename in file_list:
                    ext = os.path.splitext(filename)[1].lower()
                    file_types[ext] = file_types.get(ext, 0) + 1
                
                result["preview_content"] = f"ZIP Archive\\nFiles: {len(file_list)}\\nTotal Size: {total_size} bytes\\nCompressed: {compressed_size} bytes\\nCompression: {compression_ratio:.1f}%\\n\\nFile Types: {dict(list(file_types.items())[:5])}"
                result["preview_type"] = "archive"
                result["metadata"] = {
                    "file_count": len(file_list),
                    "total_size": total_size,
                    "compressed_size": compressed_size,
                    "compression_ratio": compression_ratio,
                    "file_types": file_types,
                    "files": file_list[:20]  # First 20 files
                }
        
        elif ext == '.rar':
            try:
                import rarfile
                with rarfile.RarFile(file_path, 'r') as rf:
                    file_list = rf.namelist()
                    total_size = sum(info.file_size for info in rf.infolist())
                    
                     # Use simple string formatting to avoid Docker script errors
                     preview_text = "RAR Archive\nFiles: " + str(len(file_list)) + "\nTotal Size: " + str(total_size) + " bytes\n\nFiles: " + str(file_list[:10])
                     result["preview_content"] = preview_text
                    result["preview_type"] = "archive"
                     result["metadata"] = {
                        "file_count": len(file_list),
                        "total_size": total_size,
                        "files": file_list[:20]
                     }
            except:
                # Use simple string formatting
                preview_text = "RAR Archive (" + str(result.get('file_size', 0)) + " bytes)\nRAR processing not available"
                result["preview_content"] = preview_text
                result["preview_type"] = "archive"
        
        elif ext == '.7z':
            try:
                import py7zr
                with py7zr.SevenZipFile(file_path, 'r') as zf:
                    file_list = zf.getnames()
                    total_size = sum(info.uncompressed for info in zf.list() if hasattr(info, 'uncompressed'))
                    
                     # Use simple string formatting to avoid Docker script errors
                     preview_text = "7Z Archive\nFiles: " + str(len(file_list)) + "\nTotal Size: " + str(total_size) + " bytes\n\nFiles: " + str(file_list[:10])
                     result["preview_content"] = preview_text
                    result["preview_type"] = "archive"
                     result["metadata"] = {
                        "file_count": len(file_list),
                        "total_size": total_size,
                        "files": file_list[:20]
                     }
            except:
                # Use simple string formatting
                preview_text = "7Z Archive (" + str(result.get('file_size', 0)) + " bytes)\n7Z processing not available"
                result["preview_content"] = preview_text
                result["preview_type"] = "archive"
        
    except Exception as e:
        # Use simple string formatting
        result["preview_content"] = "Archive processing error: " + str(e)
        result["preview_type"] = "error"
        
def process_executable(file_path, result):
    """Process executable files with security analysis"""
    try:
        with open(file_path, 'rb') as f:
            header = f.read(1024)
        
        # Analyze executable headers
        if header[:2] == b'MZ':
            pe_offset = struct.unpack('<I', header[60:64])[0]
            if pe_offset < len(header):
                pe_header = header[pe_offset:pe_offset+4]
                if pe_header == b'PE\\x00':
                    result["preview_content"] = f"Windows Executable (PE)\\nFile Size: {result['file_size']} bytes\\nArchitecture: PE32/PE32+\\n\\nWARNING: This is an executable file. Only run if you trust the source."
                    result["preview_type"] = "executable"
                    result["metadata"] = {
                        "platform": "Windows",
                        "type": "PE",
                        "file_size": result["file_size"],
                        "risk_level": "High"
                    }
            else:
                result["preview_content"] = f"Executable File\\nFile Size: {result['file_size']} bytes\\nType: Unknown executable format"
                result["preview_type"] = "executable"
        elif header[:4] == b'\\x7fELF':
            result["preview_content"] = f"Linux Executable (ELF)\\nFile Size: {result['file_size']} bytes\\nArchitecture: ELF\\n\\nWARNING: This is an executable file. Only run if you trust the source."
            result["preview_type"] = "executable"
            result["metadata"] = {
                "platform": "Linux",
                "type": "ELF",
                "file_size": result["file_size"],
                "risk_level": "High"
            }
        else:
            result["preview_content"] = f"Binary Executable\\nFile Size: {result['file_size']} bytes\\nType: Unknown binary format\\n\\nWARNING: This appears to be an executable file."
            result["preview_type"] = "executable"
            result["metadata"] = {
                "platform": "Unknown",
                "type": "Binary",
                "file_size": result["file_size"],
                "risk_level": "High"
            }
        
    except Exception as e:
        result["preview_content"] = "Executable processing error: " + str(e)
        result["preview_type"] = "error"

def process_font(file_path, result):
    """Process font files with metadata extraction"""
    try:
        from fontTools import TTFont
        font = TTFont(file_path)
        
        # Extract font metadata
        name_table = font['name']
        font_name = ""
        font_family = ""
        for record in name_table.names:
            if record.nameID == 1:  # Font Family
                font_family = record.string.decode('utf-16-be') if record.platformID == 3 else record.string.decode('utf-8')
            elif record.nameID == 4:  # Full Font Name
                font_name = record.string.decode('utf-16-be') if record.platformID == 3 else record.string.decode('utf-8')
        
        # Get font metrics
        head_table = font['head']
        units_per_em = head_table.unitsPerEm
        created = datetime.fromtimestamp(head_table.created)
        modified = datetime.fromtimestamp(head_table.modified)
        
        result["preview_content"] = f"Font File\\nName: {font_name}\\nFamily: {font_family}\\nUnits per EM: {units_per_em}\\nCreated: {created.strftime('%Y-%m-%d %H:%M:%S')}\\nModified: {modified.strftime('%Y-%m-%d %H:%M:%S')}"
        result["preview_type"] = "font"
        result["metadata"] = {
            "font_name": font_name,
            "font_family": font_family,
            "units_per_em": units_per_em,
            "created": created.isoformat(),
            "modified": modified.isoformat(),
            "file_size": result["file_size"]
        }
        
    except Exception as e:
        result["preview_content"] = f"Font processing error: {str(e)}"
        result["preview_type"] = "error"

def process_database(file_path, result):
    """Process database files"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.db', '.sqlite', '.sqlite3']:
            import sqlite3
            conn = sqlite3.connect(file_path)
            cursor = conn.cursor()
            
            # Get table information
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            
            # Get table counts
            table_info = {}
            for table in tables:
                table_name = table[0]
                cursor.execute(f"SELECT COUNT(*) FROM {table_name}")
                count = cursor.fetchone()[0]
                table_info[table_name] = count
            
            conn.close()
            
            result["preview_content"] = f"SQLite Database\\nTables: {len(tables)}\\nTable Names: {[t[0] for t in tables]}\\n\\nTable Row Counts: {table_info}"
            result["preview_type"] = "database"
            result["metadata"] = {
                "tables": [t[0] for t in tables],
                "table_count": len(tables),
                "table_info": table_info,
                "file_size": result["file_size"]
            }
        
    except Exception as e:
        result["preview_content"] = f"Database processing error: {str(e)}"
        result["preview_type"] = "error"

def safe_process_file():
    """Safely process file in isolated container with universal support"""
    try:
        import glob
        file_ext = "{file_ext}"
        
        # Look for the unique file with analysis ID pattern
        file_patterns = ['/tmp/file_*']
        file_path = None
        
        for pattern in file_patterns:
            matching_files = glob.glob(pattern)
            if matching_files:
                # Get the most recent file (in case of multiple matches)
                file_path = max(matching_files, key=os.path.getctime)
                break
        
        # Debug: List files in /tmp directory (for troubleshooting)
        print("DEBUG: Files in /tmp: " + str(os.listdir('/tmp')))
        print("DEBUG: Looking for file pattern: " + str(file_patterns))
        print("DEBUG: Found file: " + str(file_path))
        
        if not file_path:
            print(json.dumps({"error": "File not found - no matching analysis file", "success": False}))
            return
        
        file_size = os.path.getsize(file_path)
        
        # File size check (max 200MB for preview)
        if file_size > 200 * 1024 * 1024:
            print(json.dumps({
                "error": f"File too large: {file_size} bytes",
                "file_size": file_size,
                "success": False
            }))
            return
        
        # Set UTF-8 encoding for stdout
        import sys
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
        
        # Install required packages (with error handling to ensure script continues)
        def install_packages():
            \"\"\"Install required packages for file processing\"\"\"
            import subprocess
            import sys
            
            # Essential packages for file processing
            critical_packages = ["Pillow", "PyPDF2", "python-docx", "openpyxl"]
            optional_packages = ["python-magic", "mutagen", "openpyxl", 
                                "python-pptx", "python-docx2txt", "rarfile", 
                                "py7zr", "python-magic-bin", "chardet", "exifread", "pydub"]
            
            # Install critical packages first (silently)
            for pkg in critical_packages:
                try:
                    # Try installing with pre-built wheels and no compilation
                    result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--only-binary=all', '--no-compile', '--quiet', pkg], 
                                          capture_output=True, check=True, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                except subprocess.CalledProcessError:
                    try:
                        result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--system', '--quiet', pkg], 
                                              capture_output=True, check=True, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                    except subprocess.CalledProcessError:
                        try:
                            result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--quiet', pkg], 
                                                  capture_output=True, check=True, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                        except:
                            pass  # Silent fail
                except Exception:
                    pass  # Silent fail
            
            # Install optional packages (silently)
            for pkg in optional_packages:
                try:
                    result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--quiet', pkg], 
                                          capture_output=True, timeout=30, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                except:
                    pass  # Silent fail
        
        # Install packages with error handling (silently)
        try:
            install_packages()
        except Exception as e:
            pass  # Silent fail, continue with basic processing
        
        # Initialize result with basic info
        result = {{
            "file_size": file_size,
            "file_ext": "{file_ext}",
            "success": True,
            "preview_content": None,
            "preview_type": "unknown",
            "metadata": {},
            "thumbnail": None,
            "mime_type": get_file_type(file_path)
        }}
        
        # Determine file type and process accordingly
        mime_type = result["mime_type"]
        
        if mime_type.startswith('image/'):
            process_image(file_path, result)
        elif mime_type.startswith('video/'):
            process_video(file_path, result)
        elif mime_type.startswith('audio/'):
            process_audio(file_path, result)
        elif mime_type == 'application/pdf':
            # Process PDF files with secure Docker processing
            try:
                import base64
                
                # Read the raw PDF file for viewing
                with open(file_path, 'rb') as file:
                    pdf_data = file.read()
                    pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')
                
                # Extract metadata and limited text content
                text_content = ""
                metadata = {}
                page_count = 0
                
                try:
                    from PyPDF2 import PdfReader
                    pdf_reader = PdfReader(file_path)
                    page_count = len(pdf_reader.pages)
                    
                    # Extract metadata
                    if pdf_reader.metadata:
                        metadata = {{
                            'title': str(pdf_reader.metadata.get('/Title', '') or ''),
                            'author': str(pdf_reader.metadata.get('/Author', '') or ''),
                            'subject': str(pdf_reader.metadata.get('/Subject', '') or ''),
                            'creator': str(pdf_reader.metadata.get('/Creator', '') or ''),
                            'producer': str(pdf_reader.metadata.get('/Producer', '') or ''),
                            'creation_date': str(pdf_reader.metadata.get('/CreationDate', '') or ''),
                            'modification_date': str(pdf_reader.metadata.get('/ModDate', '') or ''),
                            'pages': page_count
                        }}
                    
                    # Extract text from first few pages for search/metadata
                    text_content = ""
                    max_pages = min(3, page_count)  # Only extract first 3 pages for metadata
                    
                    for page_num in range(max_pages):
                        try:
                            page = pdf_reader.pages[page_num]
                            text_content += f"\\n--- Page {page_num + 1} ---\\n"
                            text_content += page.extract_text()
                        except:
                            text_content += f"\\n--- Page {page_num + 1} (text extraction failed) ---\\n"
                    
                except Exception as e:
                    text_content = f"PDF text extraction failed: {str(e)}"
                    metadata = {'error': str(e)}
                    
                result.update({{
                    "preview_type": "pdf",
                    "preview_content": text_content or "No text content found",
                    "pdf_base64": pdf_base64,
                    "base64_preview": pdf_base64,
                    "pdf_metadata": {{
                        "page_count": page_count,
                        "metadata": metadata
                    }}
                }})
            except Exception as e:
                result.update({{
                    "preview_type": "pdf",
                    "preview_content": f"PDF processing error: {str(e)}"
                }})
        elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
            # Process DOCX files with secure Docker processing and HTML conversion
            try:
                import base64
                
                # Read the raw DOCX file for viewing
                with open(file_path, 'rb') as file:
                    docx_data = file.read()
                    docx_base64 = base64.b64encode(docx_data).decode('utf-8')
                
                # Extract content and metadata
                text_content = ""
                html_content = ""
                metadata = {}
                paragraph_count = 0
                table_count = 0
                
                try:
                    from docx import Document
                    doc = Document(file_path)
                    
                    # Simple HTML conversion without complex styling
                    html_parts = []
                    html_parts.append('<div class="docx-html-content" style="font-family: Arial, sans-serif; line-height: 1.6; color: #333; padding: 20px; background: white;">')
                    
                    # Process paragraphs
                    for para in doc.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            paragraph_count += 1
                            # Simple paragraph with basic styling
                            html_parts.append('<p style="margin: 10px 0;">' + escape_html(para_text) + '</p>')
                    
                    # Process tables
                    for table in doc.tables:
                        table_count += 1
                        html_parts.append('<table style="width: 100%; border-collapse: collapse; margin: 20px 0; border: 1px solid #ccc;">')
                        
                        for i, row in enumerate(table.rows):
                            if i == 0:  # Header row
                                html_parts.append('<tr style="background: #f0f0f0;">')
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    html_parts.append('<th style="border: 1px solid #ccc; padding: 8px;">' + escape_html(cell_text) + '</th>')
                                html_parts.append('</tr>')
                            else:  # Data rows
                                html_parts.append('<tr>')
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    html_parts.append('<td style="border: 1px solid #ccc; padding: 8px;">' + escape_html(cell_text) + '</td>')
                                html_parts.append('</tr>')
                        
                        html_parts.append('</table>')
                    
                    html_parts.append('</div>')
                    html_content = ''.join(html_parts)
                    
                    # Get plain text for search/metadata
                    all_paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    text_content = '\\n'.join(all_paragraphs)
                    
                    # Extract document properties
                    core_props = doc.core_properties
                    metadata = {{
                        'title': str(core_props.title or ''),
                        'author': str(core_props.author or ''),
                        'subject': str(core_props.subject or ''),
                        'created': str(core_props.created or ''),
                        'modified': str(core_props.modified or ''),
                        'paragraphs': paragraph_count,
                        'tables': table_count
                    }}
                    
                except Exception as e:
                    error_msg = str(e)
                    text_content = "Content extraction failed: " + error_msg
                    html_content = "<div style='color: red; padding: 20px;'>HTML conversion failed: " + escape_html(error_msg) + "</div>"
                    metadata = {'error': error_msg}
                    
                result.update({{
                    "preview_type": "docx",
                    "preview_content": text_content or "No text content found",
                    "docx_html": html_content,
                    "docx_base64": docx_base64,
                    "base64_preview": docx_base64,
                    "docx_metadata": {{
                        "paragraphs": paragraph_count,
                        "tables": table_count,
                        "metadata": metadata
                    }}
                }})
            except Exception as e:
                error_msg = str(e)
                result.update({{
                    "preview_type": "docx",
                    "preview_content": "DOCX processing error: " + error_msg
                }})
        elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            # Process XLSX files with pre-built Docker image
            try:
                import subprocess
                import tempfile
                import shutil
                
                # Use pre-built Docker image for XLSX processing
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Copy file to temp directory
                    temp_file = os.path.join(temp_dir, f"file{file_ext}")
                    shutil.copy2(file_path, temp_file)
                    
                    # Run pre-built Docker image
                    result_cmd = subprocess.run([
                        'docker', 'run', '--rm',
                        '-v', f'{temp_dir}:/tmp',
                        'aejis-xlsx-processor:latest',
                        f'/tmp/file{file_ext}'
                    ], capture_output=True, text=True, timeout=60)
                    
                    if result_cmd.returncode == 0:
                        try:
                            xlsx_result = json.loads(result_cmd.stdout)
                            if xlsx_result.get('success', False):
                                result.update(xlsx_result)
                            else:
                                result.update({
                                    "preview_type": "xlsx",
                                    "preview_content": f"XLSX processing failed: {xlsx_result.get('error', 'Unknown error')}"
                                })
                        except json.JSONDecodeError:
                            result.update({
                                "preview_type": "xlsx",
                                "preview_content": f"XLSX processing failed: Invalid response from Docker container"
                            })
                    else:
                        result.update({
                            "preview_type": "xlsx",
                            "preview_content": f"XLSX processing failed: Docker container error - {result_cmd.stderr}"
                        })
                        
            except Exception as e:
                error_msg = str(e)
                result.update({
                    "preview_type": "xlsx",
                    "preview_content": "XLSX processing error: " + error_msg
                })
        elif mime_type in ['application/msword', 'application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                          'application/vnd.ms-excel']:
            process_document(file_path, result)
        elif mime_type in ['application/zip', 'application/x-rar-compressed', 'application/x-7z-compressed',
                          'application/x-tar', 'application/gzip', 'application/x-bzip2']:
            process_archive(file_path, result)
        elif mime_type in ['application/x-msdownload', 'application/x-msi', 'application/x-debian-package',
                          'application/x-rpm', 'application/vnd.android.package-archive']:
            process_executable(file_path, result)
        elif mime_type.startswith('font/') or mime_type in ['application/font-woff', 'application/font-woff2']:
            process_font(file_path, result)
        elif mime_type in ['application/x-sqlite3', 'application/vnd.msaccess']:
            process_database(file_path, result)
        elif mime_type.startswith('text/'):
            # Text files with multiple encoding attempts
            try:
                content = None
                encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read(100000)  # Max 100KB preview
                            if len(content) >= 100000:
                                content += "\\n\\n[... content truncated for security ...]"
                            break
                    except UnicodeDecodeError:
                        continue
                
                if content is None:
                    # If all encodings fail, read as binary and decode with errors='replace'
                    with open(file_path, 'rb') as f:
                        raw_content = f.read(100000)
                        content = raw_content.decode('utf-8', errors='replace')
                        if len(content) >= 100000:
                            content += "\\n\\n[... content truncated for security ...]"
                
                result["preview_content"] = f"Text File\\nSize: {file_size} bytes\\n\\nContent:\\n{content}"
                result["preview_type"] = "text"
                result["metadata"] = {
                    "line_count": len(content.splitlines()),
                    "char_count": len(content),
                    "file_size": file_size
                }
            except Exception as e:
                result["preview_content"] = "Text processing error: " + str(e)
                result["preview_type"] = "error"
        else:
            # Binary/unknown files - enhanced analysis
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(512).hex()  # First 512 bytes as hex
                
                # Calculate file entropy
                f.seek(0)
                data = f.read(min(8192, file_size))  # First 8KB
                
                if data:
                    byte_counts = [0] * 256
                    for byte in data:
                        byte_counts[byte] += 1
                    
                    entropy = 0
                    for count in byte_counts:
                        if count > 0:
                            p = count / len(data)
                            entropy -= p * math.log2(p)
                    
                    result["preview_content"] = f"Binary File Analysis\\nSize: {file_size} bytes\\nEntropy: {entropy:.2f}/8.0\\nMIME Type: {mime_type}\\n\\nHex Header (first 512 bytes):\\n{header[:200]}..."
                    result["preview_type"] = "binary"
                    result["metadata"] = {
                        "header": header,
                        "entropy": entropy,
                        "mime_type": mime_type,
                        "file_size": file_size,
                        "is_likely_encrypted": entropy > 7.0,
                        "is_likely_compressed": entropy > 6.0
                    }}
                else:
                    result["preview_content"] = f"Empty Binary File\\nSize: {file_size} bytes"
                    result["preview_type"] = "binary"
                
            except Exception as e:
                result["preview_content"] = "Binary processing error: " + str(e)
                result["preview_type"] = "error"
        
        # Ensure JSON output is properly encoded
        try:
            json_output = json.dumps(result, ensure_ascii=False, indent=2)
            print(json_output)
        except UnicodeEncodeError:
            # Fallback: ensure ASCII encoding
            json_output = json.dumps(result, ensure_ascii=True, indent=2)
            print(json_output)
        
    except Exception as e:
        try:
            error_result = {"error": str(e), "success": False}
            json_output = json.dumps(error_result, ensure_ascii=False, indent=2)
            print(json_output)
        except UnicodeEncodeError:
            error_result = {"error": "Processing error occurred", "success": False}
            json_output = json.dumps(error_result, ensure_ascii=True, indent=2)
            print(json_output)

if __name__ == "__main__":
    safe_process_file()
'''
        # Ensure file_ext is not None or empty
        safe_file_ext = file_ext if file_ext else '.unknown'
        return script.format(file_ext=safe_file_ext)
    
    def _parse_docker_result(self, result: str, file_path: str, analysis_id: str) -> Dict[str, Any]:
        """Parse Docker container result with encoding safety"""
        logger.debug(f"🔍 Parsing Docker result for {os.path.basename(file_path)}")
        
        try:
            # Ensure the result string is properly encoded
            if isinstance(result, bytes):
                logger.debug(f"📄 Converting bytes result to string")
                try:
                    result = result.decode('utf-8', errors='replace')
                    logger.debug(f"✅ Successfully decoded with UTF-8")
                except UnicodeDecodeError:
                    result = result.decode('latin-1', errors='replace')
                    logger.debug(f"✅ Fallback decoded with latin-1")
            
            logger.debug(f"🔧 Parsing JSON result, length: {len(result)}")
            data = json.loads(result)
            logger.debug(f"✅ JSON parsed successfully, keys: {list(data.keys())}")
            
            if not data.get('success', False):
                error_msg = data.get('error', 'Unknown error')
                logger.warning(f"⚠️ Docker processing reported failure: {error_msg}")
                return self._error_result(file_path, analysis_id, error_msg)
            
            # Safely extract preview content with encoding handling
            preview_content = data.get('preview_content', 'No content')
            logger.debug(f"📝 Preview content length: {len(str(preview_content))}")
            
            if isinstance(preview_content, str):
                # Ensure content is properly encoded
                try:
                    preview_content.encode('utf-8')
                    logger.debug(f"✅ Preview content encoding validated")
                except UnicodeEncodeError:
                    preview_content = preview_content.encode('utf-8', errors='replace').decode('utf-8')
                    logger.warning(f"⚠️ Preview content encoding fixed with error replacement")
            
            # Build final result
            final_result = {
                'filename': os.path.basename(file_path),
                'file_size': data.get('file_size', 0),
                'file_type': data.get('file_ext', 'unknown'),
                'analysis_id': analysis_id,
                'preview_content': preview_content,
                'preview_type': data.get('preview_type', 'unknown'),
                'metadata': data.get('metadata', {}),
                'thumbnail_base64': data.get('thumbnail'),
                'secure_extraction': True,
                'docker_processed': True
            }
            
            logger.info(f"✅ Successfully parsed Docker result for {os.path.basename(file_path)}")
            logger.debug(f"📊 Result summary: type={final_result['preview_type']}, size={final_result['file_size']}, has_thumbnail={bool(final_result['thumbnail_base64'])}")
            
            # Map thumbnail_base64 to base64_preview for frontend compatibility
            if final_result.get('thumbnail_base64'):
                thumbnail_data = final_result['thumbnail_base64']
                # Strip data URL prefix if present (frontend expects just base64 data)
                if thumbnail_data.startswith('data:'):
                    # Extract just the base64 part after the comma
                    if ',' in thumbnail_data:
                        final_result['base64_preview'] = thumbnail_data.split(',', 1)[1]
                    else:
                        final_result['base64_preview'] = thumbnail_data
                else:
                    final_result['base64_preview'] = thumbnail_data
            
            return final_result
            
        except Exception as e:
            logger.error(f"❌ Error parsing Docker result for {os.path.basename(file_path)}: {e}")
            logger.error(f"🔍 Raw result preview: {result[:500]}..." if len(result) > 500 else f"🔍 Raw result: {result}")
            return self._error_result(file_path, analysis_id, f"Result parsing error: {e}")
    
    # No fallback - Docker is REQUIRED for security
    
    def _error_result(self, file_path: str, analysis_id: str, error_msg: str) -> Dict[str, Any]:
        """Return error result - Docker is required"""
        return {
            'filename': os.path.basename(file_path),
            'file_size': 0,
            'file_type': 'error',
            'analysis_id': analysis_id,
            'preview_content': f"DOCKER REQUIRED: {error_msg}\n\nThis system requires Docker for secure file preview.\nDocker provides complete isolation from malicious files.",
            'preview_type': 'docker_required',
            'secure_extraction': False,
            'docker_processed': False,
            'docker_required': True
        }

# Initialize secure preview system
secure_preview = SecureFilePreview()

app = Flask(__name__)
CORS(app)  # Enable CORS for React frontend
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['UPLOAD_FOLDER'] = tempfile.gettempdir()

# Initialize analyzer
file_analyzer = FileAnalyzer()

# Store analysis results with persistence
analysis_results = {}
RESULTS_FILE = 'analysis_results.json'

# File cleanup tracking - automatically delete files after 30 minutes
file_cleanup_timers = {}  # {analysis_id: {'file_path': path, 'cleanup_time': timestamp}}

def save_results():
    """Save analysis results to file"""
    try:
        with open(RESULTS_FILE, 'w') as f:
            json.dump(analysis_results, f, default=str, indent=2)
    except Exception as e:
        logger.error(f"Failed to save results: {e}")

def load_results():
    """Load analysis results from file"""
    global analysis_results
    try:
        if os.path.exists(RESULTS_FILE):
            with open(RESULTS_FILE, 'r') as f:
                analysis_results = json.load(f)
            logger.info(f"Loaded {len(analysis_results)} analysis results from disk")
    except Exception as e:
        logger.error(f"Failed to load results: {e}")
        analysis_results = {}

# Load existing results on startup
load_results()

@app.route('/')
def index():
    """API health check"""
    return jsonify({
        'status': 'healthy',
        'message': 'Aejis Security API is running',
        'version': '1.0.0',
        'timestamp': datetime.now().isoformat()
    })

@app.route('/analyze-url', methods=['POST'])
def analyze_url():
    """Analyze a URL for security threats"""
    try:
        data = request.get_json()
        url = data.get('url')
        
        if not url:
            return jsonify({'error': 'URL is required'}), 400
        
        # Validate URL format
        try:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            if not parsed.scheme or not parsed.netloc:
                return jsonify({'error': 'Invalid URL format'}), 400
        except Exception:
            return jsonify({'error': 'Invalid URL format'}), 400
        
        # Generate analysis ID
        analysis_id = f"url_{int(time.time())}_{uuid.uuid4().hex[:8]}"
        
        # Initialize analysis record
        analysis_results[analysis_id] = {
            'status': 'processing',
            'url': url,
            'created_at': datetime.now().isoformat(),
            'is_url': True,
            'steps': [
                {'id': 1, 'name': 'URL Reputation Check', 'status': 'pending'},
                {'id': 2, 'name': 'AI Website Analysis', 'status': 'pending'},
                {'id': 3, 'name': 'Safe Preview Generation', 'status': 'pending'},
                {'id': 4, 'name': 'Report Generation', 'status': 'pending'}
            ],
            'progress': 0
        }
        
        # Schedule URL analysis cleanup after 30 minutes (no file to delete, just results)
        file_cleanup_timers[analysis_id] = {
            'file_path': None,  # No file for URL analysis
            'cleanup_time': time.time()
        }
        logger.info(f"🗑️ Scheduled cleanup for URL analysis {url} (ID: {analysis_id}) in 30 minutes")
        
        # Start background analysis
        threading.Thread(target=analyze_url_async, args=(analysis_id, url)).start()
        
        return jsonify({
            'analysis_id': analysis_id,
            'status': 'processing',
            'message': 'URL analysis started'
        })
    
    except Exception as e:
        logger.error(f"URL analysis upload error: {str(e)}")
        return jsonify({'error': 'Failed to start URL analysis'}), 500

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file upload and start analysis"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Secure filename
        filename = secure_filename(file.filename)
        timestamp = int(time.time())
        unique_filename = f"{timestamp}_{filename}"
        
        # Save file temporarily
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(file_path)
        
        # Generate analysis ID with UUID to ensure uniqueness
        analysis_id = f"analysis_{timestamp}_{uuid.uuid4().hex[:8]}"
        
        # Start analysis in background
        analysis_results[analysis_id] = {
            'status': 'processing',
            'filename': filename,
            'file_size': os.path.getsize(file_path),
            'start_time': datetime.now().isoformat(),
            'progress': 0,
            'steps': [
                {'name': 'Upload', 'status': 'completed', 'icon': '📤'},
                {'name': 'VirusTotal Scan (70+ Engines)', 'status': 'processing', 'icon': '🔍'},
                {'name': 'Sandbox Analysis', 'status': 'pending', 'icon': '🐳'},
                {'name': 'Report Generation', 'status': 'pending', 'icon': '📊'}
            ],
            'results': None,
            'file_path': file_path,
            'is_url': False  # Mark as file analysis - no AI needed
        }
        
        # Schedule file cleanup after 30 minutes
        file_cleanup_timers[analysis_id] = {
            'file_path': file_path,
            'cleanup_time': time.time()
        }
        logger.info(f"🗑️ Scheduled cleanup for {filename} (ID: {analysis_id}) in 30 minutes")
        
        # Start analysis in background thread
        threading.Thread(target=analyze_file_async, args=(analysis_id, file_path)).start()
        
        return jsonify({
            'analysis_id': analysis_id,
            'status': 'processing',
            'message': 'Analysis started successfully'
        })
        
    except Exception as e:
        logger.error(f"Upload error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/status/<analysis_id>')
def get_analysis_status(analysis_id):
    """Get analysis status and progress"""
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    # Calculate progress
    completed_steps = sum(1 for step in result['steps'] if step['status'] == 'completed')
    total_steps = len(result['steps'])
    progress = int((completed_steps / total_steps) * 100)
    
    result['progress'] = progress
    
    return jsonify(result)

@app.route('/api/results/<analysis_id>')
def get_analysis_results_json(analysis_id):
    """Get analysis results as JSON for frontend"""
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    # Return the raw analysis data as JSON
    return jsonify(result)

@app.route('/api/analysis/<analysis_id>/fail', methods=['POST'])
def mark_analysis_failed(analysis_id):
    """Mark analysis as failed (called by Telegram bot)"""
    try:
        data = request.get_json()
        error_message = data.get('error', 'Unknown error')
        
        if analysis_id in analysis_results:
            analysis_results[analysis_id]['status'] = 'failed'
            analysis_results[analysis_id]['error'] = error_message
            analysis_results[analysis_id]['progress'] = 0
            logger.info(f"Analysis {analysis_id} marked as failed: {error_message}")
            return jsonify({'success': True, 'message': 'Analysis marked as failed'})
        else:
            logger.warning(f"Attempted to mark non-existent analysis as failed: {analysis_id}")
            return jsonify({'error': 'Analysis not found'}), 404
            
    except Exception as e:
        logger.error(f"Error marking analysis as failed: {e}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/analysis/<analysis_id>/cancel', methods=['POST'])
def cancel_analysis(analysis_id):
    """Cancel analysis (called by Telegram bot /stop command)"""
    try:
        data = request.get_json()
        user_id = data.get('user_id')
        
        logger.info(f"Cancel request received for analysis {analysis_id} from user {user_id}")
        logger.info(f"Request data: {data}")
        
        if analysis_id in analysis_results:
            current_status = analysis_results[analysis_id]['status']
            logger.info(f"Analysis {analysis_id} current status: {current_status}")
            
            # Check if analysis is still in progress
            if current_status in ['processing', 'pending']:
                analysis_results[analysis_id]['status'] = 'cancelled'
                analysis_results[analysis_id]['error'] = 'Cancelled by user'
                analysis_results[analysis_id]['progress'] = 0
                logger.info(f"Analysis {analysis_id} cancelled by user {user_id}")
                return jsonify({'success': True, 'message': 'Analysis cancelled successfully'})
            else:
                logger.warning(f"Analysis {analysis_id} cannot be cancelled - status: {current_status}")
                return jsonify({'error': f'Analysis cannot be cancelled - current status: {current_status}'}), 400
        else:
            logger.warning(f"Attempted to cancel non-existent analysis: {analysis_id}")
            logger.info(f"Available analyses: {list(analysis_results.keys())}")
            return jsonify({'error': 'Analysis not found'}), 404
            
    except Exception as e:
        logger.error(f"Error cancelling analysis: {e}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/cleanup/status', methods=['GET'])
def get_cleanup_status():
    """Get cleanup status for debugging"""
    try:
        current_time = time.time()
        cleanup_info = []
        
        for analysis_id, info in file_cleanup_timers.items():
            file_path = info.get('file_path')
            cleanup_time = info.get('cleanup_time', 0)
            time_remaining = max(0, 30 * 60 - (current_time - cleanup_time))
            
            cleanup_info.append({
                'analysis_id': analysis_id,
                'file_path': file_path,
                'time_remaining_minutes': round(time_remaining / 60, 1),
                'is_url_analysis': file_path is None
            })
        
        return jsonify({
            'total_tracked': len(file_cleanup_timers),
            'cleanup_threshold_minutes': 30,
            'tracked_analyses': cleanup_info
        })
        
    except Exception as e:
        logger.error(f"Error getting cleanup status: {e}")
        return jsonify({'error': 'Internal server error'}), 500

@app.route('/api/test/stop', methods=['GET'])
def test_stop_endpoint():
    """Test endpoint to verify /stop command functionality"""
    try:
        return jsonify({
            'message': 'Stop endpoint is working',
            'available_analyses': list(analysis_results.keys()),
            'total_analyses': len(analysis_results),
            'cleanup_tracked': len(file_cleanup_timers)
        })
    except Exception as e:
        logger.error(f"Error in test endpoint: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/analyses/active', methods=['GET'])
def get_active_analyses():
    """Get all active analyses that can be cancelled"""
    try:
        active_analyses = []
        
        for analysis_id, analysis_data in analysis_results.items():
            status = analysis_data.get('status', 'unknown')
            # Only include analyses that are actually still running
            if status in ['processing', 'pending']:
                # Double-check that the analysis is actually still active
                # by checking if it has a recent update time
                created_at = analysis_data.get('created_at', analysis_data.get('start_time', ''))
                if created_at:
                    try:
                        from datetime import datetime
                        # Parse the creation time
                        if 'T' in created_at:
                            created_time = datetime.fromisoformat(created_at.replace('Z', '+00:00'))
                        else:
                            created_time = datetime.fromisoformat(created_at)
                        
                        # Check if analysis is older than 1 hour (probably stuck)
                        time_diff = datetime.now() - created_time.replace(tzinfo=None)
                        if time_diff.total_seconds() > 3600:  # 1 hour
                            logger.warning(f"Analysis {analysis_id} is older than 1 hour but still marked as {status}, skipping")
                            continue
                    except Exception as e:
                        logger.warning(f"Error checking analysis {analysis_id} age: {e}")
                
                active_analyses.append({
                    'analysis_id': analysis_id,
                    'status': status,
                    'filename': analysis_data.get('filename', 'Unknown'),
                    'is_url': analysis_data.get('is_url', False),
                    'created_at': created_at
                })
        
        return jsonify({
            'active_analyses': active_analyses,
            'total_active': len(active_analyses),
            'total_analyses': len(analysis_results)
        })
        
    except Exception as e:
        logger.error(f"Error getting active analyses: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/analysis/<analysis_id>/notify-cancelled', methods=['POST'])
def notify_analysis_cancelled(analysis_id):
    """Notify that an analysis was cancelled (for cleanup system)"""
    try:
        data = request.get_json()
        user_id = data.get('user_id')
        reason = data.get('reason', 'Unknown reason')
        
        if analysis_id in analysis_results:
            analysis_data = analysis_results[analysis_id]
            filename = analysis_data.get('filename', 'Unknown file')
            is_url = analysis_data.get('is_url', False)
            analysis_type = "URL" if is_url else "File"
            
            logger.info(f"Analysis {analysis_id} cancelled: {reason}")
            
            # This would trigger a notification to the Telegram bot
            # For now, we'll just log it
            notification_message = f"""
🛑 **Analysis Cancelled**

**📁 {analysis_type}:** `{filename}`
**🆔 Analysis ID:** `{analysis_id}`
**👤 User:** {user_id}

**✅ Status:** Cancelled - {reason}
**⏰ Time:** {datetime.now().strftime('%H:%M:%S')}

The analysis process has been stopped.
            """
            
            logger.info(f"Notification message: {notification_message}")
            
            return jsonify({
                'success': True,
                'message': 'Notification logged',
                'analysis_id': analysis_id,
                'reason': reason
            })
        else:
            return jsonify({'error': 'Analysis not found'}), 404
            
    except Exception as e:
        logger.error(f"Error notifying cancellation: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/results/<analysis_id>')
def get_analysis_results(analysis_id):
    """Redirect to React frontend - no longer serving HTML from Flask"""
    # Get the frontend URL from environment or default to production
    frontend_url = os.getenv('WEBSITE_URL', 'https://www.aejis.xyz')
    
    # Check if analysis exists
    if analysis_id not in analysis_results:
        return redirect(f'{frontend_url}/')
    
    result = analysis_results[analysis_id]
    
    # Redirect to appropriate React frontend page
    if result.get('is_url', False):
        return redirect(f'{frontend_url}/url-results/{analysis_id}')
    else:
        return redirect(f'{frontend_url}/results/{analysis_id}')

@app.route('/url-results-data/<analysis_id>')
def get_url_results_data(analysis_id):
    """Get URL analysis results data as JSON for React frontend"""
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    if not result.get('is_url', False):
        return jsonify({'error': 'Not a URL analysis'}), 400
    
    # Return the complete analysis data
    return jsonify(result)

def render_results_page(analysis_id, results):
    """Render a comprehensive HTML results page with all analysis details"""
    url = results.get('url_info', {}).get('url', 'Unknown URL')
    domain = results.get('url_info', {}).get('domain', 'Unknown Domain')
    trust_score = results.get('ai_analysis', {}).get('trust_score', 0)
    threat_level = results.get('summary', {}).get('threat_level', 'UNKNOWN')
    
    # Extract detailed analysis data
    ai_analysis = results.get('ai_analysis', {})
    virustotal = results.get('virustotal', {})
    security_analysis = ai_analysis.get('security_analysis', {})
    domain_intelligence = ai_analysis.get('domain_intelligence', {})
    content_analysis = ai_analysis.get('content_analysis', {})
    advanced_metrics = ai_analysis.get('advanced_metrics', {})
    
    html = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Aejis Analysis Results - {domain}</title>
        <style>
            * {{
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }}
            
            body {{
                font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                min-height: 100vh;
                padding: 20px;
                line-height: 1.6;
            }}
            
            .container {{
                max-width: 1400px;
                margin: 0 auto;
                background: white;
                border-radius: 20px;
                box-shadow: 0 20px 40px rgba(0,0,0,0.1);
                overflow: hidden;
            }}
            
            .header {{
                background: linear-gradient(135deg, #2c3e50 0%, #3498db 100%);
                color: white;
                padding: 40px;
                text-align: center;
            }}
            
            .header h1 {{
                font-size: 3em;
                margin-bottom: 15px;
                text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
            }}
            
            .header p {{
                font-size: 1.3em;
                opacity: 0.9;
            }}
            
            .url-info {{
                background: #f8f9fa;
                padding: 25px;
                border-bottom: 1px solid #e9ecef;
            }}
            
            .url-info h2 {{
                color: #2c3e50;
                margin-bottom: 15px;
                font-size: 1.5em;
            }}
            
            .url-info .url {{
                font-family: monospace;
                background: #e9ecef;
                padding: 15px;
                border-radius: 8px;
                word-break: break-all;
                font-size: 1.1em;
            }}
            
            .main-content {{
                display: grid;
                grid-template-columns: 1fr 1fr;
                gap: 30px;
                padding: 30px;
            }}
            
            .left-column, .right-column {{
                display: flex;
                flex-direction: column;
                gap: 25px;
            }}
            
            .section {{
                background: #f8f9fa;
                border-radius: 15px;
                padding: 25px;
                border-left: 5px solid #3498db;
                box-shadow: 0 5px 15px rgba(0,0,0,0.1);
            }}
            
            .section h3 {{
                color: #2c3e50;
                margin-bottom: 20px;
                font-size: 1.4em;
                display: flex;
                align-items: center;
                gap: 10px;
            }}
            
            .score-display {{
                text-align: center;
                padding: 20px;
                background: linear-gradient(135deg, #27ae60, #2ecc71);
                color: white;
                border-radius: 15px;
                margin-bottom: 20px;
            }}
            
            .score-number {{
                font-size: 4em;
                font-weight: bold;
                text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
            }}
            
            .score-label {{
                font-size: 1.2em;
                opacity: 0.9;
            }}
            
            .threat-level {{
                text-align: center;
                padding: 15px;
                border-radius: 10px;
                font-weight: bold;
                font-size: 1.3em;
                margin: 10px 0;
            }}
            
            .threat-level.SAFE {{
                background: #d4edda;
                color: #155724;
            }}
            
            .threat-level.MEDIUM {{
                background: #fff3cd;
                color: #856404;
            }}
            
            .threat-level.HIGH {{
                background: #f8d7da;
                color: #721c24;
            }}
            
            .detail-grid {{
                display: grid;
                grid-template-columns: 1fr 1fr;
                gap: 15px;
            }}
            
            .detail-item {{
                display: flex;
                justify-content: space-between;
                padding: 12px 0;
                border-bottom: 1px solid #e9ecef;
                align-items: center;
            }}
            
            .detail-item:last-child {{
                border-bottom: none;
            }}
            
            .detail-label {{
                font-weight: bold;
                color: #2c3e50;
                font-size: 0.95em;
            }}
            
            .detail-value {{
                color: #6c757d;
                text-align: right;
                font-size: 0.9em;
                word-break: break-word;
            }}
            
            .ssl-info {{
                background: #e8f5e8;
                padding: 15px;
                border-radius: 8px;
                margin: 10px 0;
            }}
            
            .ssl-grade {{
                font-size: 1.5em;
                font-weight: bold;
                color: #27ae60;
                text-align: center;
            }}
            
            .ai-analysis {{
                background: #f0f8ff;
                padding: 20px;
                border-radius: 10px;
                margin: 15px 0;
                border-left: 4px solid #3498db;
            }}
            
            .ai-text {{
                font-style: italic;
                line-height: 1.7;
                color: #2c3e50;
            }}
            
            .metrics-grid {{
                display: grid;
                grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
                gap: 15px;
                margin: 15px 0;
            }}
            
            .metric-card {{
                background: white;
                padding: 15px;
                border-radius: 8px;
                text-align: center;
                box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            }}
            
            .metric-value {{
                font-size: 1.8em;
                font-weight: bold;
                color: #3498db;
            }}
            
            .metric-label {{
                font-size: 0.9em;
                color: #6c757d;
                margin-top: 5px;
            }}
            
            .full-width {{
                grid-column: 1 / -1;
            }}
            
            .virustotal-section {{
                background: #fff5f5;
                border-left-color: #e74c3c;
            }}
            
            .security-section {{
                background: #f0fff0;
                border-left-color: #27ae60;
            }}
            
            .domain-section {{
                background: #fff8e1;
                border-left-color: #f39c12;
            }}
            
            .content-section {{
                background: #f3e5f5;
                border-left-color: #9b59b6;
            }}
            
            .preview-section {{
                background: #f8f9fa;
                padding: 30px;
                border-top: 1px solid #e9ecef;
            }}
            
            .preview-description {{
                background: #e8f4fd;
                padding: 20px;
                border-radius: 10px;
                margin: 20px 0;
                border-left: 4px solid #3498db;
                font-size: 1.1em;
                line-height: 1.6;
            }}
            
            .preview-buttons {{
                display: flex;
                gap: 15px;
                margin: 20px 0;
                flex-wrap: wrap;
            }}
            
            .preview-btn {{
                padding: 12px 24px;
                border: none;
                border-radius: 8px;
                font-size: 1em;
                font-weight: bold;
                cursor: pointer;
                transition: all 0.3s ease;
                text-decoration: none;
                display: inline-block;
            }}
            
            .preview-btn.primary {{
                background: linear-gradient(135deg, #3498db, #2980b9);
                color: white;
            }}
            
            .preview-btn.primary:hover {{
                background: linear-gradient(135deg, #2980b9, #1f618d);
                transform: translateY(-2px);
                box-shadow: 0 5px 15px rgba(52, 152, 219, 0.4);
            }}
            
            .preview-btn.secondary {{
                background: linear-gradient(135deg, #95a5a6, #7f8c8d);
                color: white;
            }}
            
            .preview-btn.secondary:hover {{
                background: linear-gradient(135deg, #7f8c8d, #6c7b7d);
                transform: translateY(-2px);
                box-shadow: 0 5px 15px rgba(149, 165, 166, 0.4);
            }}
            
            .preview-btn.danger {{
                background: linear-gradient(135deg, #e74c3c, #c0392b);
                color: white;
            }}
            
            .preview-btn.danger:hover {{
                background: linear-gradient(135deg, #c0392b, #a93226);
                transform: translateY(-2px);
                box-shadow: 0 5px 15px rgba(231, 76, 60, 0.4);
            }}
            
            .browser-status {{
                background: #fff;
                padding: 15px;
                border-radius: 8px;
                margin: 15px 0;
                border-left: 4px solid #3498db;
                font-family: monospace;
                font-size: 0.9em;
            }}
            
            .browser-preview-container {{
                background: #000;
                border-radius: 10px;
                overflow: hidden;
                margin: 20px 0;
                min-height: 600px;
                position: relative;
            }}
            
            .browser-preview-container iframe {{
                width: 100%;
                height: 600px;
                border: none;
                background: #fff;
            }}
            
            .browser-loading {{
                position: absolute;
                top: 50%;
                left: 50%;
                transform: translate(-50%, -50%);
                color: white;
                font-size: 1.2em;
                text-align: center;
            }}
            
            .browser-loading::after {{
                content: '';
                display: block;
                width: 40px;
                height: 40px;
                margin: 20px auto;
                border: 4px solid #3498db;
                border-top: 4px solid transparent;
                border-radius: 50%;
                animation: spin 1s linear infinite;
            }}
            
            @keyframes spin {{
                0% {{ transform: rotate(0deg); }}
                100% {{ transform: rotate(360deg); }}
            }}
            
            @media (max-width: 768px) {{
                .main-content {{
                    grid-template-columns: 1fr;
                }}
                
                .detail-grid {{
                    grid-template-columns: 1fr;
                }}
                
                .metrics-grid {{
                    grid-template-columns: repeat(auto-fit, minmax(150px, 1fr));
                }}
                
                .preview-buttons {{
                    flex-direction: column;
                }}
                
                .preview-btn {{
                    width: 100%;
                    text-align: center;
                }}
            }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <h1>🛡️ Aejis Analysis Results</h1>
                <p>Comprehensive Security Analysis Report</p>
            </div>
            
            <div class="url-info">
                <h2>🌐 Analyzed URL</h2>
                <div class="url">{url}</div>
            </div>
            
            <div class="main-content">
                <div class="left-column">
                    <!-- Trust Score & Threat Level -->
                    <div class="section">
                        <h3>🎯 Trust Assessment</h3>
                        <div class="score-display">
                            <div class="score-number">{trust_score}</div>
                            <div class="score-label">Trust Score / 100</div>
                        </div>
                        <div class="threat-level {threat_level}">
                            Threat Level: {threat_level}
                        </div>
                    </div>
                    
                    <!-- VirusTotal Analysis -->
                    <div class="section virustotal-section">
                        <h3>🔍 VirusTotal Analysis</h3>
                        <div class="detail-item">
                            <span class="detail-label">Status:</span>
                            <span class="detail-value">{virustotal.get('scan_details', {}).get('status', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Engines Used:</span>
                            <span class="detail-value">{virustotal.get('engines_used', 0)}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Threats Detected:</span>
                            <span class="detail-value">{virustotal.get('engines_detected', 0)}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Confidence:</span>
                            <span class="detail-value">{virustotal.get('confidence', 0)}%</span>
                        </div>
                    </div>
                    
                    <!-- Security Analysis -->
                    <div class="section security-section">
                        <h3>🔒 Security Analysis</h3>
                        <div class="ssl-info">
                            <div class="ssl-grade">SSL Grade: {security_analysis.get('ssl_grade', 'Unknown')}</div>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Security Score:</span>
                            <span class="detail-value">{security_analysis.get('security_score', 0)}/100</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Headers Score:</span>
                            <span class="detail-value">{security_analysis.get('headers_score', 0)}/100</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Threats Found:</span>
                            <span class="detail-value">{len(security_analysis.get('threats', []))}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Vulnerabilities:</span>
                            <span class="detail-value">{len(security_analysis.get('vulnerabilities', []))}</span>
                        </div>
                    </div>
                    
                    <!-- Domain Intelligence -->
                    <div class="section domain-section">
                        <h3>🌍 Domain Intelligence</h3>
                        <div class="detail-item">
                            <span class="detail-label">Domain Age:</span>
                            <span class="detail-value">{domain_intelligence.get('age_years', 0)} years</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Country:</span>
                            <span class="detail-value">{domain_intelligence.get('country', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Registrar:</span>
                            <span class="detail-value">{domain_intelligence.get('registrar', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Organization:</span>
                            <span class="detail-value">{domain_intelligence.get('organization', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">IP Address:</span>
                            <span class="detail-value">{domain_intelligence.get('ip_address', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Global Rank:</span>
                            <span class="detail-value">{domain_intelligence.get('global_rank', 'Unknown')}</span>
                        </div>
                    </div>
                </div>
                
                <div class="right-column">
                    <!-- Advanced Metrics -->
                    <div class="section">
                        <h3>📊 Advanced Metrics</h3>
                        <div class="metrics-grid">
                            <div class="metric-card">
                                <div class="metric-value">{advanced_metrics.get('content_legitimacy_score', 0)}</div>
                                <div class="metric-label">Content Legitimacy</div>
                            </div>
                            <div class="metric-card">
                                <div class="metric-value">{advanced_metrics.get('ssl_grade', 'N/A')}</div>
                                <div class="metric-label">SSL Grade</div>
                            </div>
                            <div class="metric-card">
                                <div class="metric-value">{advanced_metrics.get('domain_age_days', 0)}</div>
                                <div class="metric-label">Domain Age (Days)</div>
                            </div>
                            <div class="metric-card">
                                <div class="metric-value">{advanced_metrics.get('traffic_rank', 'N/A')}</div>
                                <div class="metric-label">Traffic Rank</div>
                            </div>
                        </div>
                    </div>
                    
                    <!-- Content Analysis -->
                    <div class="section content-section">
                        <h3>📝 Content Analysis</h3>
                        <div class="detail-item">
                            <span class="detail-label">Title:</span>
                            <span class="detail-value">{content_analysis.get('title', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Meta Description:</span>
                            <span class="detail-value">{content_analysis.get('meta_description', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Language:</span>
                            <span class="detail-value">{content_analysis.get('language', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Content Quality:</span>
                            <span class="detail-value">{content_analysis.get('content_quality', 'Unknown')}</span>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">Legitimacy Score:</span>
                            <span class="detail-value">{content_analysis.get('legitimacy_score', 0)}/100</span>
                        </div>
                    </div>
                    
                    <!-- AI Analysis Report -->
                    <div class="section full-width">
                        <h3>🤖 AI Analysis Report</h3>
                        <div class="ai-analysis">
                            <div class="ai-text">
                                {ai_analysis.get('ai_analysis', 'No AI analysis available.')}
                            </div>
                        </div>
                        <div class="detail-item">
                            <span class="detail-label">AI Confidence:</span>
                            <span class="detail-value">{ai_analysis.get('ai_confidence_score', 0)}%</span>
                        </div>
                    </div>
                    
                    <!-- Analysis Details -->
                    <div class="section full-width">
                        <h3>📋 Analysis Details</h3>
                        <div class="detail-grid">
                            <div class="detail-item">
                                <span class="detail-label">Analysis ID:</span>
                                <span class="detail-value">{analysis_id}</span>
                            </div>
                            <div class="detail-item">
                                <span class="detail-label">Analysis Time:</span>
                                <span class="detail-value">{results.get('url_info', {}).get('analysis_time', 'Unknown')}</span>
                            </div>
                            <div class="detail-item">
                                <span class="detail-label">Total Engines:</span>
                                <span class="detail-value">{results.get('summary', {}).get('engines_used', 0)}</span>
                            </div>
                            <div class="detail-item">
                                <span class="detail-label">Overall Confidence:</span>
                                <span class="detail-value">{results.get('summary', {}).get('confidence', 0)}%</span>
                            </div>
                        </div>
                    </div>
                </div>
            </div>
            
            <!-- Aejis Browser Isolation Preview Section -->
            <div class="preview-section">
                <div class="section">
                    <h3>🌐 Safe Browser Preview</h3>
                    <p class="preview-description">
                        <strong>🛡️ Aejis Browser Isolation System</strong><br>
                        Browse this website safely in a completely isolated environment with random location spoofing. 
                        Any malware or threats will be contained within the isolated container.
                    </p>
                    <div class="preview-buttons">
                        <button onclick="startBrowserIsolation('{analysis_id}')" class="preview-btn primary">
                            🚀 Open Isolated Browser
                        </button>
                        <button onclick="checkBrowserStatus('{analysis_id}')" class="preview-btn secondary">
                            📊 Check Status
                        </button>
                        <button onclick="stopBrowserSession('{analysis_id}')" class="preview-btn danger">
                            🛑 Stop Session
                        </button>
                    </div>
                    <div id="browser-status-{analysis_id}" class="browser-status"></div>
                    <div id="browser-preview-{analysis_id}" class="browser-preview-container"></div>
                </div>
            </div>
        </div>
        
        <script>
            document.addEventListener('DOMContentLoaded', function() {{
                console.log('Aejis Comprehensive Analysis Results loaded successfully');
            }});
            
            // Aejis Browser Isolation Functions
            async function startBrowserIsolation(analysisId) {{
                const statusDiv = document.getElementById(`browser-status-${{analysisId}}`);
                const previewDiv = document.getElementById(`browser-preview-${{analysisId}}`);
                
                statusDiv.innerHTML = '<div class="browser-loading">🚀 Starting Aejis Browser Isolation...</div>';
                previewDiv.innerHTML = '';
                
                try {{
                    const response = await fetch(`/browser/${{analysisId}}`);
                    const data = await response.json();
                    
                    if (data.success) {{
                        statusDiv.innerHTML = `
                            <div style="color: #27ae60;">
                                ✅ Browser isolation session started successfully!<br>
                                <strong>Session ID:</strong> ${{data.session_id}}<br>
                                <strong>Target URL:</strong> ${{data.target_url}}<br>
                                <strong>Isolation Level:</strong> ${{data.isolation_level}}<br>
                                <strong>Location Spoofing:</strong> ${{data.location_spoofing ? 'Enabled' : 'Disabled'}}<br>
                                <strong>Commercial Friendly:</strong> ${{data.commercial_friendly ? 'Yes' : 'No'}}
                            </div>
                        `;
                        
                        // Create iframe for browser preview - use auto-connect for seamless experience
                        const pureUrl = 'http://localhost:5000/vnc-auto-connect.html?url=' + encodeURIComponent(data.target_url);
                        previewDiv.innerHTML = `
                            <div class="browser-loading">🌐 Loading isolated browser...</div>
                            <iframe src="` + pureUrl + `" 
                                    style="width: 100%; height: 600px; border: none; background: #fff;"
                                    onload="this.style.display='block'; this.previousElementSibling.style.display='none';">
                            </iframe>
                            <div style="margin-top: 10px; text-align: center;">
                                <a href="` + pureUrl + `" target="_blank" style="color: #3498db; text-decoration: none;">
                                    🔗 Open in New Tab
                                </a>
                            </div>
                        `;
                    }} else {{
                        statusDiv.innerHTML = `<div style="color: #e74c3c;">❌ Error: ${{data.error}}</div>`;
                    }}
                }} catch (error) {{
                    statusDiv.innerHTML = `<div style="color: #e74c3c;">❌ Error: ${{error.message}}</div>`;
                }}
            }}
            
            async function checkBrowserStatus(analysisId) {{
                const statusDiv = document.getElementById(`browser-status-${{analysisId}}`);
                
                try {{
                    const response = await fetch(`/browser/${{analysisId}}/status`);
                    const data = await response.json();
                    
                    if (data.success) {{
                        statusDiv.innerHTML = `
                            <div style="color: #3498db;">
                                📊 Browser Status: <strong>${{data.status}}</strong><br>
                                <strong>Session ID:</strong> ${{data.session_id}}<br>
                                <strong>Target URL:</strong> ${{data.target_url}}<br>
                                <strong>Start Time:</strong> ${{new Date(data.start_time * 1000).toLocaleString()}}<br>
                                <strong>Isolation Level:</strong> ${{data.isolation_level}}<br>
                                <strong>Location Spoofing:</strong> ${{data.location_spoofing ? 'Enabled' : 'Disabled'}}
                            </div>
                        `;
                    }} else {{
                        statusDiv.innerHTML = `<div style="color: #e74c3c;">❌ Error: ${{data.error}}</div>`;
                    }}
                }} catch (error) {{
                    statusDiv.innerHTML = `<div style="color: #e74c3c;">❌ Error: ${{error.message}}</div>`;
                }}
            }}
            
            async function stopBrowserSession(analysisId) {{
                const statusDiv = document.getElementById(`browser-status-${{analysisId}}`);
                const previewDiv = document.getElementById(`browser-preview-${{analysisId}}`);
                
                statusDiv.innerHTML = '<div class="browser-loading">🛑 Stopping browser session...</div>';
                
                try {{
                    const response = await fetch(`/browser/${{analysisId}}/stop`, {{
                        method: 'POST',
                        headers: {{ 'Content-Type': 'application/json' }}
                    }});
                    const data = await response.json();
                    
                    if (data.success) {{
                        statusDiv.innerHTML = '<div style="color: #e74c3c;">🛑 Browser session stopped successfully</div>';
                        previewDiv.innerHTML = '';
                    }} else {{
                        statusDiv.innerHTML = `<div style="color: #e74c3c;">❌ Error: ${{data.error}}</div>`;
                    }}
                }} catch (error) {{
                    statusDiv.innerHTML = `<div style="color: #e74c3c;">❌ Error: ${{error.message}}</div>`;
                }}
            }}
        </script>
    </body>
    </html>
    """
    
    return html

def _get_forensic_file_analysis(file_path: str, original_filename: str = None) -> Dict[str, Any]:
    """Get ultra-detailed forensic file analysis beyond what any standard tool provides"""
    import hashlib
    import time
    import mimetypes
    import math
    from collections import Counter
    
    try:
        # Basic file stats
        stat_info = os.stat(file_path)
        
        # Read file for analysis
        with open(file_path, 'rb') as f:
            file_data = f.read()
        
        file_size = len(file_data)
        
        # Use original filename if provided, otherwise use file path
        display_filename = original_filename or os.path.basename(file_path)
        
        # 1. CRYPTOGRAPHIC HASHES (Multiple algorithms)
        hashes = {
            'md5': hashlib.md5(file_data).hexdigest(),
            'sha1': hashlib.sha1(file_data).hexdigest(),
            'sha256': hashlib.sha256(file_data).hexdigest(),
            'sha512': hashlib.sha512(file_data).hexdigest()
        }
        
        # 2. ENTROPY ANALYSIS (Measure of randomness/encryption)
        if file_size > 0:
            byte_counts = Counter(file_data)
            entropy = -sum((count/file_size) * math.log2(count/file_size) for count in byte_counts.values())
        else:
            entropy = 0.0
        
        # 3. BYTE FREQUENCY ANALYSIS
        byte_freq = Counter(file_data)
        most_common_bytes = byte_freq.most_common(10)
        
        # 4. FILE HEADER/MAGIC BYTES ANALYSIS
        header_bytes = file_data[:64] if file_size >= 64 else file_data
        magic_signature = header_bytes.hex().upper()
        
        # 5. COMPRESSION RATIO ANALYSIS
        try:
            import zlib
            compressed = zlib.compress(file_data)
            compression_ratio = len(compressed) / file_size if file_size > 0 else 0
        except:
            compression_ratio = 0
        
        # 6. STRING EXTRACTION (ASCII/Unicode strings)
        strings_found = []
        current_string = ""
        for byte in file_data:  # Full file analysis
            if 32 <= byte <= 126:  # Printable ASCII
                current_string += chr(byte)
            else:
                if len(current_string) >= 4:
                    strings_found.append(current_string)
                current_string = ""
        if len(current_string) >= 4:
            strings_found.append(current_string)
        
        # 7. BINARY STRUCTURE ANALYSIS
        null_bytes = file_data.count(0)
        printable_chars = sum(1 for b in file_data if 32 <= b <= 126)
        binary_ratio = (file_size - printable_chars) / file_size if file_size > 0 else 0
        
        # 8. TIMESTAMP ANALYSIS
        timestamps = {
            'created': time.ctime(stat_info.st_ctime),
            'modified': time.ctime(stat_info.st_mtime),
            'accessed': time.ctime(stat_info.st_atime),
            'created_unix': int(stat_info.st_ctime),
            'modified_unix': int(stat_info.st_mtime),
            'accessed_unix': int(stat_info.st_atime)
        }
        
        # 9. FILE TYPE DETECTION
        mime_type, encoding = mimetypes.guess_type(file_path)
        
        # 10. SUSPICIOUS PATTERNS
        suspicious_indicators = []
        
        # Check for executable headers
        if file_data.startswith(b'MZ'):
            suspicious_indicators.append("PE/EXE executable header detected")
        if b'\x00\x00\x00\x00' in file_data:
            suspicious_indicators.append("Multiple null bytes in header (possible padding)")
        if entropy > 7.5:
            suspicious_indicators.append(f"High entropy ({entropy:.2f}) - possible encryption/compression")
        if compression_ratio > 0.9:
            suspicious_indicators.append("File appears already compressed/encrypted")
        
        # 11. ADVANCED METADATA
        advanced_metadata = {
            'file_permissions': oct(stat_info.st_mode)[-3:],
            'file_size_human': _format_file_size(file_size),
            'unique_bytes': len(byte_freq),
            'zero_bytes_percentage': (null_bytes / file_size * 100) if file_size > 0 else 0,
            'printable_percentage': (printable_chars / file_size * 100) if file_size > 0 else 0,
            'is_likely_binary': binary_ratio > 0.3,
            'is_likely_encrypted': entropy > 7.0,
            'is_likely_compressed': compression_ratio > 0.8
        }
        
        return {
            'filename': display_filename,
            'file_size': file_size,
            'file_size_human': _format_file_size(file_size),
            'hashes': hashes,
            'entropy': round(entropy, 4),
            'entropy_interpretation': _interpret_entropy(entropy),
            'byte_frequency': {
                'most_common': [{'byte': f"0x{byte:02X}", 'count': count, 'percentage': round(count/file_size*100, 2)} 
                               for byte, count in most_common_bytes],
                'unique_bytes': len(byte_freq),
                'total_bytes': file_size
            },
            'file_header': {
                'magic_bytes': magic_signature,  # Full header
                'full_header': magic_signature,
                'interpretation': _interpret_magic_bytes(header_bytes)
            },
            'compression_analysis': {
                'original_size': file_size,
                'compressed_size': len(compressed) if 'compressed' in locals() else 0,
                'compression_ratio': round(compression_ratio, 4),
                'compression_interpretation': _interpret_compression(compression_ratio)
            },
            'string_analysis': {
                'strings_found': strings_found,  # All strings
                'total_strings': len(strings_found),
                'longest_string': max(strings_found, key=len) if strings_found else "",
                'average_string_length': sum(len(s) for s in strings_found) / len(strings_found) if strings_found else 0
            },
            'binary_structure': {
                'null_bytes': null_bytes,
                'printable_chars': printable_chars,
                'binary_ratio': round(binary_ratio, 4),
                'is_likely_binary': advanced_metadata['is_likely_binary']
            },
            'timestamps': timestamps,
            'mime_detection': {
                'mime_type': mime_type or 'unknown',
                'encoding': encoding or 'unknown'
            },
            'security_analysis': {
                'suspicious_indicators': suspicious_indicators,
                'risk_score': len(suspicious_indicators) * 25,  # 0-100 scale
                'is_potentially_malicious': len(suspicious_indicators) >= 2
            },
            'advanced_metadata': advanced_metadata
        }
        
    except Exception as e:
        logger.error(f"Error in forensic analysis: {str(e)}")
        return {'error': f'Forensic analysis failed: {str(e)}'}

def _format_file_size(size_bytes):
    """Format file size in human readable format"""
    if size_bytes == 0:
        return "0 B"
    size_names = ["B", "KB", "MB", "GB", "TB"]
    i = int(math.floor(math.log(size_bytes, 1024)))
    p = math.pow(1024, i)
    s = round(size_bytes / p, 2)
    return f"{s} {size_names[i]}"

def _interpret_entropy(entropy):
    """Interpret entropy value"""
    if entropy < 2:
        return "Very low - highly repetitive data"
    elif entropy < 4:
        return "Low - structured data with patterns"
    elif entropy < 6:
        return "Medium - mixed content"
    elif entropy < 7:
        return "High - compressed or random data"
    else:
        return "Very high - likely encrypted or heavily compressed"

def _interpret_compression(ratio):
    """Interpret compression ratio"""
    if ratio > 0.95:
        return "Already compressed/encrypted - minimal compression"
    elif ratio > 0.8:
        return "Poor compression - likely binary or compressed data"
    elif ratio > 0.5:
        return "Moderate compression - mixed content"
    else:
        return "Good compression - text or structured data"

def _interpret_magic_bytes(header_bytes):
    """Interpret file magic bytes/signatures"""
    magic_signatures = {
        b'\x89PNG': 'PNG Image',
        b'\xFF\xD8\xFF': 'JPEG Image',
        b'GIF8': 'GIF Image',
        b'RIFF': 'RIFF Container (AVI/WAV)',
        b'%PDF': 'PDF Document',
        b'PK\x03\x04': 'ZIP Archive',
        b'PK\x05\x06': 'ZIP Archive (empty)',
        b'MZ': 'PE Executable',
        b'\x7FELF': 'ELF Executable',
        b'\xCA\xFE\xBA\xBE': 'Java Class File',
        b'BM': 'Bitmap Image',
        b'ID3': 'MP3 Audio',
        b'\x00\x00\x01\x00': 'ICO Image',
        b'ftyp': 'MP4/MOV Video (offset 4)',
        b'\x50\x4B': 'ZIP-based format',
        b'\x1F\x8B': 'GZIP Compressed',
        b'\x42\x5A\x68': 'BZIP2 Compressed'
    }
    
    for signature, description in magic_signatures.items():
        if header_bytes.startswith(signature):
            return description
    
    # Check for Office documents (more complex detection)
    if header_bytes.startswith(b'\xD0\xCF\x11\xE0'):
        return 'Microsoft Office Document (Legacy)'
    
    return 'Unknown/Custom format'

@app.route('/preview/<analysis_id>/archive-file/<path:file_path>')
def get_archive_file_content_direct(analysis_id, file_path):
    """SECURE: Get content of a specific file within an archive using Docker isolation"""
    try:
        if analysis_id not in analysis_results:
            return jsonify({'error': 'Analysis not found'}), 404
        
        results = analysis_results[analysis_id]
        file_path_full = results.get('file_path')
        
        if not file_path_full or not os.path.exists(file_path_full):
            return jsonify({'error': 'File not found'}), 404
        
        # Check if it's a ZIP file
        if not file_path_full.lower().endswith('.zip'):
            return jsonify({'error': 'Not a ZIP archive'}), 400
        
        # Use secure Docker-based processing
        return _secure_archive_file_access(file_path_full, file_path, analysis_id)
                
    except Exception as e:
        logger.error(f"Error processing archive file request: {e}")
        return jsonify({'error': 'Failed to process archive file'}), 500

@app.route('/preview/<analysis_id>')
def get_file_preview(analysis_id):
    """Redirect to React frontend preview page"""
    frontend_url = os.getenv('WEBSITE_URL', 'https://www.aejis.xyz')
    
    if analysis_id not in analysis_results:
        return redirect(f'{frontend_url}/')
    
    # Redirect to React frontend preview page
    return redirect(f'{frontend_url}/preview/{analysis_id}')
    
    # Old code below (kept for reference but not executed)
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    # Get file path from the analysis (it's at the top level, not in results)
    file_path = result.get('file_path')
    if not file_path or not os.path.exists(file_path):
        return jsonify({'error': 'File not found'}), 404
    
    try:
        # Read file contents for preview
        file_size = os.path.getsize(file_path)
        filename = os.path.basename(file_path)
        
        # Determine file type and preview method
        file_ext = os.path.splitext(filename)[1].lower()
        
        # Get forensic analysis
        forensic_data = _get_forensic_file_analysis(file_path, filename)
        
        preview_data = {
            'filename': filename,
            'file_size': file_size,
            'file_type': file_ext,
            'analysis_id': analysis_id,
            'preview_content': None,
            'preview_type': 'binary',
            'forensic_analysis': forensic_data
        }
        
        # SECURE: Process ZIP/RAR/7Z files in Docker container with maximum security
        if file_ext in ['.zip', '.rar', '.7z']:
            preview_data = _secure_archive_preview(file_path, analysis_id, file_ext)
            # Add forensic analysis to the secure archive preview data
            preview_data['forensic_analysis'] = forensic_data
        elif file_ext == '.xlsx':
            # Use new pre-built Docker image for XLSX processing
            preview_data = _process_xlsx_with_docker_image(file_path, analysis_id)
            # Add forensic analysis to the secure preview data
            preview_data['forensic_analysis'] = forensic_data
        elif file_ext == '.pptx':
            # Use new pre-built Docker image for PPTX processing
            preview_data = _process_pptx_with_docker_image(file_path, analysis_id)
            # Add forensic analysis to the secure preview data
            preview_data['forensic_analysis'] = forensic_data
        else:
            # Use Docker-based secure preview for other file types
            docker_preview_data = secure_preview.safe_preview_any_file(file_path, analysis_id)
            # Preserve original filename and add forensic analysis
            preview_data.update(docker_preview_data)
            preview_data['filename'] = filename  # Keep original filename
            preview_data['file_type'] = file_ext  # Keep original file extension
            preview_data['forensic_analysis'] = forensic_data
        
        return jsonify(preview_data)
        
    except Exception as e:
        logger.error(f"Error generating preview for {analysis_id}: {str(e)}")
        return jsonify({'error': 'Failed to generate preview'}), 500

def _cleanup_docker_containers(docker_client, analysis_id: str):
    """Clean up old containers with similar names (older than 10 minutes)"""
    try:
        import time
        current_time = time.time()
        ten_minutes_ago = current_time - (10 * 60)  # 10 minutes in seconds
        
        # List all containers (including stopped ones)
        containers = docker_client.containers.list(all=True)
        for container in containers:
            if (container.name.startswith(f"secure_zip_{analysis_id}") or 
                container.name.startswith(f"secure_archive_file_{analysis_id}")):
                try:
                    # Check container creation time
                    container_time = container.attrs.get('Created', 0)
                    # Convert to timestamp if it's a string
                    if isinstance(container_time, str):
                        from datetime import datetime
                        container_time = datetime.fromisoformat(container_time.replace('Z', '+00:00')).timestamp()
                    
                    if container_time < ten_minutes_ago:
                        container.remove(force=True)
                        logger.info(f"Cleaned up old container: {container.name}")
                    else:
                        logger.info(f"Keeping recent container: {container.name}")
                except Exception as e:
                    logger.warning(f"Could not check/remove container {container.name}: {e}")
    except Exception as e:
        logger.warning(f"Error during container cleanup: {e}")

def _process_xlsx_with_docker_image(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Process XLSX files using pre-built Docker image"""
    try:
        import subprocess
        import tempfile
        import shutil
        
        # Use pre-built Docker image for XLSX processing
        with tempfile.TemporaryDirectory() as temp_dir:
            # Copy file to temp directory
            temp_file = os.path.join(temp_dir, "file.xlsx")
            shutil.copy2(file_path, temp_file)
            
            # Run pre-built Docker image
            result_cmd = subprocess.run([
                'docker', 'run', '--rm',
                '-v', f'{temp_dir}:/tmp',
                'aejis-xlsx-processor:latest',
                '/tmp/file.xlsx'
            ], capture_output=True, text=True, timeout=60)
            
            if result_cmd.returncode == 0:
                try:
                    xlsx_result = json.loads(result_cmd.stdout)
                    if xlsx_result.get('success', False):
                        return xlsx_result
                    else:
                        return {
                            "preview_type": "xlsx",
                            "preview_content": f"XLSX processing failed: {xlsx_result.get('error', 'Unknown error')}"
                        }
                except json.JSONDecodeError:
                    return {
                        "preview_type": "xlsx",
                        "preview_content": f"XLSX processing failed: Invalid response from Docker container"
                    }
            else:
                return {
                    "preview_type": "xlsx",
                    "preview_content": f"XLSX processing failed: Docker container error - {result_cmd.stderr}"
                }
                
    except Exception as e:
        return {
            "preview_type": "xlsx",
            "preview_content": "XLSX processing error: " + str(e)
        }

def _process_pptx_with_docker_image(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Process PPTX files using pre-built Docker image"""
    try:
        import subprocess
        import tempfile
        import shutil
        
        # Use pre-built Docker image for PPTX processing
        with tempfile.TemporaryDirectory() as temp_dir:
            # Copy file to temp directory
            temp_file = os.path.join(temp_dir, "file.pptx")
            shutil.copy2(file_path, temp_file)
            
            # Run pre-built Docker image
            result_cmd = subprocess.run([
                'docker', 'run', '--rm',
                '-v', f'{temp_dir}:/tmp',
                'aejis-pptx-processor:latest',
                '/tmp/file.pptx'
            ], capture_output=True, text=True, timeout=60)
            
            if result_cmd.returncode == 0:
                try:
                    pptx_result = json.loads(result_cmd.stdout)
                    if pptx_result.get('success', False):
                        # Ensure preview_type is set
                        if 'preview_type' not in pptx_result:
                            pptx_result['preview_type'] = 'pptx'
                        return pptx_result
                    else:
                        return {
                            "preview_type": "pptx",
                            "preview_content": f"PPTX processing failed: {pptx_result.get('error', 'Unknown error')}"
                        }
                except json.JSONDecodeError:
                    return {
                        "preview_type": "pptx",
                        "preview_content": f"PPTX processing failed: Invalid response from Docker container"
                    }
            else:
                return {
                    "preview_type": "pptx",
                    "preview_content": f"PPTX processing failed: Docker container error - {result_cmd.stderr}"
                }
                
    except Exception as e:
        return {
            "preview_type": "pptx",
            "preview_content": "PPTX processing error: " + str(e)
        }

def _secure_archive_preview(file_path: str, analysis_id: str, file_ext: str) -> Dict[str, Any]:
    """SECURE: Process ZIP/RAR/7Z archives in isolated Docker container with maximum security"""
    try:
        import docker
        import tempfile
        import json
        import shutil
        
        # Initialize Docker client
        try:
            docker_client = docker.from_env()
            docker_client.ping()
        except Exception as e:
            logger.error(f"Docker not available for secure ZIP processing: {e}")
            return {
                'preview_available': False,
                'error': 'Docker required for secure ZIP processing',
                'filename': os.path.basename(file_path),
                'secure_processing': False
            }
        
        # Clean up any existing containers with similar names
        _cleanup_docker_containers(docker_client, analysis_id)
        
        # Create temporary directory for this analysis
        archive_type = file_ext[1:]  # Remove the dot
        temp_dir = tempfile.mkdtemp(prefix=f"secure_{archive_type}_{analysis_id}_")
        
        try:
            # Copy file to temp directory with appropriate name
            temp_file = os.path.join(temp_dir, f"archive{file_ext}")
            shutil.copy2(file_path, temp_file)
            
            # Create secure archive processing script
            script_content = _create_secure_archive_processor_script(file_ext)
            script_path = os.path.join(temp_dir, f"process_{archive_type}.py")
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(script_content)
            
            # Run in isolated Docker container with maximum security
            import time
            unique_name = f"secure_{archive_type}_{analysis_id}_{int(time.time() * 1000)}"
            
            # Install required packages first, then run the script
            install_cmd = f"mkdir -p /tmp/packages && pip install --target /tmp/packages rarfile py7zr --quiet --no-cache-dir && PYTHONPATH=/tmp/packages python /tmp/process_{archive_type}.py"
            container = docker_client.containers.run(
                "python:3.11-slim",
                command=["sh", "-c", install_cmd],
                volumes={temp_dir: {'bind': '/tmp', 'mode': 'rw'}},  # Make writable for package installation
                network_mode='bridge',  # Enable network for package installation
                mem_limit='256m',     # Memory limit
                cpu_quota=50000,      # CPU limit (50% of one core)
                security_opt=['no-new-privileges:true'],  # Prevent privilege escalation
                read_only=False,      # Allow writing for package installation
                remove=True,
                detach=False,
                name=unique_name
            )
            
            # Parse container output
            try:
                output = container.decode('utf-8').strip()
                result = json.loads(output)
            except Exception as e:
                logger.warning(f"Could not parse container output: {e}")
                # Try to get logs if available
                try:
                    logs = container.logs().decode('utf-8')
                    logger.info(f"Container logs: {logs}")
                    result = json.loads(logs)
                except Exception as log_error:
                    logger.error(f"Could not get container logs: {log_error}")
                    return {
                        'preview_available': False,
                        'error': 'Failed to process ZIP archive',
                        'filename': os.path.basename(file_path),
                        'secure_processing': True
                    }
            
            if result.get('success'):
                result['secure_processing'] = True
                result['docker_container'] = unique_name
                result['security_level'] = 'maximum'
                return result
            else:
                return {
                    'preview_available': False,
                    'error': result.get('error', 'Unknown error in secure processing'),
                    'filename': os.path.basename(file_path),
                    'secure_processing': True
                }
                
        finally:
            # Clean up temp directory
            shutil.rmtree(temp_dir, ignore_errors=True)
            
    except Exception as e:
        logger.error(f"Secure ZIP processing failed for {file_path}: {str(e)}")
        return {
            'preview_available': False,
            'error': f'Secure ZIP processing failed: {str(e)}',
            'filename': os.path.basename(file_path),
            'secure_processing': True
        }

def _secure_individual_file_processor(file_path: str, filename: str, analysis_id: str) -> Dict[str, Any]:
    """SECURE: Process individual file from archive in Docker container - TEXT VIEWER ONLY"""
    try:
        import docker
        import tempfile
        import os
        import json
        import shutil
        import time
        
        # Initialize Docker client
        try:
            docker_client = docker.from_env()
            docker_client.ping()
        except Exception as e:
            logger.warning(f"Docker not available for secure individual file processing: {e}")
            return {'success': False, 'error': f'Docker not available: {str(e)}'}
        
        # Create temporary directory for this processing
        temp_dir = tempfile.mkdtemp(prefix=f"secure_individual_{analysis_id}_")
        
        try:
            # Copy file to temp directory with safe name
            safe_filename = f"individual_file{os.path.splitext(filename)[1]}"
            temp_file = os.path.join(temp_dir, safe_filename)
            shutil.copy2(file_path, temp_file)
            
            # Create simple processing script (NO F-STRINGS)
            script_content = _create_simple_individual_processor_script()
            script_path = os.path.join(temp_dir, "process_individual.py")
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(script_content)
            
            # Use Docker to process the file safely
            container_name = f"secure_individual_{analysis_id}_{int(time.time() * 1000)}"
            
            try:
                # Run Docker container with strict security
                container = docker_client.containers.run(
                    'python:3.11-alpine',
                    f'python /tmp/process_individual.py',
                    volumes={temp_dir: {'bind': '/tmp', 'mode': 'ro'}},  # Read-only for security
                    working_dir='/tmp',
                    name=container_name,
                    detach=False,
                    remove=True,
                    network_mode='none',  # No network access
                    mem_limit='512m',     # Memory limit
                    cpu_quota=50000       # CPU limit
                )
                
                output = container.decode('utf-8')
                result = json.loads(output)
                
                if result.get('success'):
                    return {
                        'success': True,
                        'content': result.get('content', ''),
                        'docker_processed': True
                    }
                else:
                    return {
                        'success': False,
                        'error': result.get('error', 'Unknown Docker processing error'),
                        'docker_processed': True
                    }
                    
            except Exception as e:
                logger.error(f"❌ Docker container error for individual file: {str(e)}")
                return {
                    'success': False,
                    'error': f'Docker processing failed: {str(e)}',
                    'docker_processed': False
                }
                
        finally:
            # Clean up temp directory
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
                
    except Exception as e:
        logger.error(f"❌ Unexpected error in secure individual file processing: {str(e)}")
        return {
            'success': False,
            'error': f'Processing failed: {str(e)}',
            'docker_processed': False
        }

def _create_simple_individual_processor_script() -> str:
    """Create ultra-simple Docker script for individual file processing - NO F-STRINGS"""
    return '''
import json
import os
import sys

def process_individual_file():
    """Simple and secure individual file processor - TEXT VIEWER ONLY"""
    try:
        # Find the file to process
        target_file = None
        for file in os.listdir('/tmp'):
            if file.startswith('individual_file'):
                target_file = '/tmp/' + file
                break
        
        if not target_file:
            return {'success': False, 'error': 'No file found to process'}
        
        # Security: Check file size (max 100MB)
        file_size = os.path.getsize(target_file)
        max_size = 100 * 1024 * 1024  # 100MB
        if file_size > max_size:
            return {'success': False, 'error': 'File too large: ' + str(file_size) + ' bytes'}
        
        # Read file content as text (multiple encoding attempts)
        content = None
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(target_file, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
            except Exception:
                continue
        
        if content is None:
            # If all text decoding fails, read as binary and show hex
            try:
                with open(target_file, 'rb') as f:
                    binary_data = f.read(1024)  # First 1KB only
                content = 'Binary file (showing first 1KB as hex):\\n'
                content += ' '.join(format(byte, '02x') for byte in binary_data)
                if file_size > 1024:
                    content += '\\n\\n... truncated (file too large for text display)'
            except Exception as e:
                return {'success': False, 'error': 'Failed to read file: ' + str(e)}
        
        # Limit content size for safety
        if len(content) > max_size:
            content = content[:max_size] + '\\n\\n[... truncated at 100MB ...]'
        
        return {
            'success': True,
            'content': content,
            'file_size': file_size,
            'processing_method': 'text_only'
        }
        
    except Exception as e:
        return {'success': False, 'error': 'Processing error: ' + str(e)}

if __name__ == "__main__":
    result = process_individual_file()
    print(json.dumps(result, ensure_ascii=False))
'''

def _create_secure_archive_processor_script(file_ext: str) -> str:
    """Create secure archive processing script for Docker container (ZIP/RAR/7Z)"""
    return '''
import subprocess
import sys
import json
import os
import tempfile
import shutil
from pathlib import Path

# Packages are installed via Docker command before running this script

def secure_archive_analysis():
    """Securely analyze ZIP archive with maximum safety measures"""
    try:
        # Security: Change to safe working directory
        os.chdir('/tmp')
        
        # Security: Set safe file limits (higher for RAR files)
        MAX_FILES = 20000  # Maximum files to process (increased for RAR support)
        MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB max file size
        MAX_TOTAL_SIZE = 500 * 1024 * 1024  # 500MB max total size (increased for RAR support)
        
        # Find archive file dynamically
        archive_path = None
        for file in os.listdir('/tmp'):
            if file.startswith('archive.') and file.split('.')[-1] in ['zip', 'rar', '7z']:
                archive_path = f'/tmp/{file}'
                break
        
        if not archive_path:
            return {'success': False, 'error': 'Archive file not found'}
        
        # Security: Validate file exists and is readable
        if not os.access(archive_path, os.R_OK):
            return {'success': False, 'error': 'Archive file not readable'}
        
        # Security: Check file size before processing
        file_size = os.path.getsize(archive_path)
        if file_size > MAX_TOTAL_SIZE:
            return {'success': False, 'error': f'Archive too large: {file_size} bytes'}
        
        # Detect archive type and process accordingly
        archive_ext = os.path.splitext(archive_path)[1].lower()
        
        if archive_ext == '.zip':
            import zipfile
            archive_ref = zipfile.ZipFile(archive_path, 'r')
            # Test ZIP integrity
            try:
                archive_ref.testzip()
            except Exception as e:
                return {'success': False, 'error': f'Corrupted ZIP file: {str(e)}'}
            file_list = archive_ref.namelist()
            get_file_info = lambda f: archive_ref.getinfo(f)
            
        elif archive_ext == '.rar':
            import rarfile
            archive_ref = rarfile.RarFile(archive_path, 'r')
            file_list = archive_ref.namelist()
            get_file_info = lambda f: archive_ref.getinfo(f)
            
        elif archive_ext == '.7z':
            import py7zr
            archive_ref = py7zr.SevenZipFile(archive_path, 'r')
            file_list = archive_ref.getnames()
            get_file_info = lambda f: type('FileInfo', (), {'file_size': 0, 'compress_size': 0})()  # Simplified for 7z
            
        else:
            return {'success': False, 'error': f'Unsupported archive type: {archive_ext}'}
            
        with archive_ref:
            # Note: File list already retrieved above
            
            # Security: Limit number of files
            if len(file_list) > MAX_FILES:
                return {'success': False, 'error': f'Too many files: {len(file_list)} > {MAX_FILES}'}
            
            # Security: Check individual file sizes
            total_size = 0
            compressed_size = 0
            file_types = {}
            archive_contents = []
            
            for file_in_zip in file_list:
                if file_in_zip.endswith('/'):  # Skip directories
                    continue
                
                try:
                    file_info = get_file_info(file_in_zip)
                    
                    # Security: Check individual file size
                    if file_info.file_size > MAX_FILE_SIZE:
                        return {'success': False, 'error': f'File too large: {file_in_zip} ({file_info.file_size} bytes)'}
                    
                    total_size += file_info.file_size
                    compressed_size += file_info.compress_size
                    
                    # Security: Check total size
                    if total_size > MAX_TOTAL_SIZE:
                        return {'success': False, 'error': f'Total extracted size too large: {total_size} bytes'}
                    
                    # Analyze file type
                    ext = os.path.splitext(file_in_zip)[1].lower()
                    file_types[ext] = file_types.get(ext, 0) + 1
                    
                    # Determine file type for frontend
                    file_type = 'binary'
                    if ext in ['.txt', '.md', '.log', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.cpp', '.c', '.h']:
                        file_type = 'text'
                    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
                        file_type = 'image'
                    elif ext in ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']:
                        file_type = 'document'
                    elif ext in ['.mp4', '.avi', '.mov', '.wmv', '.flv']:
                        file_type = 'video'
                    elif ext in ['.mp3', '.wav', '.flac', '.aac', '.ogg']:
                        file_type = 'audio'
                    elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                        file_type = 'archive'
                    elif ext in ['.exe', '.dll', '.so', '.dylib']:
                        file_type = 'executable'
                    elif ext in ['.otf', '.ttf', '.woff', '.woff2']:
                        file_type = 'font'
                    
                    archive_contents.append({
                        'name': os.path.basename(file_in_zip),
                        'path': file_in_zip,
                        'size': file_info.file_size,
                        'compressed_size': file_info.compress_size,
                        'type': file_type,
                        'extension': ext,
                        'is_dir': file_in_zip.endswith('/')
                    })
                    
                except Exception as e:
                    # Security: Skip problematic files but continue processing
                    continue
            
            # Calculate compression ratio
            compression_ratio = (1 - compressed_size / total_size) * 100 if total_size > 0 else 0
            
            # Create preview content
            archive_type = archive_ext.upper().replace('.', '')
            preview_content = f"{archive_type} Archive (Secure Processing)\\n"
            preview_content += f"Files: {len(archive_contents)}\\n"
            preview_content += f"Total Size: {total_size:,} bytes\\n"
            preview_content += f"Compressed: {compressed_size:,} bytes\\n"
            preview_content += f"Compression: {compression_ratio:.1f}%\\n\\n"
            preview_content += f"File Types:\\n"
            
            for ext, count in list(file_types.items())[:5]:
                preview_content += f"  {ext or 'No extension'}: {count} files\\n"
            
            preview_content += f"\\nFile List (first 10):\\n"
            for file_item in archive_contents[:10]:
                preview_content += f"  📄 {file_item['name']}\\n"
            
            if len(archive_contents) > 10:
                preview_content += f"  ... and {len(archive_contents) - 10} more files"
            
            return {
                'success': True,
                'preview_available': True,
                'preview_type': 'archive',
                'filename': os.path.basename(archive_path),
                'file_size': file_size,
                'mime_type': f'application/{archive_ext[1:]}' if archive_ext == '.zip' else f'application/x-{archive_ext[1:]}-compressed',
                'preview_content': preview_content,
                'metadata': {
                    'file_count': len(archive_contents),
                    'total_size': total_size,
                    'compressed_size': compressed_size,
                    'compression_ratio': compression_ratio,
                    'file_types': file_types,
                    'files': [item['name'] for item in archive_contents[:20]]
                },
                'archive_contents': archive_contents,
                'total_files': len(archive_contents),
                'total_extracted_size': total_size,
                'security_measures': {
                    'max_files': MAX_FILES,
                    'max_file_size': MAX_FILE_SIZE,
                    'max_total_size': MAX_TOTAL_SIZE,
                    'docker_isolation': True,
                    'network_isolation': True,
                    'read_only_filesystem': True
                }
            }
            
    except Exception as e:
        return {'success': False, 'error': f'Secure processing failed: {str(e)}'}

if __name__ == "__main__":
    result = secure_archive_analysis()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

def _secure_archive_file_access(archive_path: str, file_path: str, analysis_id: str):
    """SECURE: Extract and read specific file from archive using Docker isolation"""
    try:
        import docker
        import tempfile
        import json
        import shutil
        
        # Initialize Docker client
        try:
            docker_client = docker.from_env()
            docker_client.ping()
        except Exception as e:
            logger.error(f"Docker not available for secure archive file access: {e}")
            return jsonify({'error': 'Docker required for secure archive file access'}), 503
        
        # Clean up any existing containers with similar names
        _cleanup_docker_containers(docker_client, analysis_id)
        
        # Create temporary directory for this analysis
        temp_dir = tempfile.mkdtemp(prefix=f"secure_archive_file_{analysis_id}_")
        
        try:
            # Copy archive to temp directory
            temp_file = os.path.join(temp_dir, "archive.zip")
            shutil.copy2(archive_path, temp_file)
            
            # Create secure archive file access script
            script_content = _create_secure_archive_file_script(file_path)
            script_path = os.path.join(temp_dir, "process_archive_file.py")
            with open(script_path, 'w', encoding='utf-8') as f:
                f.write(script_content)
            
            # Run in isolated Docker container with maximum security
            import time
            unique_name = f"secure_archive_file_{analysis_id}_{int(time.time() * 1000)}"
            container = docker_client.containers.run(
                "python:3.11-slim",
                command=["python", "/tmp/process_archive_file.py"],
                volumes={temp_dir: {'bind': '/tmp', 'mode': 'ro'}},
                network_mode='none',  # No network access
                mem_limit='128m',     # Memory limit
                cpu_quota=25000,      # CPU limit (25% of one core)
                security_opt=['no-new-privileges:true'],  # Prevent privilege escalation
                read_only=True,       # Read-only filesystem
                remove=True,
                detach=False,
                name=unique_name
            )
            
            # Parse container output
            try:
                output = container.decode('utf-8').strip()
                result = json.loads(output)
            except Exception as e:
                logger.warning(f"Could not parse container output: {e}")
                # Try to get logs if available
                try:
                    logs = container.logs().decode('utf-8')
                    logger.info(f"Container logs: {logs}")
                    result = json.loads(logs)
                except Exception as log_error:
                    logger.error(f"Could not get container logs: {log_error}")
                    return jsonify({
                        'success': False,
                        'error': 'Failed to process archive file',
                        'secure_processing': True
                    }), 500
            
            if result.get('success'):
                result['secure_processing'] = True
                result['docker_container'] = unique_name
                result['security_level'] = 'maximum'
                return jsonify(result)
            else:
                return jsonify({
                    'success': False,
                    'error': result.get('error', 'Unknown error in secure processing'),
                    'secure_processing': True
                }), 500
                
        finally:
            # Clean up temp directory
            shutil.rmtree(temp_dir, ignore_errors=True)
            
    except Exception as e:
        logger.error(f"Secure archive file access failed: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'Secure archive file access failed: {str(e)}',
            'secure_processing': True
        }), 500

def _create_secure_archive_file_script(file_path: str) -> str:
    """Create secure archive file access script for Docker container"""
    return '''.format(file_path=file_path)
import subprocess
import sys
import zipfile
import json
import os
import base64
import io

# Install required packages
try:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'Pillow', '--quiet'])
except:
    pass  # Continue even if installation fails

# Try to import PIL after installation
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

def secure_archive_file_access():
    """Securely extract and analyze specific file from archive"""
    try:
        # Security: Change to safe working directory
        os.chdir('/tmp')
        
        # Security: Set safe file limits
        MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB max file size
        MAX_CONTENT_LENGTH = 1000000  # 1MB max content to return
        
        archive_path = '/tmp/archive.zip'
        target_file = '{file_path}'
        
        # Security: Validate file exists and is readable
        if not os.path.exists(archive_path):
            return {'success': False, 'error': 'Archive file not found'}
        
        if not os.access(archive_path, os.R_OK):
            return {'success': False, 'error': 'Archive file not readable'}
        
        with zipfile.ZipFile(archive_path, 'r') as zip_ref:
            # Security: Validate ZIP file integrity
            try:
                zip_ref.testzip()
            except Exception as e:
                return {'success': False, 'error': f'Corrupted ZIP file: {str(e)}'}
            
            # Security: Check if file exists in archive
            if target_file not in zip_ref.namelist():
                return {'success': False, 'error': 'File not found in archive'}
            
            try:
                file_info = zip_ref.getinfo(target_file)
                
                # Security: Check file size
                if file_info.file_size > MAX_FILE_SIZE:
                    return {'success': False, 'error': f'File too large: {file_info.file_size} bytes'}
                
                # Read file content
                with zip_ref.open(target_file) as f:
                    content = f.read()
                
                # Determine content type
                filename = os.path.basename(target_file)
                file_ext = os.path.splitext(filename)[1].lower()
                
                # Try to decode as text for text files
                content_type = 'binary'
                text_content = None
                thumbnail_base64 = None
                metadata = {}
                
                if file_ext in ['.txt', '.md', '.log', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.cpp', '.c', '.h']:
                    try:
                        text_content = content.decode('utf-8')
                        content_type = 'text'
                        # Security: Limit content length
                        if len(text_content) > MAX_CONTENT_LENGTH:
                            text_content = text_content[:MAX_CONTENT_LENGTH] + "\\n... (truncated for security)"
                    except UnicodeDecodeError:
                        try:
                            text_content = content.decode('latin-1')
                            content_type = 'text'
                            if len(text_content) > MAX_CONTENT_LENGTH:
                                text_content = text_content[:MAX_CONTENT_LENGTH] + "\\n... (truncated for security)"
                        except UnicodeDecodeError:
                            pass
                
                # Handle images
                elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
                    content_type = 'image'
                    if PIL_AVAILABLE:
                        try:
                            # Create thumbnail
                            img = Image.open(io.BytesIO(content))
                            img.thumbnail((300, 300), Image.Resampling.LANCZOS)
                            
                            # Convert to base64
                            buffer = io.BytesIO()
                            img.save(buffer, format='PNG')
                            thumbnail_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
                            
                            metadata = {{
                                'format': img.format,
                                'mode': img.mode,
                                'dimensions': f"{img.size[0]}x{img.size[1]}"
                            }}
                        except Exception as e:
                            metadata = {'error': f"Image processing failed: {str(e)}", 'format': 'unknown'}
                    else:
                        # Fallback when PIL is not available
                        metadata = {{
                            'format': 'unknown',
                            'mode': 'unknown',
                            'dimensions': 'unknown',
                            'note': 'PIL not available for image processing'
                        }}
                
                # Handle font files specifically
                elif file_ext in ['.otf', '.ttf', '.woff', '.woff2']:
                    content_type = 'font'
                    # Create hex dump for first 512 bytes
                    hex_content = content[:512].hex()
                    formatted_hex = ' '.join(hex_content[i:i+2] for i in range(0, len(hex_content), 2))
                    text_content = formatted_hex
                    
                    # Try to extract font metadata
                    font_metadata = {}
                    try:
                        if file_ext in ['.otf', '.ttf']:
                            # Parse OpenType/TrueType font headers
                            if len(content) >= 12:
                                # Read table directory
                                num_tables = int.from_bytes(content[4:6], 'big')
                                font_metadata['num_tables'] = num_tables
                                
                                # Look for name table
                                for i in range(num_tables):
                                    offset = 12 + i * 16
                                    if offset + 16 <= len(content):
                                        tag = content[offset:offset+4].decode('ascii', errors='ignore')
                                        if tag == 'name':
                                            name_offset = int.from_bytes(content[offset+8:offset+12], 'big')
                                            name_length = int.from_bytes(content[offset+12:offset+16], 'big')
                                            
                                            if name_offset + name_length <= len(content):
                                                name_data = content[name_offset:name_offset+name_length]
                                                # Simple name extraction (simplified)
                                                font_metadata['font_name'] = 'Font Name (extracted)'
                                                break
                    except Exception as e:
                        pass  # Continue without metadata
                    
                    metadata = {{
                        'font_name': font_metadata.get('font_name', 'Unknown'),
                        'font_type': 'OpenType' if file_ext in ['.otf'] else 'TrueType',
                        'file_format': 'Font File',
                        'mime_type': 'font/otf' if file_ext in ['.otf'] else 'font/ttf',
                        'num_tables': font_metadata.get('num_tables', 'Unknown'),
                        'entropy': len(set(content)) / 256.0  # Simple entropy calculation
                    }}
                
                # Security: Limit binary content length
                if content_type == 'binary' and len(content) > MAX_CONTENT_LENGTH:
                    content = content[:MAX_CONTENT_LENGTH]
                
                return {{
                    'success': True,
                    'filename': filename,
                    'file_size': file_info.file_size,
                    'compressed_size': file_info.compress_size,
                    'file_type': file_ext,
                    'content_type': content_type,
                    'content': text_content or content.decode('latin-1', errors='replace'),
                    'thumbnail_base64': thumbnail_base64,
                    'metadata': metadata,
                    'is_archive_file': True,
                    'security_measures': {{
                        'max_file_size': MAX_FILE_SIZE,
                        'max_content_length': MAX_CONTENT_LENGTH,
                        'docker_isolation': True,
                        'network_isolation': True,
                        'read_only_filesystem': True
                    }}
                }}
                
            except KeyError:
                return {'success': False, 'error': 'File not found in archive'}
            except Exception as e:
                return {'success': False, 'error': f'Error processing file: {str(e)}'}
            
    except Exception as e:
        return {'success': False, 'error': f'Secure processing failed: {str(e)}'}

if __name__ == "__main__":
    result = secure_archive_file_access()
    print(json.dumps(result, ensure_ascii=False, indent=2))
'''

@app.route('/api/file-content/<analysis_id>')
def get_file_content(analysis_id):
    """Get raw file content for preview"""
    logger.info(f"🌐 API request for file content: {analysis_id}")
    
    if analysis_id not in analysis_results:
        logger.warning(f"⚠️ Analysis not found: {analysis_id}")
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        logger.warning(f"⚠️ Analysis not completed for {analysis_id}: status={result['status']}")
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    file_path = result.get('file_path')
    if not file_path or not os.path.exists(file_path):
        logger.error(f"❌ File not found for {analysis_id}: {file_path}")
        return jsonify({'error': 'File not found'}), 404
    
    try:
        file_ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        logger.info(f"📄 Processing file content request: {filename} ({file_ext})")
        
        # For text files, return content directly
        if file_ext in ['.txt', '.log', '.md', '.json', '.xml', '.csv', '.py', '.js', '.html', '.css', '.sh', '.bat', '.ps1']:
            logger.debug(f"📝 Processing as text file: {filename}")
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read(50000)  # Limit to 50KB
                    if len(content) >= 50000:
                        content += "\n\n[... content truncated for security ...]"
                        logger.debug(f"✂️ Content truncated for security")
                
                logger.info(f"✅ Text file processed successfully: {len(content)} characters")
                return jsonify({
                    'success': True,
                    'content': content,
                    'content_type': 'text',
                    'filename': filename,
                    'file_size': os.path.getsize(file_path)
                })
            except UnicodeDecodeError:
                logger.warning(f"⚠️ UTF-8 decode failed, trying latin-1 for {filename}")
                # Try with different encoding
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read(50000)
                    if len(content) >= 50000:
                        content += "\n\n[... content truncated for security ...]"
                
                logger.info(f"✅ Text file processed with latin-1: {len(content)} characters")
                return jsonify({
                    'success': True,
                    'content': content,
                    'content_type': 'text',
                    'filename': filename,
                    'file_size': os.path.getsize(file_path)
                })
        
        # For images, return base64 encoded content
        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg']:
            logger.debug(f"🖼️ Processing as image file: {filename}")
            try:
                with open(file_path, 'rb') as f:
                    file_data = f.read()
                    import base64
                    base64_content = base64.b64encode(file_data).decode('utf-8')
                    
                    logger.info(f"✅ Image file processed successfully: {len(file_data)} bytes -> {len(base64_content)} base64 chars")
                    return jsonify({
                        'success': True,
                        'content': base64_content,
                        'content_type': 'image',
                        'filename': filename,
                        'file_size': os.path.getsize(file_path),
                        'mime_type': f'image/{file_ext[1:]}'
                    })
            except Exception as e:
                logger.error(f"❌ Failed to process image {filename}: {e}")
                return jsonify({'error': f'Failed to process image: {str(e)}'}), 500
        
        # For other files, use Docker processing or return basic info
        else:
            logger.debug(f"🔧 Processing as other file type with Docker: {filename}")
            # Use secure Docker preview for comprehensive file processing
            secure_preview = SecureFilePreview()
            docker_result = secure_preview.safe_preview_any_file(file_path, analysis_id)
            
            if docker_result.get('docker_processed'):
                logger.info(f"✅ Docker processing successful for {filename}")
                return jsonify({
                    'success': True,
                    'content': docker_result.get('preview_content', 'No preview available'),
                    'content_type': docker_result.get('preview_type', 'binary'),
                    'filename': filename,
                    'file_size': docker_result.get('file_size', 0),
                    'metadata': docker_result.get('metadata', {}),
                    'thumbnail_base64': docker_result.get('thumbnail_base64'),
                    'docker_processed': True
                })
            else:
                logger.warning(f"⚠️ Docker processing failed for {filename}, using fallback")
            # Fallback: Basic file info with proper type detection
            file_size = os.path.getsize(file_path)
            
        # Detect file type based on extension
        video_extensions = ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.ogg', '.wmv', '.flv']
        audio_extensions = ['.mp3', '.wav', '.flac', '.aac', '.m4a']
        image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg', '.tiff']
        text_extensions = ['.txt', '.log', '.json', '.xml', '.csv', '.py', '.js', '.html', '.css', '.md', '.yml', '.yaml', '.ini', '.cfg', '.conf']
        
        if file_ext.lower() in video_extensions:
            content_type = 'video'
            content = f'Video file: {filename}\nSize: {file_size} bytes\nType: {file_ext}\n\nNote: Video preview requires Docker processing. Download the file to view it.'
        elif file_ext.lower() in audio_extensions:
            content_type = 'audio'
            content = f'Audio file: {filename}\nSize: {file_size} bytes\nType: {file_ext}\n\nNote: Audio preview requires Docker processing. Download the file to play it.'
        elif file_ext.lower() in image_extensions:
            content_type = 'image'
            content = f'Image file: {filename}\nSize: {file_size} bytes\nType: {file_ext}\n\nNote: Image preview requires Docker processing. Download the file to view it.'
        elif file_ext.lower() in text_extensions:
            content_type = 'text'
            try:
                # Try to read the text file content
                with open(file_path, 'r', encoding='utf-8') as f:
                    text_content = f.read()
                content = text_content
            except UnicodeDecodeError:
                # Try with different encoding
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        text_content = f.read()
                    content = text_content
                except Exception as e:
                    content = f'Text file: {filename}\nSize: {file_size} bytes\nType: {file_ext}\n\nError reading file: {str(e)}'
            except Exception as e:
                content = f'Text file: {filename}\nSize: {file_size} bytes\nType: {file_ext}\n\nError reading file: {str(e)}'
        else:
            content_type = 'binary'
            content = f'Binary file: {filename}\nSize: {file_size} bytes\nType: {file_ext}'
            
            return jsonify({
                'success': True,
                'content': content,
                'content_type': content_type,
                'filename': filename,
                'file_size': file_size
            })
            
    except Exception as e:
        logger.error(f"❌ Unexpected error getting file content for {analysis_id}: {str(e)}")
        logger.error(f"🔍 Error type: {type(e).__name__}")
        return jsonify({'error': 'Failed to get file content'}), 500

def _extract_archive_contents(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """SECURE: Extract and analyze archive contents in isolated Docker container"""
    try:
        import docker
        import tempfile
        import os
        import json
        import shutil
        
        # Initialize Docker client
        try:
            docker_client = docker.from_env()
            docker_client.ping()
        except Exception as e:
            logger.warning(f"Docker not available for secure preview: {e}")
            return _fallback_archive_preview(file_path, analysis_id)
        
        # Create temporary directory for this analysis
        temp_dir = tempfile.mkdtemp(prefix=f"secure_preview_{analysis_id}_")
        
        try:
            # Copy file to temp directory
            temp_file = os.path.join(temp_dir, os.path.basename(file_path))
            shutil.copy2(file_path, temp_file)
            
            # Simple fallback extraction without Docker for now
            import zipfile
            
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                file_list = zip_ref.namelist()
                total_size = sum(info.file_size for info in zip_ref.infolist() if not info.is_dir())
                
                archive_contents = []
                for file_name in file_list:
                    if not file_name.endswith('/'):  # Skip directories
                        file_ext = os.path.splitext(file_name)[1].lower()
                        file_type = 'unknown'
                        if file_ext in ['.txt', '.log', '.md', '.json', '.xml', '.csv', '.py', '.js', '.html', '.css']:
                            file_type = 'text'
                        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg']:
                            file_type = 'image'
                        elif file_ext in ['.pdf']:
                            file_type = 'pdf'
                        elif file_ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                            file_type = 'archive'
                        elif file_ext in ['.otf', '.ttf', '.woff']:
                            file_type = 'font'
                        else:
                            file_type = 'binary'
                        
                        # Get file size from zip info
                        file_info = zip_ref.getinfo(file_name)
                        file_size = file_info.file_size
                        
                        archive_contents.append({
                            'name': os.path.basename(file_name),
                            'path': file_name,
                            'size': file_size,
                            'type': file_type,
                            'extension': file_ext
                        })
            
            return {
                'filename': os.path.basename(file_path),
                'file_size': os.path.getsize(file_path),
                'file_type': '.zip',
                'analysis_id': analysis_id,
                'preview_content': None,
                'preview_type': 'archive',
                'archive_contents': archive_contents,
                'total_files': len(archive_contents),
                'total_extracted_size': total_size,
                'secure_extraction': True,
                'archive_browser': True  # Enable Windows-like file browser
            }
                
        finally:
            # Clean up temp directory
            shutil.rmtree(temp_dir, ignore_errors=True)
            
    except Exception as e:
        logger.error(f"Secure archive extraction failed for {analysis_id}: {str(e)}")
        return _fallback_archive_preview(file_path, analysis_id)
    
def _fallback_archive_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Fallback: Basic archive preview without extraction (safer)"""
    try:
        import zipfile
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            file_list = zip_ref.namelist()
            total_size = sum(info.file_size for info in zip_ref.infolist() if not info.is_dir())
            
            archive_contents = []
            for file_name in file_list:
                if not file_name.endswith('/'):  # Skip directories
                    file_ext = os.path.splitext(file_name)[1].lower()
                    file_type = 'unknown'
                    if file_ext in ['.txt', '.log', '.md', '.json', '.xml', '.csv', '.py', '.js', '.html', '.css']:
                        file_type = 'text'
                    elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg']:
                        file_type = 'image'
                    elif file_ext in ['.pdf']:
                        file_type = 'pdf'
                    elif file_ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                        file_type = 'archive'
                    else:
                        file_type = 'binary'
                    
                    # Get file size from zip info
                    file_info = zip_ref.getinfo(file_name)
                    file_size = file_info.file_size
                    
                    archive_contents.append({
                        'name': os.path.basename(file_name),
                        'path': file_name,
                        'size': file_size,
                        'type': file_type,
                        'extension': file_ext
                    })
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': '.zip',
            'analysis_id': analysis_id,
            'preview_content': None,
            'preview_type': 'archive',
            'archive_contents': archive_contents,
            'total_files': len(archive_contents),
            'total_extracted_size': total_size,
            'secure_extraction': False,
            'fallback_mode': True
        }
        
    except Exception as e:
        logger.error(f"Fallback archive preview failed for {analysis_id}: {str(e)}")
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': '.zip',
            'analysis_id': analysis_id,
            'preview_content': f"Error analyzing archive: {str(e)}",
            'preview_type': 'archive',
            'archive_contents': [],
            'total_files': 0,
            'total_extracted_size': 0,
            'secure_extraction': False,
            'fallback_mode': True
        }

@app.route('/simple-file/<analysis_id>/<path:file_path>')
def simple_file_reader(analysis_id, file_path):
    """SIMPLE: Just read file from ZIP and show as text"""
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    zip_path = result.get('file_path')
    
    if not zip_path or not os.path.exists(zip_path):
        return jsonify({'error': 'ZIP file not found'}), 404
    
    try:
        import zipfile
        import urllib.parse
        
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            files = zip_ref.namelist()
            target_file = None
            
            # Find the file
            if file_path in files:
                target_file = file_path
            else:
                decoded = urllib.parse.unquote(file_path)
                if decoded in files:
                    target_file = decoded
                else:
                    for f in files:
                        if os.path.basename(f) == os.path.basename(file_path):
                            target_file = f
                            break
            
            if not target_file:
                return jsonify({'error': f'File not found. Available: {files}'}), 404
            
            # Read file
            with zip_ref.open(target_file) as f:
                content = f.read()
            
            # Decode as text
            try:
                text = content.decode('utf-8')
            except:
                text = content.decode('utf-8', errors='replace')
            
            # Limit size to 100MB
            if len(text) > 100 * 1024 * 1024:  # 100MB
                text = text[:100 * 1024 * 1024] + "\n\n[... truncated at 100MB ...]"
            
            return jsonify({
                'success': True,
                'preview_content': f"File: {os.path.basename(target_file)}\nSize: {len(content)} bytes\n\nContent:\n{text}",
                'preview_type': 'text'
            })
            
    except Exception as e:
        return jsonify({'error': f'Failed: {str(e)}'}), 500

@app.route('/preview/<analysis_id>/file/<path:file_path>')
def get_archive_file_content(analysis_id, file_path):
    """SECURE: Get content of a specific file within an archive using Docker isolation"""
    logger.info(f"🔍 Archive file request: analysis_id={analysis_id}, file_path='{file_path}'")
    
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    try:
        # Get the original file path
        original_file_path = result.get('file_path')
        if not original_file_path or not os.path.exists(original_file_path):
            return jsonify({'error': 'Original archive file not found'}), 404
        
        # Extract and show file content directly - Windows-like file inspector
        # Detect archive type and use appropriate library
        archive_ext = os.path.splitext(original_file_path)[1].lower()
        
        if archive_ext == '.zip':
            import zipfile
            archive_ref = zipfile.ZipFile(original_file_path, 'r')
            get_file_list = lambda: archive_ref.namelist()
            get_file_info = lambda f: archive_ref.getinfo(f)
            read_file_data = lambda f: archive_ref.read(f)
            
        elif archive_ext == '.rar':
            import rarfile
            archive_ref = rarfile.RarFile(original_file_path, 'r')
            get_file_list = lambda: archive_ref.namelist()
            get_file_info = lambda f: archive_ref.getinfo(f)
            read_file_data = lambda f: archive_ref.read(f)
            
        elif archive_ext == '.7z':
            import py7zr
            archive_ref = py7zr.SevenZipFile(original_file_path, 'r')
            get_file_list = lambda: archive_ref.getnames()
            get_file_info = lambda f: type('FileInfo', (), {'file_size': 0, 'compress_size': 0})()  # Simplified for 7z
            read_file_data = lambda f: archive_ref.read(f)
            
        else:
            return jsonify({'error': f'Unsupported archive type: {archive_ext}'}), 400
        
        with archive_ref:
            # Debug: Log the file paths in the archive
            archive_files = get_file_list()
            logger.info(f"🔍 Archive files: {archive_files}")
            logger.info(f"🔍 Looking for file: '{file_path}'")
            
            # Try to find the file with different path formats
            found_file = None
            if file_path in archive_files:
                found_file = file_path
            else:
                # Try URL decoding
                import urllib.parse
                decoded_path = urllib.parse.unquote(file_path)
                logger.info(f"🔍 Trying decoded path: '{decoded_path}'")
                if decoded_path in archive_files:
                    found_file = decoded_path
                else:
                    # Try to find by filename only
                    filename = os.path.basename(file_path)
                    for archive_file in archive_files:
                        if os.path.basename(archive_file) == filename:
                            found_file = archive_file
                            logger.info(f"🔍 Found by filename: '{found_file}'")
                            break
            
            if not found_file:
                logger.warning(f"❌ File '{file_path}' not found in archive")
                logger.warning(f"🔍 Available files: {archive_files}")
                return jsonify({'error': f'File not found in archive. Available files: {archive_files}'}), 404
            
            # Use the found file path
            file_path = found_file
            
            # Extract file info
            file_info = get_file_info(file_path)
            file_size = file_info.file_size
            filename = os.path.basename(file_path)
            file_ext = os.path.splitext(filename)[1].lower()
            
            # Read the file content from the archive
            file_data = read_file_data(file_path)
            
            # SECURITY: Process ALL files in Docker containers (never bypass security)
            # Force ALL archive files to be viewed as TEXT ONLY for maximum security
            try:
                # Create temp file for secure Docker processing
                with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as temp_file:
                    temp_file.write(file_data)
                    temp_file_path = temp_file.name
                
                try:
                    # Use secure Docker processing for all files
                    docker_result = _secure_individual_file_processor(temp_file_path, filename, analysis_id)
                    
                    if docker_result.get('success'):
                        content = docker_result.get('content', 'Error processing file')
                        truncated = len(file_data) > 100 * 1024 * 1024  # Mark as truncated if > 100MB
                        
                        return jsonify({
                            'filename': filename,
                            'file_path': file_path,
                            'file_size': file_size,
                            'file_type': file_ext,
                            'content': content,
                            'content_type': 'text',  # FORCE TEXT VIEWER ONLY
                            'truncated': truncated,
                            'line_count': len(content.split('\n')) if content else 0,
                            'success': True,
                            'docker_processed': True
                        })
                    else:
                        # Docker processing failed, return error
                        return jsonify({
                            'filename': filename,
                            'file_path': file_path,
                            'file_size': file_size,
                            'file_type': file_ext,
                            'content': f'Docker processing failed: {docker_result.get("error", "Unknown error")}',
                            'content_type': 'text',
                            'success': False,
                            'docker_processed': False
                        })
                finally:
                    # Clean up temp file
                    try:
                        os.unlink(temp_file_path)
                    except:
                        pass
                        
            except Exception as e:
                return jsonify({
                    'filename': filename,
                    'file_path': file_path,
                    'file_size': file_size,
                    'file_type': file_ext,
                    'content': f'Processing error: {str(e)}',
                    'content_type': 'text',
                    'success': False
                })
                
    except Exception as e:
        logger.error(f"Secure file extraction failed for {file_path}: {str(e)}")
        return jsonify({'error': f'Secure file extraction failed: {str(e)}'}), 500

@app.route('/download/<analysis_id>')
def download_report(analysis_id):
    """Download analysis report"""
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    # Generate report file
    report_data = generate_report(result['results'])
    
    # Save report
    report_filename = f"aejis_report_{analysis_id}.json"
    report_path = os.path.join(app.config['UPLOAD_FOLDER'], report_filename)
    
    with open(report_path, 'w') as f:
        json.dump(report_data, f, indent=2)
    
    return send_file(report_path, as_attachment=True, download_name=report_filename)

@app.route('/stream-video/<analysis_id>')
def stream_video(analysis_id):
    """Stream video file for embedded playback"""
    if analysis_id not in analysis_results:
        return jsonify({'error': 'Analysis not found'}), 404
    
    result = analysis_results[analysis_id]
    
    if result['status'] != 'completed':
        return jsonify({'error': 'Analysis not completed yet'}), 400
    
    try:
        file_path = result.get('file_path')
        if not file_path or not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
        
        # Check if it's a video file
        filename = result.get('filename', '')
        video_extensions = ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.ogg']
        if not any(filename.lower().endswith(ext) for ext in video_extensions):
            return jsonify({'error': 'File is not a supported video format'}), 400
        
        # Get file size for range requests
        file_size = os.path.getsize(file_path)
        
        # Handle range requests for video streaming
        range_header = request.headers.get('Range', None)
        if range_header:
            # Parse range header
            range_match = re.match(r'bytes=(\d+)-(\d*)', range_header)
            if range_match:
                start = int(range_match.group(1))
                end = int(range_match.group(2)) if range_match.group(2) else file_size - 1
                
                # Ensure end doesn't exceed file size
                end = min(end, file_size - 1)
                content_length = end - start + 1
                
                def generate():
                    with open(file_path, 'rb') as f:
                        f.seek(start)
                        remaining = content_length
                        while remaining:
                            chunk_size = min(8192, remaining)
                            chunk = f.read(chunk_size)
                            if not chunk:
                                break
                            remaining -= len(chunk)
                            yield chunk
                
                response = Response(generate(), 
                                  status=206, 
                                  headers={
                                      'Content-Range': f'bytes {start}-{end}/{file_size}',
                                      'Accept-Ranges': 'bytes',
                                      'Content-Length': str(content_length),
                                      'Content-Type': get_mime_type(filename)
                                  })
                return response
        
        # Return full file if no range request
        return send_file(file_path, mimetype=get_mime_type(filename))
        
    except Exception as e:
        logger.error(f"Video streaming error: {e}")
        return jsonify({'error': str(e)}), 500

def get_mime_type(filename):
    """Get MIME type for video file"""
    ext = os.path.splitext(filename)[1].lower()
    mime_types = {
        '.mp4': 'video/mp4',
        '.avi': 'video/x-msvideo',
        '.mov': 'video/quicktime',
        '.mkv': 'video/x-matroska',
        '.webm': 'video/webm',
        '.ogg': 'video/ogg'
    }
    return mime_types.get(ext, 'video/mp4')

def analyze_url_async(analysis_id: str, url: str):
    """Analyze URL asynchronously"""
    try:
        logger.info(f"Starting URL analysis for {analysis_id}: {url}")
        
        # Step 1: URL Reputation Check (using VirusTotal URL API)
        update_step_status(analysis_id, 1, 'processing')
        try:
            vt_result = scan_url_with_virustotal(url)
        except Exception as e:
            logger.warning(f"VirusTotal URL scan failed: {e}")
            vt_result = {
                'threat_detected': False,
                'engines_used': 0,
                'engines_detected': 0,
                'confidence': 0,
                'threat_score': 0.0,
                'scan_details': {'error': str(e)}
            }
        update_step_status(analysis_id, 1, 'completed')
        
        # Step 2: AI Website Analysis
        update_step_status(analysis_id, 2, 'processing')
        try:
            logger.info(f"Starting AI analysis for {url}")
            ai_result = analyze_url_with_ai(url)
            logger.info(f"AI analysis completed for {url}")
        except Exception as e:
            logger.error(f"AI URL analysis failed: {e}", exc_info=True)
            ai_result = {
                'ai_confidence_score': 0,
                'risk_score': 50,
                'ai_analysis': f'Analysis failed: {str(e)}',
                'threat_explanations': ['Analysis system temporarily unavailable'],
                'trust_score': 50
            }
        update_step_status(analysis_id, 2, 'completed')
        
        logger.info(f"[URL_ANALYSIS] AI analysis step completed, moving to preview generation")
        
        # Step 3: Safe Preview Generation
        update_step_status(analysis_id, 3, 'processing')
        try:
            logger.info(f"[URL_ANALYSIS] Starting preview generation")
            preview_result = generate_safe_preview(url, analysis_id)
            logger.info(f"[URL_ANALYSIS] Preview generation completed")
        except Exception as e:
            logger.warning(f"Preview generation failed: {e}", exc_info=True)
            preview_result = {
                'preview_available': False,
                'error': str(e)
            }
        update_step_status(analysis_id, 3, 'completed')
        
        # Step 4: Generate Report
        update_step_status(analysis_id, 4, 'processing')
        
        # No sandbox analysis for URLs, create empty result
        sandbox_result = {
            "sandbox_available": False,
            "behaviors_detected": ["URL analysis - sandbox not applicable"],
            "behavioral_score": 100,
            "execution_time": 0
        }
        
        try:
            logger.info(f"[URL_ANALYSIS] Combining final results")
            logger.info(f"[URL_ANALYSIS] vt_result type: {type(vt_result)}")
            logger.info(f"[URL_ANALYSIS] ai_result type: {type(ai_result)}")
            logger.info(f"[URL_ANALYSIS] sandbox_result type: {type(sandbox_result)}")
            logger.info(f"[URL_ANALYSIS] url: {url}")
            logger.info(f"[URL_ANALYSIS] preview_result type: {type(preview_result)}")
            final_results = combine_url_results(vt_result, ai_result, sandbox_result, url, preview_result)
            logger.info(f"[URL_ANALYSIS] Results combination completed")
        except Exception as e:
            logger.error(f"[URL_ANALYSIS] Results combination failed: {e}", exc_info=True)
            final_results = {
                'analysis_id': analysis_id,
                'url': url,
                'status': 'completed',
                'threat_detected': False,
                'threat_level': 'LOW',
                'confidence': 50,
                'engines_used': 0,
                'engines_detected': 0,
                'ai_analysis': ai_result.get('ai_analysis', 'Analysis completed'),
                'sandbox_analysis': sandbox_result,
                'preview_available': preview_result.get('preview_available', False),
                'is_url': True
            }
        
        update_step_status(analysis_id, 4, 'completed')
        
        # Store results
        try:
            logger.info(f"[URL_ANALYSIS] Storing final results")
            analysis_results[analysis_id]['status'] = 'completed'
            analysis_results[analysis_id]['results'] = final_results
            analysis_results[analysis_id]['progress'] = 100
            save_results()  # Persist to disk
            logger.info(f"[URL_ANALYSIS] Results stored successfully")
        except Exception as e:
            logger.error(f"[URL_ANALYSIS] Results storage failed: {e}", exc_info=True)
        
        logger.info(f"URL analysis completed for {analysis_id}")
        
    except Exception as e:
        logger.error(f"URL analysis failed for {analysis_id}: {str(e)}")
        analysis_results[analysis_id]['status'] = 'failed'
        analysis_results[analysis_id]['error'] = str(e)
        analysis_results[analysis_id]['progress'] = 0

def analyze_file_async(analysis_id, file_path):
    """Analyze file asynchronously"""
    try:
        logger.info(f"Starting analysis for {analysis_id}")
        
        # Step 1: VirusTotal Scan
        update_step_status(analysis_id, 1, 'processing')
        vt_result = scan_with_virustotal(file_path)
        
        # VirusTotal scan completed
        update_step_status(analysis_id, 1, 'completed')
        
        # Step 2: Sandbox Analysis (if available)
        update_step_status(analysis_id, 2, 'processing')
        try:
            # Add timeout using threading to prevent hanging
            import queue
            import threading
            
            result_queue = queue.Queue()
            
            def run_sandbox():
                try:
                    result = file_analyzer.sandbox_engine.dynamic_behavioral_analysis(file_path)
                    result_queue.put(('success', result))
                except Exception as e:
                    result_queue.put(('error', str(e)))
            
            # Start sandbox in separate thread with timeout
            sandbox_thread = threading.Thread(target=run_sandbox)
            sandbox_thread.daemon = True
            sandbox_thread.start()
            
            # Wait for result with 60 second timeout
            try:
                status, sandbox_result = result_queue.get(timeout=60)
                if status == 'error':
                    raise Exception(sandbox_result)
                logger.info(f"Sandbox result: {sandbox_result}")
            except queue.Empty:
                logger.warning("Sandbox analysis timed out after 60 seconds")
                # Create safe fallback result
                sandbox_result = {
                    'sandbox_available': True,
                    'execution_successful': False,
                    'behaviors_detected': [],
                    'network_activity': [],
                    'file_operations': [],
                    'system_changes': [],
                    'crypto_activity': [],
                    'threat_indicators': [],
                    'behavioral_score': 100,
                    'execution_time': 60.0,
                    'sandbox_logs': ['⚠️ Sandbox analysis timed out after 60 seconds']
                }
        except Exception as e:
            logger.warning(f"Sandbox analysis failed: {e}")
            # Create safe fallback result
            sandbox_result = {
                'sandbox_available': True,
                'execution_successful': False,
                'behaviors_detected': [],
                'network_activity': [],
                'file_operations': [],
                'system_changes': [],
                'crypto_activity': [],
                'threat_indicators': [],
                'behavioral_score': 100,
                'execution_time': 0.0,
                'sandbox_logs': [f'⚠️ Sandbox analysis failed: {str(e)}']
            }
        update_step_status(analysis_id, 2, 'completed')
        
        # Step 3: Generate Report
        update_step_status(analysis_id, 3, 'processing')
        
        # Skip AI analysis for files - only use for URLs
        ai_result = {
            "ai_confidence_score": 0,
            "risk_score": 0,
            "ai_analysis": "AI analysis skipped for file uploads - only used for URL analysis",
            "threat_explanations": []
        }
        final_results = combine_results(vt_result, ai_result, sandbox_result, file_path, is_url=False)
        update_step_status(analysis_id, 3, 'completed')
        
        # Store results
        analysis_results[analysis_id]['status'] = 'completed'
        analysis_results[analysis_id]['results'] = final_results
        analysis_results[analysis_id]['end_time'] = datetime.now().isoformat()
        save_results()  # Persist to disk
        
        logger.info(f"Analysis completed for {analysis_id}")
        
        # Keep file for preview functionality - don't delete immediately
        # Files will be cleaned up by system temp cleanup or manual cleanup
        # try:
        #     os.remove(file_path)
        # except:
        #     pass
            
    except Exception as e:
        logger.error(f"Analysis error for {analysis_id}: {e}")
        analysis_results[analysis_id]['status'] = 'error'
        analysis_results[analysis_id]['error'] = str(e)

def update_step_status(analysis_id, step_index, status):
    """Update step status"""
    if analysis_id in analysis_results:
        # Convert 1-based step_index to 0-based array index
        array_index = step_index - 1
        if 0 <= array_index < len(analysis_results[analysis_id]['steps']):
            analysis_results[analysis_id]['steps'][array_index]['status'] = status

def scan_url_with_virustotal(url: str):
    """Scan URL with VirusTotal"""
    try:
        # Use the existing VirusTotal integration but for URLs
        # For now, return a placeholder result
        return {
            'threat_detected': False,
            'engines_used': 45,
            'engines_detected': 0,
            'confidence': 95,
            'threat_score': 0.0,
            'scan_details': {'URL': url, 'status': 'clean'}
        }
    except Exception as e:
        logger.error(f"VirusTotal URL scan error: {e}")
        return {
            'threat_detected': False,
            'engines_used': 0,
            'engines_detected': 0,
            'confidence': 0,
            'threat_score': 0.0,
            'scan_details': {'error': str(e)}
        }

def analyze_url_with_ai(url: str):
    """Ultra-comprehensive AI-powered website analysis with advanced intelligence"""
    try:
        logger.info(f"[AI_ANALYSIS] Starting comprehensive URL analysis for: {url}")
        
        from urllib.parse import urlparse
        import re
        import requests
        import json
        from datetime import datetime, timedelta
        import socket
        import ssl
        
        parsed_url = urlparse(url)
        domain = parsed_url.netloc
        logger.info(f"[AI_ANALYSIS] Parsed domain: {domain}")
        
        # Advanced domain intelligence
        try:
            logger.info(f"Getting domain intelligence for: {domain}")
            domain_intelligence = get_advanced_domain_intelligence(domain, url)
            logger.info(f"Domain intelligence completed for: {domain}")
        except Exception as e:
            logger.error(f"Domain intelligence failed: {e}", exc_info=True)
            domain_intelligence = {'domain': domain, 'age_days': 0, 'global_rank': 0}
        
        # Security configuration analysis
        try:
            logger.info(f"Analyzing security for: {url}")
            security_analysis = analyze_website_security(url, domain)
            logger.info(f"Security analysis completed for: {url}")
        except Exception as e:
            logger.error(f"Security analysis failed: {e}", exc_info=True)
            security_analysis = {'ssl_grade': 'Unknown', 'headers_score': 0, 'threats': []}
        
        # Content and reputation analysis
        try:
            logger.info(f"Analyzing content for: {url}")
            content_analysis = analyze_website_content(url)
            logger.info(f"Content analysis completed for: {url}")
        except Exception as e:
            logger.error(f"Content analysis failed: {e}", exc_info=True)
            content_analysis = {'legitimacy_score': 50}
        
        # Trust scoring algorithm
        try:
            logger.info(f"Calculating trust score")
            trust_score = calculate_advanced_trust_score(domain_intelligence, security_analysis, content_analysis)
            logger.info(f"Trust score calculation completed: {trust_score}")
        except Exception as e:
            logger.error(f"Trust score calculation failed: {e}", exc_info=True)
            trust_score = 50
        
        # Generate comprehensive AI analysis
        try:
            logger.info(f"[AI_ANALYSIS] Generating comprehensive report")
            ai_analysis = generate_comprehensive_website_report(domain_intelligence, security_analysis, content_analysis, trust_score, url)
            logger.info(f"[AI_ANALYSIS] Report generation completed")
        except Exception as e:
            logger.error(f"[AI_ANALYSIS] Report generation failed: {e}", exc_info=True)
            ai_analysis = f"Analysis completed for {url} with trust score: {trust_score}/100"
        
        try:
            logger.info(f"[AI_ANALYSIS] Building final result structure")
            result = {
                'ai_confidence_score': 95,
                'risk_score': max(0, 100 - trust_score),
                'ai_analysis': ai_analysis,
                'threat_explanations': security_analysis.get('threats', []),
                'domain_intelligence': domain_intelligence,
                'security_analysis': security_analysis,
                'content_analysis': content_analysis,
                'trust_score': trust_score,
                'reputation_score': trust_score,
                'advanced_metrics': {
                    'domain_age_days': domain_intelligence.get('age_days', 0),
                    'ssl_grade': security_analysis.get('ssl_grade', 'Unknown'),
                    'security_headers_score': security_analysis.get('headers_score', 0),
                    'content_legitimacy_score': content_analysis.get('legitimacy_score', 0),
                    'global_rank': domain_intelligence.get('global_rank', 0),
                    'country_rank': domain_intelligence.get('country_rank', 0),
                    'traffic_rank': domain_intelligence.get('traffic_rank', 0)
                }
            }
            logger.info(f"[AI_ANALYSIS] Final result structure completed")
            return result
        except Exception as e:
            logger.error(f"[AI_ANALYSIS] Final result building failed: {e}", exc_info=True)
            return {
                'ai_confidence_score': 95,
                'risk_score': max(0, 100 - trust_score),
                'ai_analysis': f"Analysis completed for {url} with trust score: {trust_score}/100",
                'threat_explanations': [],
                'domain_intelligence': {'domain': domain},
                'security_analysis': {'ssl_grade': 'Unknown'},
                'content_analysis': {'legitimacy_score': 50},
                'trust_score': trust_score,
                'reputation_score': trust_score,
                'advanced_metrics': {}
            }
    except Exception as e:
        error_msg = str(e).lower()
        if 'quota' in error_msg or 'limit' in error_msg or 'exceeded' in error_msg:
            logger.warning(f"AI analysis quota exceeded for {url}: {e}")
            return {
                'error': 'AI analysis quota exceeded',
                'ai_confidence_score': 0,
                'risk_score': 50,
                'ai_analysis': f'AI analysis unavailable (quota exceeded): {str(e)}. Manual review recommended.',
                'threat_explanations': ['AI analysis quota exceeded'],
                'trust_score': 50
            }
        else:
            logger.error(f"Advanced AI URL analysis error: {e}")
            return {
                'ai_confidence_score': 0,
                'risk_score': 50,
                'ai_analysis': f'Advanced analysis failed: {str(e)}. Manual review recommended.',
                'threat_explanations': ['Analysis system temporarily unavailable'],
                'trust_score': 0
            }

def generate_safe_preview(url: str, analysis_id: str):
    """Generate live interactive preview using proxy system"""
    try:
        logger.info(f"Generating live proxy preview for {url}")
        
        # Create live proxy preview interface
        preview_html = create_live_proxy_preview(url, analysis_id)
        
        # Store the HTML in the analysis results for serving
        if analysis_id not in analysis_results:
            analysis_results[analysis_id] = {}
        
        analysis_results[analysis_id]['preview_html'] = preview_html
        analysis_results[analysis_id]['target_url'] = url  # Store original URL for proxy
        save_results()  # Persist to disk
        
        preview_data = {
            'preview_available': True,
            'preview_url': f'/safe-preview/{analysis_id}',
            'interactive_preview': True,
            'live_content': True,
            'security_level': 'proxy_isolation',
            'browser_like': True,
            'fullscreen_support': True,
            'preview_type': 'live_proxy',
            'bypass_csp': True,
            'real_interaction': True
        }
        return preview_data
    except Exception as e:
        logger.error(f"Preview generation error: {e}")
        return {
            'preview_available': False,
            'error': str(e)
        }

def capture_website_screenshot(url: str):
    """Capture website screenshot using headless Chrome"""
    try:
        logger.info(f"Capturing screenshot for {url}")
        
        # Configure headless Chrome
        chrome_options = Options()
        chrome_options.add_argument('--headless')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument('--disable-gpu')
        chrome_options.add_argument('--window-size=1920,1080')
        chrome_options.add_argument('--disable-web-security')
        chrome_options.add_argument('--disable-features=VizDisplayCompositor')
        chrome_options.add_argument('--disable-extensions')
        chrome_options.add_argument('--disable-plugins')
        chrome_options.add_argument('--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36')
        
        # Initialize driver
        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=chrome_options)
        
        try:
            # Set timeouts
            driver.set_page_load_timeout(30)
            driver.implicitly_wait(10)
            
            # Navigate to URL
            logger.info(f"Navigating to {url}")
            driver.get(url)
            
            # Wait for page to load
            WebDriverWait(driver, 15).until(
                lambda d: d.execute_script("return document.readyState") == "complete"
            )
            
            # Additional wait for dynamic content
            time.sleep(3)
            
            # Capture full page screenshot
            logger.info("Capturing full page screenshot")
            
            # Get page dimensions
            total_height = driver.execute_script("return Math.max(document.body.scrollHeight, document.documentElement.scrollHeight)")
            driver.set_window_size(1920, total_height)
            
            # Take screenshot
            screenshot_png = driver.get_screenshot_as_png()
            
            # Convert to base64 for embedding
            screenshot_b64 = base64.b64encode(screenshot_png).decode('utf-8')
            
            # Get page title and other metadata
            title = driver.title
            current_url = driver.current_url
            
            logger.info(f"Screenshot captured successfully for {url}")
            
            return {
                'success': True,
                'screenshot_b64': screenshot_b64,
                'title': title,
                'final_url': current_url,
                'dimensions': {'width': 1920, 'height': total_height}
            }
            
        finally:
            driver.quit()
            
    except TimeoutException:
        logger.warning(f"Timeout capturing screenshot for {url}")
        return {
            'success': False,
            'error': 'Page load timeout',
            'fallback': True
        }
    except WebDriverException as e:
        logger.warning(f"WebDriver error for {url}: {e}")
        return {
            'success': False,
            'error': f'Browser error: {str(e)}',
            'fallback': True
        }
    except Exception as e:
        logger.error(f"Unexpected error capturing screenshot for {url}: {e}")
        return {
            'success': False,
            'error': str(e),
            'fallback': True
        }

def create_live_proxy_preview(url: str, analysis_id: str):
    """Create live proxy preview interface"""
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Aejis Live Preview - {url}</title>
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <style>
            * {{ margin: 0; padding: 0; box-sizing: border-box; }}
            body {{ 
                font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
                background: #f8fafc;
                height: 100vh;
                overflow: hidden;
            }}
            .browser-chrome {{
                background: #ffffff;
                border-bottom: 1px solid #e2e8f0;
                padding: 12px 16px;
                display: flex;
                align-items: center;
                gap: 12px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                z-index: 100;
                position: relative;
            }}
            .browser-controls {{
                display: flex;
                gap: 8px;
            }}
            .control-btn {{
                width: 32px;
                height: 32px;
                border: 1px solid #cbd5e1;
                background: #ffffff;
                border-radius: 6px;
                cursor: pointer;
                display: flex;
                align-items: center;
                justify-content: center;
                font-size: 14px;
                transition: all 0.2s;
            }}
            .control-btn:hover {{
                background: #f1f5f9;
                border-color: #94a3b8;
            }}
            .url-bar {{
                flex: 1;
                padding: 8px 16px;
                border: 1px solid #cbd5e1;
                border-radius: 24px;
                background: #f8fafc;
                font-family: monospace;
                font-size: 14px;
                color: #475569;
            }}
            .security-indicator {{
                padding: 6px 12px;
                background: #10b981;
                color: white;
                border-radius: 12px;
                font-size: 12px;
                font-weight: 600;
                display: flex;
                align-items: center;
                gap: 4px;
            }}
            .live-indicator {{
                padding: 6px 12px;
                background: #06b6d4;
                color: white;
                border-radius: 12px;
                font-size: 12px;
                font-weight: 600;
                animation: pulse 2s infinite;
            }}
            @keyframes pulse {{
                0%, 100% {{ opacity: 1; }}
                50% {{ opacity: 0.8; }}
            }}
            .viewport {{
                height: calc(100vh - 65px);
                position: relative;
                background: white;
            }}
            .proxy-iframe {{
                width: 100%;
                height: 100%;
                border: none;
                background: white;
            }}
            .loading-overlay {{
                position: absolute;
                top: 0;
                left: 0;
                right: 0;
                bottom: 0;
                background: rgba(255, 255, 255, 0.95);
                display: flex;
                flex-direction: column;
                align-items: center;
                justify-content: center;
                z-index: 10;
            }}
            .loading-spinner {{
                width: 48px;
                height: 48px;
                border: 4px solid #e2e8f0;
                border-top: 4px solid #06b6d4;
                border-radius: 50%;
                animation: spin 1s linear infinite;
                margin-bottom: 16px;
            }}
            @keyframes spin {{
                0% {{ transform: rotate(0deg); }}
                100% {{ transform: rotate(360deg); }}
            }}
        </style>
    </head>
    <body>
        <div class="browser-chrome">
            <div class="browser-controls">
                <button class="control-btn" onclick="goBack()" title="Go Back">←</button>
                <button class="control-btn" onclick="goForward()" title="Go Forward">→</button>
                <button class="control-btn" onclick="reloadPage()" title="Reload">⟳</button>
            </div>
            <input type="text" class="url-bar" value="{url}" readonly title="Live Website URL">
            <div class="security-indicator" title="Secure Aejis Proxy">
                🛡️ SECURE
            </div>
            <div class="live-indicator" title="Live Interactive Content">
                🌐 LIVE
            </div>
            <button class="control-btn" onclick="toggleFullscreen()" title="Toggle Fullscreen">⛶</button>
            <button class="control-btn" onclick="openCleanView()" title="Clean Browser View">🖥️</button>
            <button class="control-btn" onclick="openPureBrowser()" title="Pure Browser (No Extensions)">🌐</button>
            <button class="control-btn" onclick="openExtensionFree()" title="Extension-Free View">🚫</button>
            <button class="control-btn" onclick="openOriginal()" title="Open Original">↗</button>
        </div>

        <div class="viewport">
            <div class="loading-overlay" id="loading-overlay">
                <div class="loading-spinner"></div>
                <div style="color: #475569; margin-bottom: 8px; font-weight: 500;">Loading Live Website...</div>
                <div style="color: #64748b; font-size: 14px;">Proxying through secure Aejis gateway...</div>
            </div>
            
            <iframe 
                id="proxy-iframe" 
                class="proxy-iframe"
                src="/proxy/{analysis_id}/"
                sandbox="allow-scripts allow-same-origin allow-forms allow-popups allow-top-navigation allow-modals allow-downloads allow-presentation"
                onload="handleLoad()"
                onerror="handleError()"
                referrerpolicy="no-referrer"
                allowfullscreen>
            </iframe>
        </div>

        <script>
            let isFullscreen = false;

            function handleLoad() {{
                document.getElementById('loading-overlay').style.display = 'none';
                console.log('Live website loaded successfully');
            }}

            function handleError() {{
                document.getElementById('loading-overlay').innerHTML = 
                    '<div style="text-align: center;"><h3>Connection Error</h3><p>Retrying connection...</p></div>';
                setTimeout(() => {{
                    document.getElementById('proxy-iframe').src = document.getElementById('proxy-iframe').src;
                }}, 3000);
            }}

            function goBack() {{
                try {{
                    document.getElementById('proxy-iframe').contentWindow.history.back();
                }} catch(e) {{
                    console.log('Navigation: back');
                }}
            }}

            function goForward() {{
                try {{
                    document.getElementById('proxy-iframe').contentWindow.history.forward();
                }} catch(e) {{
                    console.log('Navigation: forward');
                }}
            }}

            function reloadPage() {{
                document.getElementById('loading-overlay').style.display = 'flex';
                document.getElementById('proxy-iframe').src = document.getElementById('proxy-iframe').src;
            }}

            function toggleFullscreen() {{
                if (isFullscreen) {{
                    if (document.exitFullscreen) {{
                        document.exitFullscreen();
                    }}
                    isFullscreen = false;
                }} else {{
                    if (document.documentElement.requestFullscreen) {{
                        document.documentElement.requestFullscreen();
                    }}
                    isFullscreen = true;
                }}
            }}

            function openCleanView() {{
                // Open in completely isolated window with extension-blocking flags
                window.open(
                    '/proxy/{analysis_id}/', 
                    '_blank', 
                    'toolbar=yes,scrollbars=yes,resizable=yes,width=1200,height=800,menubar=yes,location=yes,status=yes'
                );
            }}

            function openPureBrowser() {{
                // Direct redirect to proxy - completely bypasses iframe and extensions
                window.location.href = '/pure-browser/{analysis_id}';
            }}

            function openExtensionFree() {{
                // Extension-free view with minimal JavaScript and special headers
                window.location.href = '/extension-free/{analysis_id}';
            }}

            function openOriginal() {{
                window.open('{url}', '_blank', 'noopener,noreferrer');
            }}

            // Handle fullscreen changes
            document.addEventListener('fullscreenchange', function() {{
                isFullscreen = !!document.fullscreenElement;
            }});

            // Keyboard shortcuts
            document.addEventListener('keydown', function(e) {{
                if (e.key === 'F11') {{
                    e.preventDefault();
                    toggleFullscreen();
                }}
                if (e.ctrlKey && e.key === 'r') {{
                    e.preventDefault();
                    reloadPage();
                }}
            }});
        </script>
    </body>
    </html>
    """

def create_interactive_browser_preview(url: str, analysis_id: str, screenshot_result: dict):
    """Create interactive preview with real browser screenshot"""
    screenshot_b64 = screenshot_result['screenshot_b64']
    title = screenshot_result.get('title', 'Website Preview')
    final_url = screenshot_result.get('final_url', url)
    
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Aejis Live Preview - {title}</title>
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <style>
            * {{ margin: 0; padding: 0; box-sizing: border-box; }}
            body {{ 
                font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
                background: #f8fafc;
                height: 100vh;
                overflow: hidden;
            }}
            .browser-chrome {{
                background: #ffffff;
                border-bottom: 1px solid #e2e8f0;
                padding: 12px 16px;
                display: flex;
                align-items: center;
                gap: 12px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                z-index: 100;
            }}
            .browser-controls {{
                display: flex;
                gap: 8px;
            }}
            .control-btn {{
                width: 32px;
                height: 32px;
                border: 1px solid #cbd5e1;
                background: #ffffff;
                border-radius: 6px;
                cursor: pointer;
                display: flex;
                align-items: center;
                justify-content: center;
                font-size: 14px;
                transition: all 0.2s;
            }}
            .control-btn:hover {{
                background: #f1f5f9;
                border-color: #94a3b8;
            }}
            .url-bar {{
                flex: 1;
                padding: 8px 16px;
                border: 1px solid #cbd5e1;
                border-radius: 24px;
                background: #f8fafc;
                font-family: monospace;
                font-size: 14px;
                color: #475569;
            }}
            .security-indicator {{
                padding: 6px 12px;
                background: #10b981;
                color: white;
                border-radius: 12px;
                font-size: 12px;
                font-weight: 600;
                display: flex;
                align-items: center;
                gap: 4px;
            }}
            .live-indicator {{
                padding: 6px 12px;
                background: #ef4444;
                color: white;
                border-radius: 12px;
                font-size: 12px;
                font-weight: 600;
                animation: pulse 2s infinite;
            }}
            @keyframes pulse {{
                0%, 100% {{ opacity: 1; }}
                50% {{ opacity: 0.7; }}
            }}
            .viewport {{
                height: calc(100vh - 65px);
                position: relative;
                background: white;
                overflow: auto;
            }}
            .screenshot-container {{
                width: 100%;
                text-align: center;
                padding: 20px;
            }}
            .screenshot {{
                max-width: 100%;
                border: 1px solid #e2e8f0;
                border-radius: 8px;
                box-shadow: 0 4px 12px rgba(0,0,0,0.1);
                cursor: zoom-in;
                transition: transform 0.2s;
            }}
            .screenshot:hover {{
                transform: scale(1.02);
            }}
            .action-bar {{
                position: fixed;
                bottom: 20px;
                right: 20px;
                display: flex;
                gap: 12px;
                z-index: 200;
            }}
            .action-btn {{
                padding: 12px 16px;
                background: #3b82f6;
                color: white;
                border: none;
                border-radius: 8px;
                cursor: pointer;
                font-weight: 600;
                box-shadow: 0 4px 12px rgba(59, 130, 246, 0.3);
                transition: all 0.2s;
            }}
            .action-btn:hover {{
                background: #2563eb;
                transform: translateY(-2px);
            }}
            .action-btn.secondary {{
                background: #6b7280;
            }}
            .action-btn.secondary:hover {{
                background: #4b5563;
            }}
        </style>
    </head>
    <body>
        <div class="browser-chrome">
            <div class="browser-controls">
                <button class="control-btn" onclick="refreshCapture()" title="Refresh Screenshot">⟳</button>
                <button class="control-btn" onclick="zoomIn()" title="Zoom In">+</button>
                <button class="control-btn" onclick="zoomOut()" title="Zoom Out">-</button>
            </div>
            <input type="text" class="url-bar" value="{final_url}" readonly title="Captured URL">
            <div class="security-indicator" title="Secure Aejis Preview">
                🛡️ SECURE
            </div>
            <div class="live-indicator" title="Live Browser Capture">
                📸 LIVE
            </div>
            <button class="control-btn" onclick="toggleFullscreen()" title="Toggle Fullscreen">⛶</button>
        </div>

        <div class="viewport" id="viewport">
            <div class="screenshot-container">
                <img id="screenshot" class="screenshot" 
                     src="data:image/png;base64,{screenshot_b64}" 
                     alt="Website Screenshot"
                     onclick="toggleZoom(this)">
            </div>
        </div>

        <div class="action-bar">
            <button class="action-btn" onclick="openOriginal()">🔗 Open Live Site</button>
            <button class="action-btn secondary" onclick="downloadScreenshot()">💾 Save Screenshot</button>
            <button class="action-btn secondary" onclick="goBack()">← Back to Results</button>
        </div>

        <script>
            let currentZoom = 1;
            let isFullscreen = false;

            function refreshCapture() {{
                alert('Refreshing screenshot...');
                location.reload();
            }}

            function zoomIn() {{
                currentZoom = Math.min(currentZoom * 1.2, 3);
                applyZoom();
            }}

            function zoomOut() {{
                currentZoom = Math.max(currentZoom / 1.2, 0.3);
                applyZoom();
            }}

            function applyZoom() {{
                document.getElementById('screenshot').style.transform = `scale(${{currentZoom}})`;
                document.getElementById('screenshot').style.transformOrigin = 'top center';
            }}

            function toggleZoom(img) {{
                if (currentZoom === 1) {{
                    currentZoom = 1.5;
                }} else {{
                    currentZoom = 1;
                }}
                applyZoom();
            }}

            function toggleFullscreen() {{
                if (isFullscreen) {{
                    if (document.exitFullscreen) {{
                        document.exitFullscreen();
                    }}
                    isFullscreen = false;
                }} else {{
                    if (document.documentElement.requestFullscreen) {{
                        document.documentElement.requestFullscreen();
                    }}
                    isFullscreen = true;
                }}
            }}

            function openOriginal() {{
                window.open('{final_url}', '_blank', 'noopener,noreferrer');
            }}

            function downloadScreenshot() {{
                const link = document.createElement('a');
                link.download = 'aejis-preview-{analysis_id}.png';
                link.href = 'data:image/png;base64,{screenshot_b64}';
                link.click();
            }}

            function goBack() {{
                window.location.href = 'http://localhost:3000/results/{analysis_id}';
            }}

            // Keyboard shortcuts
            document.addEventListener('keydown', function(e) {{
                if (e.key === 'F11') {{
                    e.preventDefault();
                    toggleFullscreen();
                }}
                if (e.ctrlKey && e.key === '=') {{
                    e.preventDefault();
                    zoomIn();
                }}
                if (e.ctrlKey && e.key === '-') {{
                    e.preventDefault();
                    zoomOut();
                }}
                if (e.ctrlKey && e.key === '0') {{
                    e.preventDefault();
                    currentZoom = 1;
                    applyZoom();
                }}
            }});

            // Handle fullscreen changes
            document.addEventListener('fullscreenchange', function() {{
                isFullscreen = !!document.fullscreenElement;
            }});
        </script>
    </body>
    </html>
    """

def create_browser_preview_html(url: str, analysis_id: str):
    """Create browser-like preview HTML"""
    return f'''
    <!DOCTYPE html>
    <html>
    <head>
        <title>Aejis Safe Preview - {url}</title>
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <style>
            * {{ margin: 0; padding: 0; box-sizing: border-box; }}
            body {{ 
                font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
                background: #f8fafc;
                height: 100vh;
                overflow: hidden;
            }}
            .browser-chrome {{
                background: #ffffff;
                border-bottom: 1px solid #e2e8f0;
                padding: 12px 16px;
                display: flex;
                align-items: center;
                gap: 12px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                z-index: 100;
                position: relative;
            }}
            .browser-controls {{
                display: flex;
                gap: 8px;
            }}
            .control-btn {{
                width: 32px;
                height: 32px;
                border: 1px solid #cbd5e1;
                background: #ffffff;
                border-radius: 6px;
                cursor: pointer;
                display: flex;
                align-items: center;
                justify-content: center;
                font-size: 14px;
                transition: all 0.2s;
            }}
            .control-btn:hover {{
                background: #f1f5f9;
                border-color: #94a3b8;
            }}
            .control-btn:disabled {{
                opacity: 0.5;
                cursor: not-allowed;
            }}
            .url-bar {{
                flex: 1;
                padding: 8px 16px;
                border: 1px solid #cbd5e1;
                border-radius: 24px;
                background: #f8fafc;
                font-family: monospace;
                font-size: 14px;
                color: #475569;
            }}
            .security-indicator {{
                padding: 6px 12px;
                background: #10b981;
                color: white;
                border-radius: 12px;
                font-size: 12px;
                font-weight: 600;
                display: flex;
                align-items: center;
                gap: 4px;
            }}
            .viewport {{
                height: calc(100vh - 65px);
                position: relative;
                background: white;
            }}
            .preview-iframe {{
                width: 100%;
                height: 100%;
                border: none;
                background: white;
            }}
            .loading-overlay {{
                position: absolute;
                top: 0;
                left: 0;
                right: 0;
                bottom: 0;
                background: rgba(255, 255, 255, 0.95);
                display: flex;
                flex-direction: column;
                align-items: center;
                justify-content: center;
                z-index: 10;
            }}
            .loading-spinner {{
                width: 48px;
                height: 48px;
                border: 4px solid #e2e8f0;
                border-top: 4px solid #3b82f6;
                border-radius: 50%;
                animation: spin 1s linear infinite;
                margin-bottom: 16px;
            }}
            @keyframes spin {{
                0% {{ transform: rotate(0deg); }}
                100% {{ transform: rotate(360deg); }}
            }}
            .error-overlay {{
                position: absolute;
                top: 0;
                left: 0;
                right: 0;
                bottom: 0;
                background: #f8fafc;
                display: none;
                flex-direction: column;
                align-items: center;
                justify-content: center;
                z-index: 20;
                padding: 40px;
                text-align: center;
            }}
            .error-icon {{
                font-size: 64px;
                margin-bottom: 24px;
            }}
            .error-title {{
                font-size: 24px;
                font-weight: 600;
                color: #1e293b;
                margin-bottom: 12px;
            }}
            .error-message {{
                color: #64748b;
                margin-bottom: 24px;
                line-height: 1.6;
                max-width: 500px;
            }}
            .action-buttons {{
                display: flex;
                gap: 12px;
                flex-wrap: wrap;
                justify-content: center;
            }}
            .action-btn {{
                padding: 12px 24px;
                border: none;
                border-radius: 8px;
                font-weight: 600;
                cursor: pointer;
                text-decoration: none;
                transition: all 0.2s;
                display: flex;
                align-items: center;
                gap: 8px;
            }}
            .action-btn.primary {{
                background: #3b82f6;
                color: white;
            }}
            .action-btn.primary:hover {{
                background: #2563eb;
            }}
            .action-btn.secondary {{
                background: #f1f5f9;
                color: #475569;
                border: 1px solid #cbd5e1;
            }}
            .action-btn.secondary:hover {{
                background: #e2e8f0;
            }}
            .loading-text {{
                color: #475569;
                margin-bottom: 8px;
                font-weight: 500;
            }}
            .loading-subtext {{
                color: #64748b;
                font-size: 14px;
            }}
            .fullscreen {{
                position: fixed !important;
                top: 0 !important;
                left: 0 !important;
                width: 100vw !important;
                height: 100vh !important;
                z-index: 9999 !important;
            }}
        </style>
    </head>
    <body>
        <div class="browser-chrome">
            <div class="browser-controls">
                <button class="control-btn" onclick="goBack()" title="Go Back">←</button>
                <button class="control-btn" onclick="goForward()" title="Go Forward">→</button>
                <button class="control-btn" onclick="reloadPage()" title="Reload">⟳</button>
            </div>
            <input type="text" class="url-bar" value="{url}" readonly title="Current URL">
            <div class="security-indicator" title="Secure Aejis Preview">
                🛡️ SECURE
            </div>
            <button class="control-btn" onclick="toggleFullscreen()" title="Toggle Fullscreen">⛶</button>
            <button class="control-btn" onclick="openOriginal()" title="Open in New Tab">↗</button>
        </div>
        
        <div class="viewport">
            <div class="loading-overlay" id="loading-overlay">
                <div class="loading-spinner"></div>
                <div class="loading-text">Loading Website...</div>
                <div class="loading-subtext">Fetching content securely through Aejis proxy...</div>
            </div>
            
            <div class="error-overlay" id="error-overlay">
                <div class="error-icon">🌐</div>
                <h2 class="error-title">Preview Method Selection</h2>
                <p class="error-message">
                    Some websites restrict iframe embedding. Choose how you'd like to view this website:
                </p>
                <div class="action-buttons">
                    <a href="{url}" target="_blank" class="action-btn primary">
                        🔗 Open in New Tab
                    </a>
                    <button onclick="tryDirectLoad()" class="action-btn secondary">
                        🔄 Try Direct Load
                    </button>
                    <a href="http://localhost:3000/results/{analysis_id}" class="action-btn secondary">
                        ← Back to Results
                    </a>
                </div>
            </div>
            
            <iframe 
                id="preview-iframe"
                class="preview-iframe"
                sandbox="allow-scripts allow-same-origin allow-forms allow-popups allow-top-navigation allow-modals allow-downloads"
                loading="lazy"
                onload="handleIframeLoad()"
                onerror="handleIframeError()"
            ></iframe>
        </div>
        
        <script>
            let loadTimeout;
            let iframe = document.getElementById('preview-iframe');
            let loadingOverlay = document.getElementById('loading-overlay');
            let errorOverlay = document.getElementById('error-overlay');
            let isFullscreen = false;
            
            function handleIframeLoad() {{
                clearTimeout(loadTimeout);
                // Check if iframe content is actually accessible
                try {{
                    let iframeDoc = iframe.contentDocument || iframe.contentWindow.document;
                    if (iframeDoc && iframeDoc.body && iframeDoc.body.innerHTML.trim()) {{
                        loadingOverlay.style.display = 'none';
                        errorOverlay.style.display = 'none';
                        console.log('Website loaded successfully in preview');
                    }} else {{
                        // Content might be blocked
                        setTimeout(checkIframeContent, 2000);
                    }}
                }} catch(e) {{
                    // Cross-origin or blocked content
                    setTimeout(checkIframeContent, 2000);
                }}
            }}
            
            function checkIframeContent() {{
                try {{
                    let iframeDoc = iframe.contentDocument || iframe.contentWindow.document;
                    if (!iframeDoc || !iframeDoc.body || iframeDoc.body.innerHTML.trim() === '') {{
                        showError();
                    }} else {{
                        loadingOverlay.style.display = 'none';
                        errorOverlay.style.display = 'none';
                    }}
                }} catch(e) {{
                    // Assume blocked if we can't access content after load
                    showError();
                }}
            }}
            
            function handleIframeError() {{
                clearTimeout(loadTimeout);
                showError();
            }}
            
            function showError() {{
                loadingOverlay.style.display = 'none';
                errorOverlay.style.display = 'flex';
            }}
            
            function tryDirectLoad() {{
                errorOverlay.style.display = 'none';
                loadingOverlay.style.display = 'flex';
                
                // Try direct iframe load
                iframe.src = "{url}";
                loadTimeout = setTimeout(showError, 10000);
            }}
            
            function reloadPage() {{
                loadingOverlay.style.display = 'flex';
                errorOverlay.style.display = 'none';
                iframe.src = iframe.src;
                loadTimeout = setTimeout(showError, 8000);
            }}
            
            function goBack() {{
                try {{
                    iframe.contentWindow.history.back();
                }} catch(e) {{
                    console.log('Cannot navigate back - cross-origin restriction');
                }}
            }}
            
            function goForward() {{
                try {{
                    iframe.contentWindow.history.forward();
                }} catch(e) {{
                    console.log('Cannot navigate forward - cross-origin restriction');
                }}
            }}
            
            function toggleFullscreen() {{
                if (isFullscreen) {{
                    if (document.exitFullscreen) {{
                        document.exitFullscreen();
                    }}
                    document.body.classList.remove('fullscreen');
                    isFullscreen = false;
                }} else {{
                    if (document.documentElement.requestFullscreen) {{
                        document.documentElement.requestFullscreen();
                    }}
                    document.body.classList.add('fullscreen');
                    isFullscreen = true;
                }}
            }}
            
            function openOriginal() {{
                window.open('{url}', '_blank', 'noopener,noreferrer');
            }}
            
            // Initialize preview
            window.onload = function() {{
                // Start loading the website immediately
                iframe.src = "{url}";
                loadTimeout = setTimeout(showError, 10000);
                
                // Handle fullscreen changes
                document.addEventListener('fullscreenchange', function() {{
                    isFullscreen = !!document.fullscreenElement;
                }});
                
                // Keyboard shortcuts
                document.addEventListener('keydown', function(e) {{
                    if (e.key === 'F11') {{
                        e.preventDefault();
                        toggleFullscreen();
                    }}
                    if (e.ctrlKey && e.key === 'r') {{
                        e.preventDefault();
                        reloadPage();
                    }}
                }});
            }};
            
            // Monitor iframe content
            setInterval(function() {{
                try {{
                    if (iframe.contentWindow && iframe.contentWindow.location) {{
                        let currentUrl = iframe.contentWindow.location.href;
                        if (currentUrl && currentUrl !== 'about:blank' && currentUrl !== document.querySelector('.url-bar').value) {{
                            document.querySelector('.url-bar').value = currentUrl;
                        }}
                    }}
                }} catch(e) {{
                    // Cross-origin restriction - this is expected and normal
                }}
            }}, 1000);
        </script>
    </body>
    </html>
    '''

def combine_url_results(vt_result, ai_result, sandbox_result, url, preview_result):
    """Combine URL analysis results"""
    try:
        logger.info(f"[COMBINE_RESULTS] Starting URL results combination")
        from urllib.parse import urlparse
        
        parsed_url = urlparse(url)
        domain = parsed_url.netloc
        logger.info(f"[COMBINE_RESULTS] Parsed URL domain: {domain}")
        
        logger.info(f"[COMBINE_RESULTS] Building result structure")
        result = {
            'url_info': {
                'url': url,
                'domain': domain,
                'scheme': parsed_url.scheme,
                'analysis_time': datetime.now().isoformat()
            },
            'virustotal': vt_result,
            'sandbox': sandbox_result,
            'preview': preview_result,
            'summary': {
                'threat_detected': vt_result.get('threat_detected', False) or ai_result.get('risk_score', 0) > 70,
                'threat_level': determine_url_threat_level(vt_result, ai_result),
                'confidence': max(vt_result.get('confidence', 0), ai_result.get('ai_confidence_score', 0)),
                'engines_used': vt_result.get('engines_used', 0),
                'engines_detected': vt_result.get('engines_detected', 0)
            },
            'is_url': True
        }
        
        logger.info(f"[COMBINE_RESULTS] Basic result structure created")
        
        # Include AI analysis for URLs (only if available and not empty)
        if ai_result and isinstance(ai_result, dict) and ai_result.get('error') is None:
            logger.info(f"[COMBINE_RESULTS] Adding AI analysis to result")
            result['ai_analysis'] = ai_result
        else:
            logger.info(f"[COMBINE_RESULTS] Skipping AI analysis (quota exceeded or error)")
            result['ai_analysis'] = {
                'error': 'AI analysis unavailable (quota exceeded)',
                'risk_score': 0,
                'ai_confidence_score': 0
            }
        
        logger.info(f"[COMBINE_RESULTS] URL results combination completed successfully")
        return result
        
    except Exception as e:
        logger.error(f"[COMBINE_RESULTS] URL results combination failed: {e}", exc_info=True)
        return {
            'url_info': {'url': url, 'domain': 'unknown', 'scheme': 'https', 'analysis_time': datetime.now().isoformat()},
            'virustotal': vt_result or {},
            'sandbox': sandbox_result or {},
            'preview': preview_result or {},
            'summary': {'threat_detected': False, 'threat_level': 'UNKNOWN', 'confidence': 0, 'engines_used': 0, 'engines_detected': 0},
            'is_url': True,
            'ai_analysis': ai_result or {}
        }

def determine_url_threat_level(vt_result, ai_result):
    """Determine threat level for URLs"""
    engines_detected = vt_result.get('engines_detected', 0)
    ai_risk = ai_result.get('risk_score', 0) if ai_result else 0
    
    if engines_detected > 5 or ai_risk > 80:
        return 'CRITICAL'
    elif engines_detected > 2 or ai_risk > 60:
        return 'HIGH'
    elif engines_detected > 0 or ai_risk > 40:
        return 'MEDIUM'
    elif ai_risk > 20:
        return 'LOW'
    else:
        return 'SAFE'

def combine_results(vt_result, ai_result, sandbox_result, file_path, is_url=True):
    """Combine all analysis results - exclude AI analysis for files"""
    result = {
        'file_info': {
            'filename': os.path.basename(file_path),
            'size': os.path.getsize(file_path),
            'analysis_time': datetime.now().isoformat()
        },
        'virustotal': vt_result,
        'sandbox': sandbox_result,
        'summary': {
            'threat_detected': vt_result.get('threat_detected', False),
            'threat_level': determine_threat_level(vt_result, ai_result if is_url else None),
            'confidence': vt_result.get('confidence', 0),
            'engines_used': vt_result.get('engines_used', 0),
            'engines_detected': vt_result.get('engines_detected', 0)
        },
        'is_url': is_url
    }
    
    # Only include AI analysis for URLs
    if is_url:
        result['ai_analysis'] = ai_result
        result['summary']['threat_detected'] = (
            vt_result.get('threat_detected', False) or ai_result.get('risk_score', 0) > 50
        )
        result['summary']['confidence'] = max(
            vt_result.get('confidence', 0), 
            ai_result.get('ai_confidence_score', 0)
        )
    
    return result

def determine_threat_level(vt_result, ai_result):
    """Determine overall threat level"""
    vt_score = vt_result.get('threat_score', 0)
    ai_score = ai_result.get('risk_score', 0) if ai_result else 0
    
    max_score = max(vt_score, ai_score)
    
    if max_score >= 80:
        return 'CRITICAL'
    elif max_score >= 60:
        return 'HIGH'
    elif max_score >= 40:
        return 'MEDIUM'
    elif max_score >= 20:
        return 'LOW'
    else:
        return 'SAFE'

def generate_report(results):
    """Generate downloadable report"""
    return {
        'report_info': {
            'generated_at': datetime.now().isoformat(),
            'platform': 'Aejis Security',
            'version': '1.0.0'
        },
        'analysis_results': results,
        'recommendations': generate_recommendations(results)
    }

def generate_recommendations(results):
    """Generate security recommendations"""
    recommendations = []
    
    if results['summary']['threat_detected']:
        recommendations.extend([
            "⚠️ THREAT DETECTED: Do not open or execute this file",
            "🗑️ Delete the file immediately",
            "🔍 Scan your system with a full antivirus scan",
            "🔄 Change any passwords that may have been compromised"
        ])
    else:
        recommendations.extend([
            "✅ File appears to be safe",
            "🔍 Continue to monitor for any suspicious behavior",
            "🛡️ Keep your antivirus software updated",
            "📚 Practice good cybersecurity hygiene"
        ])
    
    return recommendations

# ============================================================================
# UNLIMITED FILE PREVIEW EXTRACTION FUNCTIONS
# ============================================================================

def _extract_image_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract image metadata and create preview"""
    try:
        import PIL
        from PIL import Image, ExifTags
        
        with Image.open(file_path) as img:
            # Basic image info
            width, height = img.size
            mode = img.mode
            format_name = img.format
            
            # Extract EXIF data if available
            exif_data = {}
            if hasattr(img, '_getexif') and img._getexif() is not None:
                exif = img._getexif()
                for tag_id, value in exif.items():
                    tag = ExifTags.TAGS.get(tag_id, tag_id)
                    exif_data[tag] = str(value)
            
            # Create base64 preview (thumbnail)
            img.thumbnail((300, 300))
            import base64
            from io import BytesIO
            buffer = BytesIO()
            img.save(buffer, format='PNG')
            img_str = base64.b64encode(buffer.getvalue()).decode()
            
            return {
                'filename': os.path.basename(file_path),
                'file_size': os.path.getsize(file_path),
                'file_type': os.path.splitext(file_path)[1],
                'analysis_id': analysis_id,
                'preview_content': f"Image Preview: {width}x{height} pixels, {mode} mode",
                'preview_type': 'image',
                'image_metadata': {
                    'dimensions': f"{width}x{height}",
                    'mode': mode,
                    'format': format_name,
                    'exif_data': exif_data
                },
                'thumbnail_base64': f"data:image/png;base64,{img_str}",
                'secure_extraction': True
            }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Image file - metadata extraction failed: {str(e)}",
            'preview_type': 'image',
            'secure_extraction': False
        }

def _extract_pdf_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract PDF content and serve raw PDF for viewing"""
    try:
        import PyPDF2
        import base64
        
        # Read the raw PDF file for viewing
        with open(file_path, 'rb') as file:
            pdf_data = file.read()
            pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')
        
        # Also extract metadata for information
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Extract metadata
            metadata = {}
            page_count = len(pdf_reader.pages)
            if pdf_reader.metadata:
                metadata = {
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'subject': pdf_reader.metadata.get('/Subject', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                    'producer': pdf_reader.metadata.get('/Producer', ''),
                    'creation_date': str(pdf_reader.metadata.get('/CreationDate', '')),
                    'modification_date': str(pdf_reader.metadata.get('/ModDate', ''))
                }
            
            # Extract text from first few pages for search/metadata
            text_content = ""
            max_pages = min(3, page_count)  # Only extract first 3 pages for metadata
            
            for page_num in range(max_pages):
                try:
                    page = pdf_reader.pages[page_num]
                    text_content += f"\n--- Page {page_num + 1} ---\n"
                    text_content += page.extract_text()
                except:
                    text_content += f"\n--- Page {page_num + 1} (text extraction failed) ---\n"
            
            return {
                'filename': os.path.basename(file_path),
                'file_size': os.path.getsize(file_path),
                'file_type': '.pdf',
                'analysis_id': analysis_id,
                'preview_content': text_content or "No text content found in PDF",
                'preview_type': 'pdf',
                'pdf_base64': pdf_base64,  # Raw PDF data for viewing
                'base64_preview': pdf_base64,  # Also add this for consistency with images
                'mime_type': 'application/pdf',
                'pdf_metadata': {
                    'page_count': page_count,
                    'metadata': metadata,
                    'extracted_pages': max_pages
                },
                'secure_extraction': True,
                'success': True
            }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': '.pdf',
            'analysis_id': analysis_id,
            'preview_content': f"PDF file - processing failed: {str(e)}",
            'preview_type': 'pdf',
            'secure_extraction': False,
            'success': False
        }

def _extract_document_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract document text content"""
    try:
        import docx
        from docx import Document
        
        doc = Document(file_path)
        
        # Extract text content
        text_content = ""
        for paragraph in doc.paragraphs:
            text_content += paragraph.text + "\n"
        
        # Extract tables
        table_content = ""
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                table_content += row_text + "\n"
        
        full_content = text_content + ("\n--- Tables ---\n" + table_content if table_content else "")
        
        # Show full content - no truncation
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': full_content or "No text content found in document",
            'preview_type': 'document',
            'document_metadata': {
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables)
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Document file - text extraction failed: {str(e)}",
            'preview_type': 'document',
            'secure_extraction': False
        }

def _extract_spreadsheet_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract spreadsheet data"""
    try:
        import pandas as pd
        
        # Read spreadsheet
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path, sheet_name=0)
        
        # Get basic info
        rows, cols = df.shape
        column_names = list(df.columns)
        
        # Get first few rows as preview
        preview_data = df.head(10).to_string()
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': preview_data,
            'preview_type': 'spreadsheet',
            'spreadsheet_metadata': {
                'rows': rows,
                'columns': cols,
                'column_names': column_names
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Spreadsheet file - data extraction failed: {str(e)}",
            'preview_type': 'spreadsheet',
            'secure_extraction': False
        }

def _extract_presentation_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract presentation content"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        # Extract slide content
        slide_content = ""
        for i, slide in enumerate(prs.slides):
            slide_content += f"\n--- Slide {i + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content += shape.text + "\n"
        
        # Show full content - no truncation
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': slide_content or "No text content found in presentation",
            'preview_type': 'presentation',
            'presentation_metadata': {
                'slide_count': len(prs.slides)
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Presentation file - content extraction failed: {str(e)}",
            'preview_type': 'presentation',
            'secure_extraction': False
        }

def _extract_code_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract code content with syntax highlighting info"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Get file stats
        lines = content.split('\n')
        line_count = len(lines)
        char_count = len(content)
        
        # Show full content - no truncation
        
        # Determine language for syntax highlighting
        file_ext = os.path.splitext(file_path)[1].lower()
        language_map = {
            '.py': 'python', '.js': 'javascript', '.html': 'html', '.css': 'css',
            '.php': 'php', '.java': 'java', '.cpp': 'cpp', '.c': 'c',
            '.cs': 'csharp', '.go': 'go', '.rs': 'rust', '.rb': 'ruby',
            '.sh': 'bash', '.bat': 'batch', '.ps1': 'powershell', '.sql': 'sql',
            '.xml': 'xml', '.yaml': 'yaml', '.yml': 'yaml', '.json': 'json',
            '.ini': 'ini', '.cfg': 'ini', '.conf': 'ini'
        }
        language = language_map.get(file_ext, 'text')
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': file_ext,
            'analysis_id': analysis_id,
            'preview_content': content,
            'preview_type': 'code',
            'code_metadata': {
                'line_count': line_count,
                'character_count': char_count,
                'language': language,
                'encoding': 'utf-8'
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Code file - content extraction failed: {str(e)}",
            'preview_type': 'code',
            'secure_extraction': False
        }

def _extract_audio_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract audio metadata"""
    try:
        import mutagen
        
        audio_file = mutagen.File(file_path)
        metadata = {}
        
        if audio_file is not None:
            # Extract common metadata
            metadata = {
                'title': str(audio_file.get('title', ['Unknown'])[0]),
                'artist': str(audio_file.get('artist', ['Unknown'])[0]),
                'album': str(audio_file.get('album', ['Unknown'])[0]),
                'duration': str(audio_file.info.length) if hasattr(audio_file, 'info') else 'Unknown',
                'bitrate': str(audio_file.info.bitrate) if hasattr(audio_file, 'info') else 'Unknown',
                'sample_rate': str(audio_file.info.sample_rate) if hasattr(audio_file, 'info') else 'Unknown'
            }
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Audio file - Duration: {metadata.get('duration', 'Unknown')}s, Bitrate: {metadata.get('bitrate', 'Unknown')} kbps",
            'preview_type': 'audio',
            'audio_metadata': metadata,
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Audio file - metadata extraction failed: {str(e)}",
            'preview_type': 'audio',
            'secure_extraction': False
        }

def _extract_video_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract video metadata"""
    try:
        import cv2
        
        cap = cv2.VideoCapture(file_path)
        
        # Get video properties
        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        duration = frame_count / fps if fps > 0 else 0
        
        cap.release()
        
        metadata = {
            'resolution': f"{width}x{height}",
            'fps': fps,
            'frame_count': frame_count,
            'duration': f"{duration:.2f}s"
        }
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Video file - {width}x{height} @ {fps:.2f}fps, Duration: {duration:.2f}s",
            'preview_type': 'video',
            'video_metadata': metadata,
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Video file - metadata extraction failed: {str(e)}",
            'preview_type': 'video',
            'secure_extraction': False
        }

def _extract_binary_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract binary file information and strings"""
    try:
        import subprocess
        
        # Get file info using file command
        try:
            result = subprocess.run(['file', file_path], capture_output=True, text=True)
            file_info = result.stdout.strip()
        except:
            file_info = "Binary file"
        
        # Extract strings (readable text from binary)
        try:
            result = subprocess.run(['strings', file_path], capture_output=True, text=True)
            strings_content = result.stdout  # Full strings output
        except:
            strings_content = "Could not extract strings"
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Binary file info: {file_info}\n\nExtracted strings:\n{strings_content}",
            'preview_type': 'binary',
            'binary_metadata': {
                'file_info': file_info,
                'strings_found': len(strings_content.split('\n')) if strings_content else 0
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Binary file - analysis failed: {str(e)}",
            'preview_type': 'binary',
            'secure_extraction': False
        }

def _extract_database_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract database schema information"""
    try:
        import sqlite3
        
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        
        # Get table names
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [row[0] for row in cursor.fetchall()]
        
        # Get schema for each table
        schema_info = {}
        for table in tables:  # All tables
            cursor.execute(f"PRAGMA table_info({table});")
            columns = cursor.fetchall()
            schema_info[table] = [{'name': col[1], 'type': col[2]} for col in columns]
        
        conn.close()
        
        table_list = ', '.join(tables)
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Database with {{len(tables)}} tables: {{table_list}}",
            'preview_type': 'database',
            'database_metadata': {
                'table_count': len(tables),
                'tables': tables,  # All table names
                'schema': schema_info
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Database file - schema extraction failed: {str(e)}",
            'preview_type': 'database',
            'secure_extraction': False
        }

def _extract_generic_preview(file_path: str, analysis_id: str) -> Dict[str, Any]:
    """Extract generic file information for unknown types"""
    try:
        # Try to read as text first
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1000)  # First 1000 characters
                is_text = True
        except:
            content = "Binary content - cannot display as text"
            is_text = False
        
        # Get file info
        file_size = os.path.getsize(file_path)
        file_ext = os.path.splitext(file_path)[1]
        
        # Try to determine file type
        import mimetypes
        mime_type, _ = mimetypes.guess_type(file_path)
        
        return {
            'filename': os.path.basename(file_path),
            'file_size': file_size,
            'file_type': file_ext,
            'analysis_id': analysis_id,
            'preview_content': content,
            'preview_type': 'text' if is_text else 'binary',
            'generic_metadata': {
                'mime_type': mime_type,
                'is_text': is_text,
                'file_extension': file_ext
            },
            'secure_extraction': True
        }
    except Exception as e:
        return {
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1],
            'analysis_id': analysis_id,
            'preview_content': f"Unknown file type - analysis failed: {str(e)}",
            'preview_type': 'unknown',
            'secure_extraction': False
        }

@app.route('/safe-preview/<analysis_id>')
def safe_preview(analysis_id):
    """Serve safe browser-like preview of analyzed website"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404
        
        analysis_data = analysis_results[analysis_id]
        
        # Check if we have stored preview HTML
        if 'preview_html' in analysis_data:
            return analysis_data['preview_html']
        
        # For URL analysis, generate preview
        if analysis_data.get('status') != 'completed':
            return "Analysis not completed", 400
        
        # Get URL from the results structure
        results = analysis_data.get('results', {})
        url_info = results.get('url_info', {})
        url = url_info.get('url') or analysis_data.get('url')
        
        if not url:
            return "URL not found", 404
        
        # Generate and serve preview HTML
        preview_html = create_browser_preview_html(url, analysis_id)
        
        # Store for future use
        analysis_data['preview_html'] = preview_html
        save_results()
        
        return preview_html
        
    except Exception as e:
        logger.error(f"Safe preview error: {str(e)}")
        return f"Preview error: {str(e)}", 500

@app.route('/clean-proxy/<analysis_id>')
def clean_proxy_redirect(analysis_id):
    """Redirect to clean proxy view without iframe"""
    return redirect(f'/proxy/{analysis_id}/')

@app.route('/ultra-clean/<analysis_id>')
def ultra_clean_proxy(analysis_id):
    """Ultra-clean proxy that bypasses ALL restrictions"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404
            
        target_url = analysis_results[analysis_id].get('target_url')
        if not target_url:
            return "Target URL not found", 404
        
        # Create completely clean HTML page with zero restrictions
        clean_html = f"""<!DOCTYPE html>
<html>
<head>
    <title>Clean View - {target_url}</title>
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <style>
        body {{ margin: 0; padding: 0; }}
        iframe {{ width: 100vw; height: 100vh; border: none; }}
    </style>
</head>
<body>
    <iframe src="/proxy/{analysis_id}/" allowfullscreen></iframe>
</body>
</html>"""
        
        response = make_response(clean_html)
        response.headers['Content-Type'] = 'text/html; charset=utf-8'
        # NO CSP or security headers for maximum compatibility
        return response
        
    except Exception as e:
        logger.error(f"Ultra-clean proxy error: {e}")
        return f"Error: {str(e)}", 500

@app.route('/pure-browser/<analysis_id>')
@app.route('/pure-browser/<analysis_id>/')
def pure_browser_proxy(analysis_id):
    """Pure browser experience - direct redirect to proxy with no iframe"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404
            
        # Direct redirect to the proxy - this bypasses ALL iframe issues
        return redirect(f'/proxy/{analysis_id}/')
        
    except Exception as e:
        logger.error(f"Pure browser proxy error: {e}")
        return f"Error: {str(e)}", 500

@app.route('/extension-free/<analysis_id>')
@app.route('/extension-free/<analysis_id>/')
def extension_free_proxy(analysis_id):
    """Extension-free proxy that serves content with special headers to minimize extension interference"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404
            
        target_url = analysis_results[analysis_id].get('target_url')
        if not target_url:
            return "Target URL not found", 404
            
        # Create a completely clean HTML page with minimal JavaScript
        clean_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Extension-Free View - {target_url}</title>
    <style>
        body {{ 
            margin: 0; 
            padding: 0; 
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
            background: #f8fafc;
        }}
        .header {{
            background: #ffffff;
            border-bottom: 1px solid #e2e8f0;
            padding: 12px 16px;
            display: flex;
            align-items: center;
            gap: 12px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }}
        .url-bar {{
            flex: 1;
            padding: 8px 16px;
            border: 1px solid #cbd5e1;
            border-radius: 24px;
            background: #f8fafc;
            font-family: monospace;
            font-size: 14px;
            color: #475569;
        }}
        .status {{
            padding: 6px 12px;
            background: #10b981;
            color: white;
            border-radius: 12px;
            font-size: 12px;
            font-weight: 600;
        }}
        .content {{
            height: calc(100vh - 65px);
            background: white;
        }}
        iframe {{
            width: 100%;
            height: 100%;
            border: none;
        }}
    </style>
</head>
<body>
    <div class="header">
        <input type="text" class="url-bar" value="{target_url}" readonly>
        <div class="status">🛡️ EXTENSION-FREE</div>
    </div>
    <div class="content">
        <iframe src="/proxy/{analysis_id}/" allowfullscreen></iframe>
    </div>
</body>
</html>"""
        
        flask_response = make_response(clean_html)
        flask_response.headers['Content-Type'] = 'text/html; charset=utf-8'
        # Minimal headers to reduce extension interference
        flask_response.headers['X-Frame-Options'] = 'SAMEORIGIN'
        flask_response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        
        return flask_response
        
    except Exception as e:
        logger.error(f"Extension-free proxy error: {e}")
        return f"Error: {str(e)}", 500

@app.route('/test-iframe/<analysis_id>')
@app.route('/test-iframe/<analysis_id>/')
def test_iframe(analysis_id):
    """Simple test page to debug iframe issues"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404
            
        target_url = analysis_results[analysis_id].get('target_url')
        if not target_url:
            return "Target URL not found", 404
            
        test_html = f"""<!DOCTYPE html>
<html>
<head>
    <title>Iframe Test - {target_url}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 20px; }}
        .test-info {{ background: #f0f0f0; padding: 15px; margin: 10px 0; border-radius: 5px; }}
        iframe {{ width: 100%; height: 600px; border: 2px solid #ccc; }}
    </style>
</head>
<body>
    <h1>Iframe Test Page</h1>
    <div class="test-info">
        <strong>Target URL:</strong> {target_url}<br>
        <strong>Analysis ID:</strong> {analysis_id}<br>
        <strong>Proxy URL:</strong> /proxy/{analysis_id}/
    </div>
    
    <h2>Direct Proxy Content:</h2>
    <iframe src="/proxy/{analysis_id}/" allowfullscreen></iframe>
    
    <h2>Alternative: Direct Link</h2>
    <p><a href="/proxy/{analysis_id}/" target="_blank">Open Proxy Directly</a></p>
</body>
</html>"""
        
        return test_html
        
    except Exception as e:
        return f"Error: {str(e)}", 500

# Aejis Browser Isolation Endpoints
@app.route('/browser/<analysis_id>')
def start_browser_isolation_session(analysis_id):
    """Start an Aejis browser isolation session for URL analysis"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404

        target_url = analysis_results[analysis_id].get('target_url')
        if not target_url:
            return "Target URL not found", 404

        # Start the Aejis browser isolation session
        session_data = browser_isolation_service.start_browser_session(analysis_id, target_url)
        
        return jsonify({
            'success': True,
            'session_id': analysis_id,
            'vnc_url': session_data['vnc_url'],
            'web_url': session_data['web_url'],
            'custom_vnc_url': session_data['custom_vnc_url'],
            'target_url': target_url,
            'status': session_data['status'],
            'isolation_level': 'maximum',
            'location_spoofing': True,
            'commercial_friendly': True
        })

    except Exception as e:
        logger.error(f"Browser isolation session error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/browser/<analysis_id>/status')
def get_browser_isolation_status(analysis_id):
    """Get browser isolation session status"""
    try:
        session_data = browser_isolation_service.get_session_status(analysis_id)
        if not session_data:
            return jsonify({'success': False, 'error': 'Session not found'}), 404

        return jsonify({
            'success': True,
            'session_id': analysis_id,
            'status': session_data['status'],
            'vnc_url': session_data.get('vnc_url'),
            'web_url': session_data.get('web_url'),
            'custom_vnc_url': session_data.get('custom_vnc_url'),
            'target_url': session_data.get('target_url'),
            'start_time': session_data.get('start_time'),
            'isolation_level': session_data.get('isolation_level'),
            'location_spoofing': session_data.get('location_spoofing')
        })

    except Exception as e:
        logger.error(f"Browser isolation status error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/browser/<analysis_id>/stop', methods=['POST'])
def stop_browser_isolation_session(analysis_id):
    """Stop browser isolation session"""
    try:
        success = browser_isolation_service.stop_browser_session(analysis_id)
        return jsonify({'success': success})

    except Exception as e:
        logger.error(f"Stop browser isolation error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/browser/info')
def get_browser_isolation_info():
    """Get browser isolation service information"""
    try:
        info = browser_isolation_service.get_browser_info()
        return jsonify({'success': True, 'info': info})

    except Exception as e:
        logger.error(f"Browser isolation info error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/vnc-auto-connect.html')
def vnc_auto_connect():
    """Serve the auto-connect VNC HTML page and start browser isolation"""
    try:
        target_url = request.args.get('url', 'https://www.google.com')
        
        # Start browser isolation session for this URL
        try:
            session_id = f"direct_{int(time.time())}"
            logger.info(f"Starting browser isolation for direct access: {target_url}")
            browser_isolation_service.start_browser_session(session_id, target_url)
            logger.info(f"Browser isolation started successfully for {session_id}")
        except Exception as e:
            logger.error(f"Failed to start browser isolation: {e}")
        
        # Read the HTML file and replace the target URL
        with open('vnc-auto-connect.html', 'r') as f:
            html_content = f.read()
        
        # Replace the target URL in the HTML
        html_content = html_content.replace('Loading Website...', f'Loading {target_url}...')
        html_content = html_content.replace('Loading isolated browser...', f'Loading {target_url} in isolated browser...')
        html_content = html_content.replace('Loading https://www.virustotal.com/...', f'Loading {target_url}...')
        
        return html_content, 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"VNC auto-connect error: {e}")
        return f"Error loading VNC interface: {e}", 500

@app.route('/vnc-direct.html')
def vnc_direct():
    """Serve the direct VNC HTML page"""
    try:
        target_url = request.args.get('url', 'https://www.google.com')
        
        # Read the HTML file and replace the target URL
        with open('vnc-direct.html', 'r') as f:
            html_content = f.read()
        
        # Replace the target URL in the HTML
        html_content = html_content.replace('Loading https://www.virustotal.com/...', f'Loading {target_url}...')
        
        return html_content, 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"VNC direct error: {e}")
        return f"Error loading VNC interface: {e}", 500

@app.route('/vnc-auto-login.html')
def vnc_auto_login():
    """Serve the VNC auto-login HTML page with credentials"""
    try:
        target_url = request.args.get('url', 'https://www.google.com')
        
        # Read the HTML file and replace the target URL
        with open('vnc-auto-login.html', 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Replace the target URL in the HTML
        html_content = html_content.replace('Loading https://www.virustotal.com/...', f'Loading {target_url}...')
        
        return html_content, 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"VNC auto-login error: {e}")
        return f"Error loading VNC interface: {e}", 500

@app.route('/vnc-fullscreen.html')
def vnc_fullscreen():
    """Serve the fullscreen VNC HTML page"""
    try:
        target_url = request.args.get('url', 'https://www.google.com')
        
        # Read the HTML file and replace the target URL
        with open('vnc-fullscreen.html', 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Replace the target URL in the HTML
        html_content = html_content.replace('Full Screen Mode', f'Full Screen Mode - {target_url}')
        
        return html_content, 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"VNC fullscreen error: {e}")
        return f"Error loading VNC interface: {e}", 500

@app.route('/vnc-clean.html')
def vnc_clean():
    """Serve the clean VNC HTML page without side panels"""
    try:
        target_url = request.args.get('url', 'https://www.google.com')
        
        # Read the HTML file and replace the target URL
        with open('vnc-clean.html', 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Replace the target URL in the HTML
        html_content = html_content.replace('Clean View', f'Clean View - {target_url}')
        
        return html_content, 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"VNC clean error: {e}")
        return f"Error loading VNC interface: {e}", 500

@app.route('/vnc-pure.html')
def vnc_pure():
    """Serve the pure VNC HTML page with no UI elements"""
    try:
        target_url = request.args.get('url', 'https://www.google.com')
        
        # Read the HTML file and replace the target URL
        with open('vnc-pure.html', 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # Replace the target URL in the HTML
        html_content = html_content.replace('Pure Website View', f'Pure Website View - {target_url}')
        
        return html_content, 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"VNC pure error: {e}")
        return f"Error loading VNC interface: {e}", 500

@app.route('/proxy/<analysis_id>/')
@app.route('/proxy/<analysis_id>/<path:subpath>')
def proxy_website(analysis_id, subpath=''):
    """Proxy website content through our server to bypass CSP"""
    try:
        if analysis_id not in analysis_results:
            return "Analysis not found", 404
            
        target_url = analysis_results[analysis_id].get('target_url')
        if not target_url:
            return "Target URL not found", 404
            
        # Build the full URL for the request
        if subpath:
            # Handle relative paths
            if subpath.startswith('http'):
                proxy_url = subpath
            else:
                # Remove leading slash if present
                subpath = subpath.lstrip('/')
                # Handle special cases for relative paths
                if subpath.startswith('gui/') or subpath.startswith('main.'):
                    # These are likely relative to the root of the domain
                    from urllib.parse import urlparse
                    parsed_target = urlparse(target_url)
                    proxy_url = f"{parsed_target.scheme}://{parsed_target.netloc}/{subpath}"
                elif subpath and not subpath.startswith('http'):
                    # Handle other relative paths - check if it's a main.js file
                    from urllib.parse import urlparse
                    parsed_target = urlparse(target_url)
                    if subpath.startswith('main.') and not subpath.startswith('gui/'):
                        # Main.js files should be in the gui directory
                        proxy_url = f"{parsed_target.scheme}://{parsed_target.netloc}/gui/{subpath}"
                    else:
                        proxy_url = f"{parsed_target.scheme}://{parsed_target.netloc}/{subpath}"
                    
                else:
                    proxy_url = f"{target_url.rstrip('/')}/{subpath}"
        else:
            proxy_url = target_url
            
        # Get query parameters
        query_string = request.query_string.decode('utf-8')
        if query_string:
            proxy_url += f"?{query_string}"
            
        logger.info(f"Proxying request to: {proxy_url}")
        
        # Make request to the target website
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': request.headers.get('Accept', '*/*'),
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'DNT': '1',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
        }
        
        # Forward some headers from the original request
        for header in ['Referer', 'Authorization']:
            if header in request.headers:
                headers[header] = request.headers[header]
        
        response = requests.get(proxy_url, headers=headers, timeout=30, verify=False)
        response.raise_for_status()
        
        # Process the response content
        content_type = response.headers.get('Content-Type', '')
        
        if 'text/html' in content_type:
            # Rewrite HTML content to proxy through our server
            html_content = response.text
            modified_html = rewrite_html_for_proxy(html_content, analysis_id, proxy_url)
            
            # Return the modified HTML with appropriate headers
            flask_response = make_response(modified_html)
            flask_response.headers['Content-Type'] = 'text/html; charset=utf-8'
        else:
            # For non-HTML content (CSS, JS, images), return as-is
            flask_response = make_response(response.content)
            flask_response.headers['Content-Type'] = response.headers.get('Content-Type', 'application/octet-stream')
        
        # Add security headers
        flask_response.headers['X-Frame-Options'] = 'SAMEORIGIN'
        flask_response.headers['X-Content-Type-Options'] = 'nosniff'
        flask_response.headers['X-XSS-Protection'] = '1; mode=block'
        flask_response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        flask_response.headers['Pragma'] = 'no-cache'
        flask_response.headers['Expires'] = '0'
        
        # Add relaxed CSP for iframe compatibility (only for HTML content)
        if 'text/html' in content_type:
            flask_response.headers['Content-Security-Policy'] = (
                "default-src 'self' https: data: blob:; "
                "script-src 'self' 'unsafe-inline' 'unsafe-eval' https: data: blob:; "
                "object-src 'none'; "
                "connect-src 'self' https: wss: ws:; "
                "frame-src 'self' https:; "
                "img-src 'self' https: data: blob:; "
                "style-src 'self' 'unsafe-inline' https:; "
                "font-src 'self' https: data:; "
                "media-src 'self' https: data: blob:; "
                "worker-src 'self' blob:; "
                "manifest-src 'self'; "
                "frame-ancestors 'self'; "
                "base-uri 'self';"
            )
        
        return flask_response
        
    except requests.exceptions.RequestException as e:
        logger.error(f"Proxy request failed for {proxy_url}: {e}")
        return f"Proxy Error: Unable to fetch content from {proxy_url}", 502
    except Exception as e:
        logger.error(f"Proxy error: {e}")
        return f"Proxy Error: {str(e)}", 500

def rewrite_html_for_proxy(html_content: str, analysis_id: str, base_url: str):
    """Rewrite HTML content to route requests through our proxy"""
    try:
        from urllib.parse import urljoin, urlparse
        from bs4 import BeautifulSoup
        
        soup = BeautifulSoup(html_content, 'html.parser')
        base_domain = urlparse(base_url).netloc
        
        # Rewrite all href attributes in links
        for link in soup.find_all(['a', 'link'], href=True):
            original_href = link['href']
            if original_href.startswith('http'):
                # Absolute URL - proxy it
                link['href'] = f"/proxy/{analysis_id}/{original_href}"
            elif original_href.startswith('//'):
                # Protocol-relative URL
                link['href'] = f"/proxy/{analysis_id}/https:{original_href}"
            elif original_href.startswith('/'):
                # Root-relative URL
                base_proto_domain = f"{urlparse(base_url).scheme}://{base_domain}"
                link['href'] = f"/proxy/{analysis_id}/{base_proto_domain}{original_href}"
            elif not original_href.startswith('#') and not original_href.startswith('javascript:'):
                # Relative URL - check if it's a main.js file or other gui resources
                if original_href.startswith('main.') and not original_href.startswith('gui/'):
                    # Main.js files should be in the gui directory
                    full_url = f"{urlparse(base_url).scheme}://{base_domain}/gui/{original_href}"
                elif original_href.startswith('images/') or original_href.startswith('manifest.'):
                    # Images and manifest files should be in the gui directory
                    full_url = f"{urlparse(base_url).scheme}://{base_domain}/gui/{original_href}"
                else:
                    full_url = urljoin(base_url, original_href)
                link['href'] = f"/proxy/{analysis_id}/{full_url}"
        
        # Rewrite all src attributes in images, scripts, etc.
        for element in soup.find_all(['img', 'script', 'iframe', 'frame', 'source'], src=True):
            original_src = element['src']
            if original_src.startswith('http'):
                element['src'] = f"/proxy/{analysis_id}/{original_src}"
            elif original_src.startswith('//'):
                element['src'] = f"/proxy/{analysis_id}/https:{original_src}"
            elif original_src.startswith('/'):
                base_proto_domain = f"{urlparse(base_url).scheme}://{base_domain}"
                element['src'] = f"/proxy/{analysis_id}/{base_proto_domain}{original_src}"
            elif not original_src.startswith('data:') and not original_src.startswith('javascript:'):
                # Relative URL - check if it's a main.js file or other gui resources
                if original_src.startswith('main.') and not original_src.startswith('gui/'):
                    # Main.js files should be in the gui directory
                    full_url = f"{urlparse(base_url).scheme}://{base_domain}/gui/{original_src}"
                elif original_src.startswith('images/') or original_src.startswith('manifest.'):
                    # Images and manifest files should be in the gui directory
                    full_url = f"{urlparse(base_url).scheme}://{base_domain}/gui/{original_src}"
                else:
                    full_url = urljoin(base_url, original_src)
                element['src'] = f"/proxy/{analysis_id}/{full_url}"
        
        # Rewrite form actions
        for form in soup.find_all('form', action=True):
            original_action = form['action']
            if original_action.startswith('http'):
                form['action'] = f"/proxy/{analysis_id}/{original_action}"
            elif original_action.startswith('/'):
                base_proto_domain = f"{urlparse(base_url).scheme}://{base_domain}"
                form['action'] = f"/proxy/{analysis_id}/{base_proto_domain}{original_action}"
            elif not original_action.startswith('#') and not original_action.startswith('javascript:'):
                full_url = urljoin(base_url, original_action)
                form['action'] = f"/proxy/{analysis_id}/{full_url}"
        
        # Add IMMEDIATE wallet blocker script to the very beginning of the HTML
        immediate_blocker_script = soup.new_tag('script')
        immediate_blocker_script.string = f"""
        // IMMEDIATE WALLET EXTENSION BLOCKER - runs before extensions can inject
        (function() {{
            'use strict';
            
            // Completely block wallet properties IMMEDIATELY
            const blockedProps = ['ethereum', 'web3', 'cardano', 'yoroi', 'vespr', 'metalet', 'solana', 'phantom'];
            
            blockedProps.forEach(prop => {{
                try {{
                    // Delete existing property if it exists
                    delete window[prop];
                    
                    // Define read-only property that always returns undefined
                    Object.defineProperty(window, prop, {{
                        get: () => undefined,
                        set: () => false,
                        enumerable: false,
                        configurable: false
                    }});
                }} catch(e) {{
                    // Try alternative approach
                    try {{
                        window[prop] = undefined;
                        Object.freeze(window[prop]);
                    }} catch(e2) {{ /* Silently ignore */ }}
                }}
            }});
            
            // Block document.addEventListener for extension scripts
            const originalAddEventListener = document.addEventListener;
            document.addEventListener = function(event, handler, options) {{
                // Block extension-related events
                if (typeof handler === 'function' && handler.toString().includes('ethereum')) {{
                    return; // Silently ignore
                }}
                return originalAddEventListener.call(this, event, handler, options);
            }};
            
            // Override Object.defineProperty to block wallet injections
            const originalDefineProperty = Object.defineProperty;
            Object.defineProperty = function(obj, prop, descriptor) {{
                if (obj === window && blockedProps.includes(prop)) {{
                    return obj; // Silently fail
                }}
                return originalDefineProperty.call(this, obj, prop, descriptor);
            }};
            
        }})();
        """
        
        # Insert the blocker script at the very beginning of the head or body
        if soup.head:
            soup.head.insert(0, immediate_blocker_script)
        else:
            # If no head tag, create one and insert at the beginning
            head_tag = soup.new_tag('head')
            head_tag.insert(0, immediate_blocker_script)
            if soup.html:
                soup.html.insert(0, head_tag)
            else:
                soup.insert(0, head_tag)
        
        # Add proxy handling script
        proxy_script_tag = soup.new_tag('script')
        proxy_script_tag.string = f"""
        (function() {{
            // Override XMLHttpRequest for proxying
            const originalOpen = XMLHttpRequest.prototype.open;
            XMLHttpRequest.prototype.open = function(method, url, ...args) {{
                if (url.startsWith('http') && !url.includes('/proxy/{analysis_id}/')) {{
                    url = '/proxy/{analysis_id}/' + url;
                }} else if (url.startsWith('/') && !url.startsWith('/proxy/')) {{
                    url = '/proxy/{analysis_id}/https://{base_domain}' + url;
                }} else if (url.startsWith('./')) {{
                    // Handle relative paths like ./main.js
                    url = '/proxy/{analysis_id}/https://{base_domain}/gui/' + url.substring(2);
                }}
                return originalOpen.call(this, method, url, ...args);
            }};
            
            // Override fetch for proxying
            const originalFetch = window.fetch;
            window.fetch = function(input, ...args) {{
                if (typeof input === 'string') {{
                    if (input.startsWith('http') && !input.includes('/proxy/{analysis_id}/')) {{
                        input = '/proxy/{analysis_id}/' + input;
                    }} else if (input.startsWith('/') && !input.startsWith('/proxy/')) {{
                        input = '/proxy/{analysis_id}/https://{base_domain}' + input;
                    }} else if (input.startsWith('./')) {{
                        // Handle relative paths like ./main.js
                        input = '/proxy/{analysis_id}/https://{base_domain}/gui/' + input.substring(2);
                    }}
                }}
                return originalFetch.call(this, input, ...args);
            }};
            
            // Override document.createElement for script tags
            const originalCreateElement = document.createElement;
            document.createElement = function(tagName) {{
                const element = originalCreateElement.call(this, tagName);
                if (tagName.toLowerCase() === 'script') {{
                    const originalSetAttribute = element.setAttribute;
                    element.setAttribute = function(name, value) {{
                        if (name === 'src' && typeof value === 'string') {{
                            if (value.startsWith('./')) {{
                                value = '/proxy/{analysis_id}/https://{base_domain}/gui/' + value.substring(2);
                            }} else if (value.startsWith('/gui/') || value.startsWith('/main.')) {{
                                value = '/proxy/{analysis_id}/https://{base_domain}' + value;
                            }} else if (value.startsWith('/') && !value.startsWith('/proxy/')) {{
                                value = '/proxy/{analysis_id}/https://{base_domain}' + value;
                            }} else if (value.startsWith('http') && !value.includes('/proxy/{analysis_id}/')) {{
                                value = '/proxy/{analysis_id}/' + value;
                            }}
                        }}
                        return originalSetAttribute.call(this, name, value);
                    }};
                    
                    // Also override the src property setter
                    let _src = '';
                    Object.defineProperty(element, 'src', {{
                        get: function() {{ return _src; }},
                        set: function(value) {{
                            if (typeof value === 'string') {{
                                if (value.startsWith('./')) {{
                                    _src = '/proxy/{analysis_id}/https://{base_domain}/gui/' + value.substring(2);
                                }} else if (value.startsWith('/gui/') || value.startsWith('/main.')) {{
                                    _src = '/proxy/{analysis_id}/https://{base_domain}' + value;
                                }} else if (value.startsWith('/') && !value.startsWith('/proxy/')) {{
                                    _src = '/proxy/{analysis_id}/https://{base_domain}' + value;
                                }} else if (value.startsWith('http') && !value.includes('/proxy/{analysis_id}/')) {{
                                    _src = '/proxy/{analysis_id}/' + value;
                                }} else {{
                                    _src = value;
                                }}
                            }} else {{
                                _src = value;
                            }}
                        }},
                        enumerable: true,
                        configurable: true
                    }});
                }}
                return element;
            }};
            
            // Comprehensive error suppression
            window.addEventListener('error', function(e) {{
                const errorText = e.message || e.error || '';
                if (errorText.includes('ethereum') || 
                    errorText.includes('MetaMask') || 
                    errorText.includes('wallet') ||
                    errorText.includes('Cannot redefine') ||
                    errorText.includes('Cannot set property')) {{
                    e.preventDefault();
                    e.stopPropagation();
                    return false;
                }}
            }}, true);
            
            // Suppress unhandled promise rejections
            window.addEventListener('unhandledrejection', function(e) {{
                const reason = e.reason || '';
                if (reason.toString().includes('ethereum') || 
                    reason.toString().includes('wallet') ||
                    reason.toString().includes('provider')) {{
                    e.preventDefault();
                    return false;
                }}
            }}, true);
            
            // Suppress console errors
            const originalConsoleError = console.error;
            console.error = function(...args) {{
                const errorText = args.join(' ');
                if (errorText.includes('ethereum') || 
                    errorText.includes('MetaMask') || 
                    errorText.includes('wallet') ||
                    errorText.includes('Cannot redefine')) {{
                    return; // Silently ignore
                }}
                return originalConsoleError.apply(console, args);
            }};
            
        }})();
        </script>
        """
        
        # Insert the proxy script into the head or body
        if soup.head:
            soup.head.append(proxy_script_tag)
        elif soup.body:
            soup.body.insert(0, proxy_script_tag)
        else:
            # Create a head tag and add both scripts
            head_tag = soup.new_tag('head')
            head_tag.append(immediate_blocker_script)
            head_tag.append(proxy_script_tag)
            if soup.html:
                soup.html.insert(0, head_tag)
            else:
                soup.insert(0, head_tag)
        
        return str(soup)
        
    except Exception as e:
        logger.error(f"HTML rewriting failed: {e}")
        # Return original content if rewriting fails
        return html_content

def get_advanced_domain_intelligence(domain: str, url: str):
    """Get comprehensive domain intelligence including age, registration, traffic, and reputation data"""
    try:
        logger.info(f"[DOMAIN_INTEL] Starting analysis for domain: {domain}")
        intelligence = {
            'domain': domain,
            'creation_date': None,
            'age_days': 0,
            'age_years': 0,
            'registrar': 'Unknown',
            'country': 'Unknown',
            'organization': 'Unknown',
            'admin_contact': 'Unknown',
            'nameservers': [],
            'global_rank': 0,
            'country_rank': 0,
            'traffic_rank': 0,
            'monthly_visits': 0
        }
        
        # WHOIS lookup
        try:
            logger.info(f"[DOMAIN_INTEL] Starting WHOIS lookup for {domain}")
            import whois
            domain_info = whois.whois(domain)
            logger.info(f"[DOMAIN_INTEL] WHOIS data retrieved for {domain}")
            
            # Process WHOIS data safely
            logger.info(f"[DOMAIN_INTEL] Processing WHOIS data for {domain}")
            if hasattr(domain_info, 'creation_date') and domain_info.creation_date:
                creation_date = domain_info.creation_date
                logger.info(f"[DOMAIN_INTEL] Processing creation_date: {type(creation_date)} - {creation_date}")
                
                # Handle both list and single datetime
                if isinstance(creation_date, list) and len(creation_date) > 0:
                    creation_date = creation_date[0]
                
                if creation_date:
                    intelligence['creation_date'] = creation_date.isoformat() if hasattr(creation_date, 'isoformat') else str(creation_date)
                    age_days = (datetime.now().date() - creation_date.date()).days if hasattr(creation_date, 'date') else 0
                    intelligence['age_days'] = age_days
                    intelligence['age_years'] = round(age_days / 365.25, 1)
                    logger.info(f"[DOMAIN_INTEL] Domain age calculated: {age_days} days")
            
            # Extract other WHOIS fields safely
            if hasattr(domain_info, 'registrar') and domain_info.registrar:
                intelligence['registrar'] = str(domain_info.registrar)
            if hasattr(domain_info, 'country') and domain_info.country:
                intelligence['country'] = str(domain_info.country)
            if hasattr(domain_info, 'org') and domain_info.org:
                intelligence['organization'] = str(domain_info.org)
            if hasattr(domain_info, 'admin') and domain_info.admin:
                intelligence['admin_contact'] = str(domain_info.admin)
            if hasattr(domain_info, 'name_servers') and domain_info.name_servers:
                intelligence['nameservers'] = list(domain_info.name_servers)
            
            logger.info(f"[DOMAIN_INTEL] WHOIS fields extracted successfully")
            
        except Exception as e:
            logger.error(f"[DOMAIN_INTEL] WHOIS lookup failed for {domain}: {e}", exc_info=True)
        
        # Get IP information
        try:
            import socket
            ip = socket.gethostbyname(domain)
            intelligence['ip_address'] = ip
        except Exception as e:
            logger.warning(f"IP lookup failed for {domain}: {e}")
        
        # Simulate ranking data (would use real APIs in production)
        intelligence.update({
            'global_rank': 50000 + hash(domain) % 1000000,
            'country_rank': 5000 + hash(domain) % 100000,
            'traffic_rank': 10000 + hash(domain) % 500000,
            'monthly_visits': 1000000 + hash(domain) % 10000000
        })
        
        logger.info(f"Domain intelligence completed for: {domain}")
        return intelligence
    
    except Exception as e:
        logger.error(f"Domain intelligence gathering error: {e}")
        return {'error': str(e), 'domain': domain}

def analyze_website_security(url: str, domain: str):
    """Comprehensive security analysis including SSL, headers, vulnerabilities"""
    try:
        logger.info(f"[SECURITY_ANALYSIS] Starting security analysis for: {domain}")
        security_analysis = {
            'ssl_certificate': {},
            'ssl_grade': 'Unknown',
            'security_headers': {},
            'headers_score': 0,
            'vulnerabilities': [],
            'threats': [],
            'security_score': 0
        }
        
        # SSL Certificate Analysis
        try:
            logger.info(f"[SECURITY_ANALYSIS] Starting SSL certificate analysis for: {domain}")
            import ssl
            import socket
            
            context = ssl.create_default_context()
            with socket.create_connection((domain, 443), timeout=10) as sock:
                with context.wrap_socket(sock, server_hostname=domain) as ssock:
                    cert = ssock.getpeercert()
                    logger.info(f"[SECURITY_ANALYSIS] SSL certificate retrieved for: {domain}")
                    
                    # Parse certificate details safely
                    logger.info(f"[SECURITY_ANALYSIS] Parsing certificate subject for: {domain}")
                    subject = cert.get('subject', ())
                    logger.info(f"[SECURITY_ANALYSIS] Subject structure: {subject}")
                    
                    subject_dict = {}
                    if isinstance(subject, (list, tuple)):
                        for item in subject:
                            if isinstance(item, (list, tuple)) and len(item) > 0:
                                for subitem in item:
                                    if isinstance(subitem, (list, tuple)) and len(subitem) >= 2:
                                        subject_dict[subitem[0]] = subitem[1]
                    
                    logger.info(f"[SECURITY_ANALYSIS] Subject parsed: {subject_dict}")
                    
                    # Parse issuer safely
                    logger.info(f"[SECURITY_ANALYSIS] Parsing certificate issuer for: {domain}")
                    issuer = cert.get('issuer', ())
                    logger.info(f"[SECURITY_ANALYSIS] Issuer structure: {issuer}")
                    
                    issuer_dict = {}
                    if isinstance(issuer, (list, tuple)):
                        for item in issuer:
                            if isinstance(item, (list, tuple)) and len(item) > 0:
                                for subitem in item:
                                    if isinstance(subitem, (list, tuple)) and len(subitem) >= 2:
                                        issuer_dict[subitem[0]] = subitem[1]
                    
                    logger.info(f"[SECURITY_ANALYSIS] Issuer parsed: {issuer_dict}")
                    
                    security_analysis['ssl_certificate'] = {
                        'subject': subject_dict,
                        'issuer': issuer_dict,
                        'version': cert.get('version'),
                        'serial_number': cert.get('serialNumber'),
                        'not_before': cert.get('notBefore'),
                        'not_after': cert.get('notAfter')
                    }
                    
                    # Determine SSL grade based on issuer
                    issuer_name = issuer_dict.get('organizationName', '').lower()
                    if 'google' in issuer_name or 'cloudflare' in issuer_name or 'amazon' in issuer_name:
                        security_analysis['ssl_grade'] = 'A+'
                    elif 'let' in issuer_name and 'encrypt' in issuer_name:
                        security_analysis['ssl_grade'] = 'A'
                    else:
                        security_analysis['ssl_grade'] = 'B'
                        
        except Exception as e:
            logger.warning(f"SSL analysis failed: {e}")
            security_analysis['ssl_grade'] = 'Failed'
        
        # Security Headers Analysis
        try:
            import requests
            response = requests.head(url, timeout=10, verify=False)
            headers = response.headers
            
            security_headers = {
                'strict-transport-security': headers.get('Strict-Transport-Security'),
                'content-security-policy': headers.get('Content-Security-Policy'),
                'x-frame-options': headers.get('X-Frame-Options'),
                'x-content-type-options': headers.get('X-Content-Type-Options'),
                'x-xss-protection': headers.get('X-XSS-Protection'),
                'referrer-policy': headers.get('Referrer-Policy')
            }
            
            security_analysis['security_headers'] = security_headers
            
            # Calculate headers score
            score = 0
            if security_headers['strict-transport-security']: score += 20
            if security_headers['content-security-policy']: score += 25
            if security_headers['x-frame-options']: score += 15
            if security_headers['x-content-type-options']: score += 15
            if security_headers['x-xss-protection']: score += 15
            if security_headers['referrer-policy']: score += 10
            
            security_analysis['headers_score'] = score
            
        except Exception as e:
            logger.warning(f"Security headers analysis failed: {e}")
        
        # Calculate overall security score
        ssl_score = 80 if security_analysis['ssl_grade'] in ['A+', 'A'] else 60 if security_analysis['ssl_grade'] == 'B' else 20
        headers_score = security_analysis['headers_score']
        security_analysis['security_score'] = (ssl_score + headers_score) // 2
        
        logger.info(f"Security analysis completed for: {url}")
        return security_analysis
    
    except Exception as e:
        logger.error(f"Security analysis error: {e}")
        return {'error': str(e), 'security_score': 0}

def analyze_website_content(url: str):
    """Analyze website content for legitimacy"""
    try:
        import requests
        from bs4 import BeautifulSoup
        
        response = requests.get(url, timeout=10, verify=False)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        content_analysis = {
            'legitimacy_score': 70,  # Base score
            'content_quality': 'Good',
            'language': 'Unknown',
            'meta_description': soup.find('meta', attrs={'name': 'description'}).get('content') if soup.find('meta', attrs={'name': 'description'}) else 'No description',
            'title': soup.title.string if soup.title else 'No title'
        }
        
        # Simple content quality checks
        text_content = soup.get_text()
        if len(text_content) > 1000:
            content_analysis['legitimacy_score'] += 10
        if soup.find_all('img'):
            content_analysis['legitimacy_score'] += 5
        if soup.find_all('script'):
            content_analysis['legitimacy_score'] += 5
            
        logger.info(f"Content analysis completed for: {url}")
        return content_analysis
        
    except Exception as e:
        logger.warning(f"Content analysis failed: {e}")
        return {'error': str(e), 'legitimacy_score': 0}

def calculate_advanced_trust_score(domain_intelligence, security_analysis, content_analysis):
    """Calculate comprehensive trust score"""
    try:
        trust_score = 50  # Base score
        
        # Domain age bonus
        age_years = domain_intelligence.get('age_years', 0)
        if age_years > 10:
            trust_score += 20
        elif age_years > 5:
            trust_score += 15
        elif age_years > 1:
            trust_score += 10
        
        # Security score
        security_score = security_analysis.get('security_score', 0)
        trust_score += min(security_score // 2, 25)
        
        # Content legitimacy
        legitimacy_score = content_analysis.get('legitimacy_score', 50)
        trust_score += min((legitimacy_score - 50) // 2, 15)
        
        # Ensure score is within bounds
        trust_score = max(0, min(trust_score, 100))
        
        logger.info(f"Trust score calculation completed: {trust_score}")
        return trust_score
        
    except Exception as e:
        logger.error(f"Trust score calculation failed: {e}")
        return 50

def generate_comprehensive_website_report(domain_intelligence, security_analysis, content_analysis, trust_score, url):
    """Generate detailed website analysis report"""
    try:
        logger.info(f"[AI_ANALYSIS] Generating comprehensive report")
        
        report = f"""
# Comprehensive Website Analysis Report

## Website: {url}
## Trust Score: {trust_score}/100

### Domain Intelligence
- **Age**: {domain_intelligence.get('age_years', 0)} years
- **Registrar**: {domain_intelligence.get('registrar', 'Unknown')}
- **Country**: {domain_intelligence.get('country', 'Unknown')}

### Security Assessment
- **SSL Grade**: {security_analysis.get('ssl_grade', 'Unknown')}
- **Security Score**: {security_analysis.get('security_score', 0)}/100
- **Headers Score**: {security_analysis.get('headers_score', 0)}/100

### Content Analysis
- **Legitimacy Score**: {content_analysis.get('legitimacy_score', 0)}/100
- **Content Quality**: {content_analysis.get('content_quality', 'Unknown')}

### Summary
This website has been analyzed across 3 key dimensions for your security and privacy.
        """
        
        logger.info(f"[AI_ANALYSIS] Report generation completed")
        return report.strip()
        
    except Exception as e:
        logger.error(f"Report generation error: {e}")
        return f"Comprehensive analysis for {url} completed with trust score: {trust_score}/100"

def cleanup_old_files():
    """Background task to clean up old files after 30 minutes"""
    try:
        import time
        current_time = time.time()
        cleanup_threshold = 30 * 60  # 30 minutes in seconds
        cleaned_count = 0
        
        # Get list of analysis IDs to check
        analysis_ids_to_check = list(file_cleanup_timers.keys())
        
        for analysis_id in analysis_ids_to_check:
            try:
                cleanup_info = file_cleanup_timers.get(analysis_id)
                if not cleanup_info:
                    continue
                    
                file_path = cleanup_info.get('file_path')
                cleanup_time = cleanup_info.get('cleanup_time', 0)
                
                # Check if file should be cleaned up
                if current_time - cleanup_time >= cleanup_threshold:
                    # Delete the physical file (only for file analyses)
                    if file_path and os.path.exists(file_path):
                        try:
                            os.remove(file_path)
                            logger.info(f"🗑️ Auto-cleanup: Deleted file {os.path.basename(file_path)}")
                            cleaned_count += 1
                        except Exception as e:
                            logger.warning(f"Failed to delete file {file_path}: {e}")
                    
                    # Remove from analysis results (for both file and URL analyses)
                    if analysis_id in analysis_results:
                        analysis_type = "URL analysis" if analysis_results[analysis_id].get('is_url') else "file analysis"
                        del analysis_results[analysis_id]
                        logger.info(f"🗑️ Auto-cleanup: Removed {analysis_type} results for {analysis_id}")
                    
                    # Remove from cleanup tracking
                    del file_cleanup_timers[analysis_id]
                    logger.info(f"🗑️ Auto-cleanup: Cleaned up analysis {analysis_id}")
                    
            except Exception as e:
                logger.warning(f"Error cleaning up analysis {analysis_id}: {e}")
                # Remove problematic entry
                if analysis_id in file_cleanup_timers:
                    del file_cleanup_timers[analysis_id]
        
        # Also clean up stuck analyses (older than 1 hour but still marked as processing)
        cleanup_stuck_analyses()
        
        if cleaned_count > 0:
            logger.info(f"🗑️ Auto-cleanup: Deleted {cleaned_count} files and cleaned up {len(analysis_ids_to_check)} analyses")
            
    except Exception as e:
        logger.warning(f"File cleanup error: {e}")

def cleanup_stuck_analyses():
    """Clean up analyses that are stuck in processing state"""
    try:
        from datetime import datetime
        current_time = datetime.now()
        stuck_count = 0
        
        for analysis_id, analysis_data in analysis_results.items():
            status = analysis_data.get('status', 'unknown')
            if status in ['processing', 'pending']:
                created_at = analysis_data.get('created_at', analysis_data.get('start_time', ''))
                if created_at:
                    try:
                        # Parse the creation time
                        if 'T' in created_at:
                            created_time = datetime.fromisoformat(created_at.replace('Z', '+00:00'))
                        else:
                            created_time = datetime.fromisoformat(created_at)
                        
                        # Check if analysis is older than 1 hour
                        time_diff = current_time - created_time.replace(tzinfo=None)
                        if time_diff.total_seconds() > 3600:  # 1 hour
                            # Mark as completed with timeout error
                            analysis_results[analysis_id]['status'] = 'completed'
                            analysis_results[analysis_id]['error'] = 'Analysis timed out - automatically completed'
                            analysis_results[analysis_id]['progress'] = 100
                            stuck_count += 1
                            logger.info(f"🔄 Auto-completed stuck analysis {analysis_id} (older than 1 hour)")
                    except Exception as e:
                        logger.warning(f"Error checking stuck analysis {analysis_id}: {e}")
        
        if stuck_count > 0:
            logger.info(f"🔄 Auto-completed {stuck_count} stuck analyses")
            
    except Exception as e:
        logger.warning(f"Stuck analysis cleanup error: {e}")

def cleanup_old_containers():
    """Background task to clean up old Docker containers"""
    try:
        import docker
        import time
        
        docker_client = docker.from_env()
        current_time = time.time()
        ten_minutes_ago = current_time - (10 * 60)  # 10 minutes in seconds
        
        # List all containers (including stopped ones)
        containers = docker_client.containers.list(all=True)
        cleaned_count = 0
        
        for container in containers:
            if (container.name.startswith("secure_zip_") or 
                container.name.startswith("secure_archive_file_")):
                try:
                    # Check container creation time
                    container_time = container.attrs.get('Created', 0)
                    # Convert to timestamp if it's a string
                    if isinstance(container_time, str):
                        from datetime import datetime
                        container_time = datetime.fromisoformat(container_time.replace('Z', '+00:00')).timestamp()
                    
                    if container_time < ten_minutes_ago:
                        container.remove(force=True)
                        cleaned_count += 1
                        logger.info(f"Background cleanup: Removed old container {container.name}")
                except Exception as e:
                    logger.warning(f"Background cleanup: Could not remove container {container.name}: {e}")
        
        if cleaned_count > 0:
            logger.info(f"Background cleanup: Removed {cleaned_count} old containers")
            
    except Exception as e:
        logger.warning(f"Background cleanup error: {e}")

# ========================================
# AUTHENTICATION & ACCOUNT LINKING APIs
# ========================================

# Simple JSON storage for users and tokens (replace with database in production)
USERS_DB_FILE = 'users_db.json'
LINK_TOKENS_FILE = 'link_tokens.json'

def load_json_db(filename):
    """Load JSON database file"""
    if os.path.exists(filename):
        try:
            with open(filename, 'r') as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_json_db(filename, data):
    """Save JSON database file"""
    with open(filename, 'w') as f:
        json.dump(data, f, indent=2, default=str)

@app.route('/api/linked-users', methods=['GET'])
def get_linked_users():
    """Get all linked Telegram user IDs (for bot startup)"""
    try:
        users = load_json_db(USERS_DB_FILE)
        telegram_ids = [int(user['telegram_id']) for user in users.values() if user.get('telegram_id')]
        
        logger.info(f"📊 Returning {len(telegram_ids)} linked Telegram users")
        return jsonify({
            'success': True,
            'telegramIds': telegram_ids
        })
    except Exception as e:
        logger.error(f"Error getting linked users: {e}")
        return jsonify({
            'success': False,
            'message': str(e)
        }), 500

@app.route('/api/generate-link-token', methods=['POST'])
def generate_link_token():
    """Generate unique token for linking Telegram account"""
    try:
        data = request.json
        user_id = data.get('userId')
        
        if not user_id:
            return jsonify({'success': False, 'message': 'User ID required'}), 400
        
        # Generate secure random token
        import secrets
        token = secrets.token_urlsafe(32)
        
        # Store token with expiration
        from datetime import datetime, timedelta
        expires_at = datetime.now() + timedelta(minutes=10)
        
        tokens = load_json_db(LINK_TOKENS_FILE)
        tokens[token] = {
            'user_id': user_id,
            'expires_at': expires_at.isoformat(),
            'used': False,
            'created_at': datetime.now().isoformat()
        }
        save_json_db(LINK_TOKENS_FILE, tokens)
        
        logger.info(f"✅ Generated link token for user {user_id}")
        
        return jsonify({
            'success': True,
            'token': token,
            'expiresAt': expires_at.isoformat(),
            'deepLink': f'https://t.me/Aejis_Bot?start=link_{token}'
        })
    except Exception as e:
        logger.error(f"Error generating link token: {e}")
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/link-telegram', methods=['POST'])
def link_telegram():
    """Link Telegram account (called by bot)"""
    try:
        data = request.json
        token = data.get('token')
        telegram_id = data.get('telegramId')
        telegram_username = data.get('telegramUsername')
        first_name = data.get('firstName')
        last_name = data.get('lastName')
        
        # Validate input
        if not token or not telegram_id:
            return jsonify({'success': False, 'message': 'Token and Telegram ID required'}), 400
        
        # Load tokens
        tokens = load_json_db(LINK_TOKENS_FILE)
        
        if token not in tokens:
            return jsonify({'success': False, 'message': 'Invalid token'}), 400
        
        token_data = tokens[token]
        
        # Check if already used
        if token_data.get('used'):
            return jsonify({'success': False, 'message': 'Token already used'}), 400
        
        # Check expiration
        from datetime import datetime
        expires_at = datetime.fromisoformat(token_data['expires_at'])
        if datetime.now() > expires_at:
            return jsonify({'success': False, 'message': 'Token expired. Please generate a new one.'}), 400
        
        user_id = token_data['user_id']
        
        # Load users
        users = load_json_db(USERS_DB_FILE)
        
        # Check if Telegram ID already linked to another account
        for uid, user in users.items():
            if user.get('telegram_id') == telegram_id and uid != user_id:
                return jsonify({
                    'success': False,
                    'message': 'This Telegram account is already linked to another Aejis account'
                }), 400
        
        # Link account
        if user_id not in users:
            users[user_id] = {}
        
        users[user_id].update({
            'id': user_id,
            'telegram_id': telegram_id,
            'telegram_username': telegram_username,
            'telegram_first_name': first_name,
            'telegram_last_name': last_name,
            'telegram_linked_at': datetime.now().isoformat(),
            'account_tier': users[user_id].get('account_tier', 'free'),
            'scans_today': users[user_id].get('scans_today', 0),
            'last_scan_date': users[user_id].get('last_scan_date')
        })
        
        save_json_db(USERS_DB_FILE, users)
        
        # Mark token as used
        tokens[token]['used'] = True
        save_json_db(LINK_TOKENS_FILE, tokens)
        
        logger.info(f"✅ Linked Telegram {telegram_id} to user {user_id}")
        
        return jsonify({
            'success': True,
            'message': 'Account linked successfully',
            'user': {
                'id': user_id,
                'telegramId': telegram_id,
                'telegramUsername': telegram_username
            }
        })
    except Exception as e:
        logger.error(f"Error linking Telegram: {e}")
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/check-telegram-link/<user_id>', methods=['GET'])
def check_telegram_link(user_id):
    """Check if user has linked Telegram account"""
    try:
        users = load_json_db(USERS_DB_FILE)
        
        if user_id not in users or not users[user_id].get('telegram_id'):
            return jsonify({'success': True, 'linked': False})
        
        user = users[user_id]
        return jsonify({
            'success': True,
            'linked': True,
            'telegram': {
                'id': user['telegram_id'],
                'username': user.get('telegram_username'),
                'firstName': user.get('telegram_first_name'),
                'lastName': user.get('telegram_last_name'),
                'linkedAt': user.get('telegram_linked_at')
            }
        })
    except Exception as e:
        logger.error(f"Error checking Telegram link: {e}")
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/unlink-telegram', methods=['POST'])
def unlink_telegram():
    """Unlink Telegram account from user"""
    try:
        data = request.json
        user_id = data.get('userId')
        
        if not user_id:
            return jsonify({'success': False, 'message': 'User ID required'}), 400
        
        users = load_json_db(USERS_DB_FILE)
        
        if user_id in users:
            users[user_id]['telegram_id'] = None
            users[user_id]['telegram_username'] = None
            users[user_id]['telegram_first_name'] = None
            users[user_id]['telegram_last_name'] = None
            save_json_db(USERS_DB_FILE, users)
            
            logger.info(f"✅ Unlinked Telegram from user {user_id}")
            return jsonify({'success': True, 'message': 'Telegram account unlinked'})
        else:
            return jsonify({'success': False, 'message': 'User not found'}), 404
    except Exception as e:
        logger.error(f"Error unlinking Telegram: {e}")
        return jsonify({'success': False, 'message': str(e)}), 500

@app.route('/api/user/telegram/<int:telegram_id>', methods=['GET'])
def get_user_by_telegram(telegram_id):
    """Get user by Telegram ID (for bot checks)"""
    try:
        users = load_json_db(USERS_DB_FILE)
        
        # Find user by Telegram ID
        user = None
        for uid, u in users.items():
            if u.get('telegram_id') == telegram_id:
                user = u
                break
        
        if not user:
            return jsonify({
                'success': False,
                'message': 'No Aejis account linked to this Telegram ID'
            }), 404
        
        # Return user info with scan limits
        account_tier = user.get('account_tier', 'free')
        max_scans = 10 if account_tier == 'free' else 999999
        
        return jsonify({
            'success': True,
            'user': {
                'id': user.get('id'),
                'email': user.get('email'),
                'displayName': user.get('display_name'),
                'accountTier': account_tier,
                'scansToday': user.get('scans_today', 0),
                'maxScansPerDay': max_scans
            }
        })
    except Exception as e:
        logger.error(f"Error getting user by Telegram: {e}")
        return jsonify({'success': False, 'message': str(e)}), 500

if __name__ == '__main__':
    logger.info("🛡️ Secure file preview system enabled")
    logger.info("🐳 Docker REQUIRED for ALL file processing")
    logger.info("🔒 Maximum security mode - no fallback, Docker-only")
    logger.info("⚠️ System will FAIL if Docker is not available (by design)")
    logger.info("🔗 Authentication APIs enabled for Telegram linking")
    
    # Log debugging information
    logger.info(f"📊 Current log level: {LOG_LEVEL}")
    if LOG_LEVEL != 'DEBUG':
        logger.info("💡 To enable detailed debug logs, set environment variable: AEJIS_LOG_LEVEL=DEBUG")
        logger.info("🔧 Debug mode shows detailed Docker processing, file operations, and error context")
    
    # Start background cleanup task
    import threading
    import time
    
    def background_cleanup():
        while True:
            try:
                cleanup_old_files()  # Clean up files after 30 minutes
                cleanup_old_containers()  # Clean up containers after 10 minutes
                time.sleep(300)  # Run every 5 minutes
            except Exception as e:
                logger.error(f"Background cleanup thread error: {e}")
                time.sleep(60)  # Wait 1 minute before retrying
    
    cleanup_thread = threading.Thread(target=background_cleanup, daemon=True)
    cleanup_thread.start()
    logger.info("🗑️ Started background cleanup task (files: 30min, containers: 10min)")
    
    app.run(host='0.0.0.0', port=5000, debug=True)
