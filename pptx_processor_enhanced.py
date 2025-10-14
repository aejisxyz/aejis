#!/usr/bin/env python3
"""
Enhanced PPTX Processor for Aejis
Processes PowerPoint presentations and converts them to high-quality images
"""

import sys
import json
import base64
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import xml.etree.ElementTree as ET
from PIL import Image, ImageDraw, ImageFont
import io
import subprocess
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import Rectangle, FancyBboxPatch
import numpy as np
import re

def extract_slide_content(slide):
    """Extract detailed content from a slide including formatting"""
    content = {
        'title': '',
        'text_boxes': [],
        'shapes': [],
        'images': [],
        'tables': []
    }
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            if shape.has_text_frame:
                # Check if it's a title placeholder
                if shape.placeholder_format.idx == 0:  # Title placeholder
                    content['title'] = shape.text.strip()
                else:
                    # Regular text box
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_content.append({
                                'text': paragraph.text.strip(),
                                'level': paragraph.level,
                                'alignment': paragraph.alignment,
                                'font_size': paragraph.font.size.pt if paragraph.font.size else 12
                            })
                    
                    if text_content:
                        content['text_boxes'].append({
                            'content': text_content,
                            'position': (shape.left, shape.top),
                            'size': (shape.width, shape.height)
                        })
        
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Extract table data
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            content['tables'].append({
                'data': table_data,
                'position': (shape.left, shape.top),
                'size': (shape.width, shape.height)
            })
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Extract image information
            content['images'].append({
                'position': (shape.left, shape.top),
                'size': (shape.width, shape.height)
            })
    
    return content

def create_enhanced_slide_image(slide_content, slide_number, slide_width, slide_height):
    """Create a high-quality slide image with proper formatting"""
    # Convert PowerPoint units to pixels (assuming 96 DPI)
    width_px = int(slide_width / 914400 * 96 * 10)  # Scale up for better quality
    height_px = int(slide_height / 914400 * 96 * 10)
    
    # Create figure with proper aspect ratio
    fig, ax = plt.subplots(figsize=(width_px/100, height_px/100))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7.5)  # 16:9 aspect ratio
    ax.set_aspect('equal')
    ax.axis('off')
    
    # Set background
    ax.add_patch(Rectangle((0, 0), 10, 7.5, facecolor='white', edgecolor='#e0e0e0', linewidth=2))
    
    # Add slide title
    if slide_content['title']:
        ax.text(5, 6.8, slide_content['title'], ha='center', va='center', 
                fontsize=24, fontweight='bold', color='#2c3e50', 
                bbox=dict(boxstyle="round,pad=0.3", facecolor='#ecf0f1', alpha=0.8))
    
    # Add slide number
    ax.text(9.5, 0.2, f'{slide_number}', ha='right', va='center', 
            fontsize=12, color='#7f8c8d', alpha=0.7)
    
    # Add text boxes
    y_position = 6.0
    for text_box in slide_content['text_boxes']:
        for text_item in text_box['content']:
            # Determine font size
            font_size = max(10, min(16, text_item['font_size'] or 14))
            
            # Determine alignment
            ha = 'left'
            if text_item['alignment'] == PP_ALIGN.CENTER:
                ha = 'center'
            elif text_item['alignment'] == PP_ALIGN.RIGHT:
                ha = 'right'
            
            # Add indentation for bullet points
            x_pos = 0.5 + (text_item['level'] * 0.5)
            
            # Create text with proper formatting
            ax.text(x_pos, y_position, text_item['text'], ha=ha, va='center',
                   fontsize=font_size, color='#34495e', wrap=True)
            
            y_position -= 0.4
    
    # Add tables
    for table in slide_content['tables']:
        if table['data']:
            # Create a simple table representation
            table_text = '\n'.join([' | '.join(row) for row in table['data'][:3]])  # Limit to 3 rows
            ax.text(5, 3, table_text, ha='center', va='center',
                   fontsize=12, color='#2c3e50',
                   bbox=dict(boxstyle="round,pad=0.5", facecolor='#f8f9fa', edgecolor='#dee2e6'))
    
    # Add images placeholder
    for img in slide_content['images']:
        ax.add_patch(Rectangle((2, 2), 6, 3, facecolor='#e9ecef', edgecolor='#adb5bd', 
                              linewidth=2, linestyle='--'))
        ax.text(5, 3.5, '📷 Image', ha='center', va='center', fontsize=16, color='#6c757d')
    
    # Convert to base64
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=300, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    img_data = buf.getvalue()
    img_base64 = base64.b64encode(img_data).decode('utf-8')
    plt.close(fig)
    
    return img_base64

def convert_pptx_to_enhanced_images(file_path):
    """Convert PPTX slides to high-quality PNG images using LibreOffice"""
    try:
        # Use LibreOffice to convert PPTX to images
        output_dir = "/tmp/pptx_images"
        os.makedirs(output_dir, exist_ok=True)
        
        # LibreOffice command to convert PPTX to PNG images
        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to", "png",
            "--outdir", output_dir,
            file_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode != 0:
            print(f"LibreOffice conversion failed: {result.stderr}", file=sys.stderr)
            return []
        
        # Find the generated PNG files
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        slide_images = []
        
        # Look for PNG files in the output directory
        for i in range(1, 20):  # Check up to 20 slides
            png_file = os.path.join(output_dir, f"{base_name}_{i}.png")
            if os.path.exists(png_file):
                # Read and encode the image
                with open(png_file, 'rb') as f:
                    img_data = f.read()
                    img_base64 = base64.b64encode(img_data).decode('utf-8')
                
                slide_images.append({
                    'slide_number': i,
                    'image_base64': img_base64,
                    'mime_type': 'image/png'
                })
            else:
                break  # No more slides
        
        # Clean up temporary files
        try:
            for i in range(1, 20):
                png_file = os.path.join(output_dir, f"{base_name}_{i}.png")
                if os.path.exists(png_file):
                    os.remove(png_file)
        except:
            pass  # Ignore cleanup errors
        
        return slide_images
        
    except Exception as e:
        print(f"Error in convert_pptx_to_enhanced_images: {e}", file=sys.stderr)
        return []

def extract_text_from_pptx(prs):
    """Extract text content from presentation"""
    text_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = f"--- Slide {i + 1} ---\n"
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
        
        text_content.append(slide_text)
    
    return "\n".join(text_content)

def generate_enhanced_html_viewer(prs, slide_images):
    """Generate enhanced HTML viewer with better styling"""
    html_parts = []
    
    # Add CSS styles
    html_parts.append("""
    <style>
    .enhanced-pptx-viewer {
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 16px;
        color: white;
    }
    
    .slide-viewer-header {
        text-align: center;
        margin-bottom: 2rem;
    }
    
    .slide-viewer-header h2 {
        margin: 0 0 1rem 0;
        font-size: 2rem;
        font-weight: 300;
    }
    
    .slide-controls {
        display: flex;
        justify-content: center;
        align-items: center;
        gap: 2rem;
        margin-bottom: 2rem;
    }
    
    .slide-btn {
        background: rgba(255, 255, 255, 0.2);
        border: 2px solid rgba(255, 255, 255, 0.3);
        color: white;
        padding: 0.75rem 1.5rem;
        border-radius: 25px;
        cursor: pointer;
        font-size: 1rem;
        transition: all 0.3s ease;
        backdrop-filter: blur(10px);
    }
    
    .slide-btn:hover {
        background: rgba(255, 255, 255, 0.3);
        transform: translateY(-2px);
    }
    
    .slide-counter {
        font-size: 1.2rem;
        font-weight: 500;
        background: rgba(255, 255, 255, 0.1);
        padding: 0.5rem 1rem;
        border-radius: 20px;
        backdrop-filter: blur(10px);
    }
    
    .slides-container {
        position: relative;
        min-height: 500px;
        margin-bottom: 2rem;
    }
    
    .slide {
        display: none;
        background: white;
        border-radius: 12px;
        padding: 2rem;
        box-shadow: 0 20px 40px rgba(0, 0, 0, 0.3);
        text-align: center;
    }
    
    .slide.active {
        display: block;
        animation: slideIn 0.5s ease-in-out;
    }
    
    @keyframes slideIn {
        from { opacity: 0; transform: translateY(20px); }
        to { opacity: 1; transform: translateY(0); }
    }
    
    .slide-image {
        max-width: 100%;
        max-height: 400px;
        border-radius: 8px;
        box-shadow: 0 10px 20px rgba(0, 0, 0, 0.2);
    }
    
    .slide-dots {
        display: flex;
        justify-content: center;
        gap: 0.5rem;
    }
    
    .dot {
        width: 12px;
        height: 12px;
        border-radius: 50%;
        background: rgba(255, 255, 255, 0.3);
        cursor: pointer;
        transition: all 0.3s ease;
    }
    
    .dot.active {
        background: white;
        transform: scale(1.2);
    }
    
    .dot:hover {
        background: rgba(255, 255, 255, 0.6);
    }
    </style>
    """)
    
    # Add HTML structure
    html_parts.append('<div class="enhanced-pptx-viewer">')
    html_parts.append('<div class="slide-viewer-header">')
    html_parts.append('<h2>📊 PowerPoint Presentation</h2>')
    html_parts.append('<div class="slide-controls">')
    html_parts.append('<button onclick="previousSlide()" class="slide-btn">← Previous</button>')
    html_parts.append('<span class="slide-counter" id="slide-counter">Slide 1 of ' + str(len(slide_images)) + '</span>')
    html_parts.append('<button onclick="nextSlide()" class="slide-btn">Next →</button>')
    html_parts.append('</div>')
    html_parts.append('</div>')
    
    html_parts.append('<div class="slides-container">')
    
    for i, slide_img in enumerate(slide_images):
        slide_class = "slide active" if i == 0 else "slide"
        html_parts.append(f'<div class="{slide_class}" id="slide-{i}">')
        html_parts.append(f'<img src="data:{slide_img["mime_type"]};base64,{slide_img["image_base64"]}" alt="Slide {slide_img["slide_number"]}" class="slide-image">')
        html_parts.append('</div>')
    
    html_parts.append('</div>')
    
    # Add navigation dots
    html_parts.append('<div class="slide-dots">')
    for i in range(len(slide_images)):
        dot_class = "dot active" if i == 0 else "dot"
        html_parts.append(f'<span class="{dot_class}" onclick="goToSlide({i})"></span>')
    html_parts.append('</div>')
    
    # Add JavaScript
    html_parts.append("""
    <script>
    let currentSlide = 0;
    const totalSlides = """ + str(len(slide_images)) + """;
    
    function showSlide(n) {
        const slides = document.querySelectorAll('.slide');
        const dots = document.querySelectorAll('.dot');
        const counter = document.getElementById('slide-counter');
        
        slides.forEach(slide => slide.classList.remove('active'));
        dots.forEach(dot => dot.classList.remove('active'));
        
        if (n >= totalSlides) currentSlide = 0;
        if (n < 0) currentSlide = totalSlides - 1;
        
        slides[currentSlide].classList.add('active');
        dots[currentSlide].classList.add('active');
        counter.textContent = `Slide ${currentSlide + 1} of ${totalSlides}`;
    }
    
    function nextSlide() {
        currentSlide++;
        showSlide(currentSlide);
    }
    
    function previousSlide() {
        currentSlide--;
        showSlide(currentSlide);
    }
    
    function goToSlide(n) {
        currentSlide = n;
        showSlide(currentSlide);
    }
    
    document.addEventListener('keydown', function(e) {
        if (e.key === 'ArrowRight') nextSlide();
        if (e.key === 'ArrowLeft') previousSlide();
    });
    </script>
    """)
    
    html_parts.append('</div>')
    return '\n'.join(html_parts)

def process_pptx(file_path):
    """Process PPTX file and extract content, metadata, and generate enhanced HTML"""
    try:
        prs = Presentation(file_path)
        
        # Extract metadata
        metadata = {
            'slides': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'metadata': {}
        }
        
        # Extract core properties
        if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
            metadata['metadata']['title'] = prs.core_properties.title
        if hasattr(prs.core_properties, 'creator') and prs.core_properties.creator:
            metadata['metadata']['creator'] = prs.core_properties.creator
        if hasattr(prs.core_properties, 'created') and prs.core_properties.created:
            metadata['metadata']['created'] = str(prs.core_properties.created)
        if hasattr(prs.core_properties, 'modified') and prs.core_properties.modified:
            metadata['metadata']['modified'] = str(prs.core_properties.modified)
        
        # Convert slides to enhanced images
        slide_images = convert_pptx_to_enhanced_images(file_path)
        
        # Generate enhanced HTML content
        html_content = generate_enhanced_html_viewer(prs, slide_images)
        
        # Extract text content
        text_content = extract_text_from_pptx(prs)
        
        # Encode original file to base64
        with open(file_path, 'rb') as f:
            file_content = f.read()
            base64_content = base64.b64encode(file_content).decode('utf-8')
        
        return {
            'success': True,
            'preview_type': 'pptx',
            'pptx_html': html_content,
            'preview_content': text_content,
            'pptx_metadata': metadata,
            'pptx_base64': base64_content,
            'slide_images': slide_images,
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': f'PPTX processing failed: {str(e)}',
            'preview_type': 'pptx'
        }

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({
            'success': False,
            'error': 'Usage: python pptx_processor_enhanced.py <file_path>',
            'preview_type': 'pptx'
        }))
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    if not os.path.exists(file_path):
        print(json.dumps({
            'success': False,
            'error': f'File not found: {file_path}',
            'preview_type': 'pptx'
        }))
        sys.exit(1)
    
    result = process_pptx(file_path)
    print(json.dumps(result))
