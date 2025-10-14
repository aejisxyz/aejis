#!/usr/bin/env python3
"""
PPTX Processor for Aejis
Processes PowerPoint presentations and converts them to HTML for preview
"""

import sys
import json
import base64
import os
from pptx import Presentation
from pptx.util import Inches
import xml.etree.ElementTree as ET
from PIL import Image
import io
import subprocess

def convert_pptx_to_images(file_path):
    """Convert PPTX slides to PNG images using matplotlib"""
    try:
        import matplotlib.pyplot as plt
        import matplotlib.patches as patches
        from matplotlib.patches import Rectangle
        import numpy as np
        
        # Load presentation
        prs = Presentation(file_path)
        slide_images = []
        
        for i, slide in enumerate(prs.slides):
            try:
                # Create figure with slide dimensions (16:9 aspect ratio)
                fig, ax = plt.subplots(figsize=(16, 9))
                ax.set_xlim(0, 16)
                ax.set_ylim(0, 9)
                ax.set_aspect('equal')
                ax.axis('off')
                
                # Set background
                ax.add_patch(Rectangle((0, 0), 16, 9, facecolor='white', edgecolor='lightgray', linewidth=2))
                
                # Add slide title
                ax.text(8, 8.5, f'Slide {i + 1}', ha='center', va='center', fontsize=24, fontweight='bold', color='#333333')
                
                # Extract and display text content
                y_position = 7.5
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if len(text) > 50:
                            # Split long text into multiple lines
                            words = text.split()
                            lines = []
                            current_line = ""
                            for word in words:
                                if len(current_line + " " + word) <= 50:
                                    current_line += " " + word if current_line else word
                                else:
                                    lines.append(current_line)
                                    current_line = word
                            if current_line:
                                lines.append(current_line)
                            
                            for line in lines:
                                ax.text(1, y_position, line, ha='left', va='center', fontsize=14, color='#666666', wrap=True)
                                y_position -= 0.4
                        else:
                            ax.text(1, y_position, text, ha='left', va='center', fontsize=16, color='#333333')
                            y_position -= 0.6
                    
                    elif hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                text = paragraph.text.strip()
                                indent = paragraph.level * 0.5
                                ax.text(1 + indent, y_position, text, ha='left', va='center', fontsize=14, color='#666666')
                                y_position -= 0.4
                
                # Add slide number at bottom
                ax.text(15, 0.5, f'{i + 1}', ha='right', va='center', fontsize=12, color='#999999')
                
                # Convert to base64
                buf = io.BytesIO()
                plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                buf.seek(0)
                img_data = buf.getvalue()
                img_base64 = base64.b64encode(img_data).decode('utf-8')
                plt.close(fig)
                
                slide_images.append({
                    'slide_number': i + 1,
                    'image_base64': img_base64,
                    'mime_type': 'image/png'
                })
                
            except Exception as e:
                print(f"Error processing slide {i + 1}: {str(e)}")
                continue
                
        return slide_images
                
    except Exception as e:
        print(f"Error converting PPTX to images: {str(e)}")
        return []

def process_pptx(file_path):
    """Process PPTX file and extract content, metadata, and generate HTML with slide images"""
    try:
        # Load the presentation
        prs = Presentation(file_path)
        
        # Extract basic metadata
        metadata = {
            'slides': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'metadata': {}
        }
        
        # Extract core properties if available
        if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
            metadata['metadata']['title'] = prs.core_properties.title
        if hasattr(prs.core_properties, 'creator') and prs.core_properties.creator:
            metadata['metadata']['creator'] = prs.core_properties.creator
        if hasattr(prs.core_properties, 'created') and prs.core_properties.created:
            metadata['metadata']['created'] = str(prs.core_properties.created)
        if hasattr(prs.core_properties, 'modified') and prs.core_properties.modified:
            metadata['metadata']['modified'] = str(prs.core_properties.modified)
        
        # Convert slides to images
        slide_images = convert_pptx_to_images(file_path)
        
        # Generate HTML content with slide images
        html_content = generate_slide_image_viewer(prs, slide_images)
        
        # Extract text content
        text_content = extract_text_from_pptx(prs)
        
        # Encode original file to base64
        with open(file_path, 'rb') as f:
            file_content = f.read()
            base64_content = base64.b64encode(file_content).decode('utf-8')
        
        return {
            'success': True,
            'preview_type': 'pptx',
            'pptx_html': html_content,
            'preview_content': text_content,
            'pptx_metadata': metadata,
            'pptx_base64': base64_content,
            'slide_images': slide_images,
            'filename': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': f'PPTX processing failed: {str(e)}',
            'preview_type': 'pptx'
        }

def generate_slide_image_viewer(prs, slide_images):
    """Generate HTML viewer with actual slide images"""
    html_parts = []
    html_parts.append('<div class="pptx-slide-viewer">')
    html_parts.append('<div class="slide-viewer-header">')
    html_parts.append('<h2>📊 PowerPoint Presentation</h2>')
    html_parts.append('<div class="slide-controls">')
    html_parts.append('<button onclick="previousSlide()" class="slide-btn">← Previous</button>')
    html_parts.append('<span id="slide-counter">Slide 1 of ' + str(len(slide_images)) + '</span>')
    html_parts.append('<button onclick="nextSlide()" class="slide-btn">Next →</button>')
    html_parts.append('</div>')
    html_parts.append('</div>')
    
    html_parts.append('<div class="slides-container">')
    
    for i, slide_img in enumerate(slide_images):
        slide_class = "slide active" if i == 0 else "slide"
        html_parts.append(f'<div class="{slide_class}" id="slide-{i}">')
        html_parts.append(f'<div class="slide-header">Slide {slide_img["slide_number"]}</div>')
        html_parts.append('<div class="slide-image-container">')
        html_parts.append(f'<img src="data:{slide_img["mime_type"]};base64,{slide_img["image_base64"]}" alt="Slide {slide_img["slide_number"]}" class="slide-image">')
        html_parts.append('</div>')
        html_parts.append('</div>')
    
    html_parts.append('</div>')
    
    # Add slide navigation dots
    html_parts.append('<div class="slide-dots">')
    for i in range(len(slide_images)):
        dot_class = "dot active" if i == 0 else "dot"
        html_parts.append(f'<span class="{dot_class}" onclick="goToSlide({i})"></span>')
    html_parts.append('</div>')
    
    # Add JavaScript for slide navigation
    html_parts.append('''
    <script>
    let currentSlide = 0;
    const totalSlides = ''' + str(len(slide_images)) + ''';
    
    function showSlide(n) {
        const slides = document.querySelectorAll('.slide');
        const dots = document.querySelectorAll('.dot');
        const counter = document.getElementById('slide-counter');
        
        // Hide all slides
        slides.forEach(slide => slide.classList.remove('active'));
        dots.forEach(dot => dot.classList.remove('active'));
        
        // Show current slide
        if (n >= totalSlides) currentSlide = 0;
        if (n < 0) currentSlide = totalSlides - 1;
        
        slides[currentSlide].classList.add('active');
        dots[currentSlide].classList.add('active');
        counter.textContent = `Slide ${currentSlide + 1} of ${totalSlides}`;
    }
    
    function nextSlide() {
        currentSlide++;
        showSlide(currentSlide);
    }
    
    function previousSlide() {
        currentSlide--;
        showSlide(currentSlide);
    }
    
    function goToSlide(n) {
        currentSlide = n;
        showSlide(currentSlide);
    }
    
    // Keyboard navigation
    document.addEventListener('keydown', function(e) {
        if (e.key === 'ArrowRight') nextSlide();
        if (e.key === 'ArrowLeft') previousSlide();
    });
    </script>
    ''')
    
    html_parts.append('</div>')
    return '\n'.join(html_parts)

def generate_html_from_pptx(prs):
    """Generate HTML representation of the PowerPoint presentation with slide viewer"""
    html_parts = []
    html_parts.append('<div class="pptx-slide-viewer">')
    html_parts.append('<div class="slide-viewer-header">')
    html_parts.append('<h2>📊 PowerPoint Presentation</h2>')
    html_parts.append('<div class="slide-controls">')
    html_parts.append('<button onclick="previousSlide()" class="slide-btn">← Previous</button>')
    html_parts.append('<span id="slide-counter">Slide 1 of ' + str(len(prs.slides)) + '</span>')
    html_parts.append('<button onclick="nextSlide()" class="slide-btn">Next →</button>')
    html_parts.append('</div>')
    html_parts.append('</div>')
    
    html_parts.append('<div class="slides-container">')
    
    for i, slide in enumerate(prs.slides):
        slide_class = "slide active" if i == 0 else "slide"
        html_parts.append(f'<div class="{slide_class}" id="slide-{i}">')
        html_parts.append(f'<div class="slide-header">Slide {i + 1}</div>')
        html_parts.append('<div class="slide-content">')
        
        # Process slide content
        slide_content = process_slide_content(slide)
        html_parts.append(slide_content)
        
        html_parts.append('</div>')
        html_parts.append('</div>')
    
    html_parts.append('</div>')
    
    # Add slide navigation dots
    html_parts.append('<div class="slide-dots">')
    for i in range(len(prs.slides)):
        dot_class = "dot active" if i == 0 else "dot"
        html_parts.append(f'<span class="{dot_class}" onclick="goToSlide({i})"></span>')
    html_parts.append('</div>')
    
    # Add JavaScript for slide navigation
    html_parts.append('''
    <script>
    let currentSlide = 0;
    const totalSlides = ''' + str(len(prs.slides)) + ''';
    
    function showSlide(n) {
        const slides = document.querySelectorAll('.slide');
        const dots = document.querySelectorAll('.dot');
        const counter = document.getElementById('slide-counter');
        
        // Hide all slides
        slides.forEach(slide => slide.classList.remove('active'));
        dots.forEach(dot => dot.classList.remove('active'));
        
        // Show current slide
        if (n >= totalSlides) currentSlide = 0;
        if (n < 0) currentSlide = totalSlides - 1;
        
        slides[currentSlide].classList.add('active');
        dots[currentSlide].classList.add('active');
        counter.textContent = `Slide ${currentSlide + 1} of ${totalSlides}`;
    }
    
    function nextSlide() {
        currentSlide++;
        showSlide(currentSlide);
    }
    
    function previousSlide() {
        currentSlide--;
        showSlide(currentSlide);
    }
    
    function goToSlide(n) {
        currentSlide = n;
        showSlide(currentSlide);
    }
    
    // Keyboard navigation
    document.addEventListener('keydown', function(e) {
        if (e.key === 'ArrowRight') nextSlide();
        if (e.key === 'ArrowLeft') previousSlide();
    });
    </script>
    ''')
    
    html_parts.append('</div>')
    return '\n'.join(html_parts)

def process_slide_content(slide):
    """Process individual slide content and convert to HTML"""
    content_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            # Text content
            text = shape.text.strip()
            if text:
                # Check if it's a title (usually larger text or first text element)
                if len(content_parts) == 0 or len(text) < 100:
                    content_parts.append(f'<h4>{escape_html(text)}</h4>')
                else:
                    content_parts.append(f'<p>{escape_html(text)}</p>')
        
        elif hasattr(shape, 'text_frame'):
            # Text frame content
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    if paragraph.level == 0:
                        content_parts.append(f'<h4>{escape_html(text)}</h4>')
                    else:
                        content_parts.append(f'<p style="margin-left: {paragraph.level * 20}px;">{escape_html(text)}</p>')
        
        elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.TABLE
            # Table content
            try:
                if hasattr(shape, 'table') and shape.table:
                    table_html = process_table(shape.table)
                    if table_html:
                        content_parts.append(table_html)
            except AttributeError:
                # Skip if shape doesn't have table attribute
                pass
    
    if not content_parts:
        content_parts.append('<p><em>No text content found in this slide</em></p>')
    
    return '\n'.join(content_parts)

def process_table(table):
    """Process table content and convert to HTML"""
    if not table or not table.rows:
        return ''
    
    html_parts = []
    html_parts.append('<table class="slide-table">')
    
    # Process header row
    if table.rows:
        html_parts.append('<thead><tr>')
        for cell in table.rows[0].cells:
            text = cell.text.strip() if cell.text else ''
            html_parts.append(f'<th>{escape_html(text)}</th>')
        html_parts.append('</tr></thead>')
    
    # Process data rows
    if len(table.rows) > 1:
        html_parts.append('<tbody>')
        for row in table.rows[1:]:
            html_parts.append('<tr>')
            for cell in row.cells:
                text = cell.text.strip() if cell.text else ''
                html_parts.append(f'<td>{escape_html(text)}</td>')
            html_parts.append('</tr>')
        html_parts.append('</tbody>')
    
    html_parts.append('</table>')
    return '\n'.join(html_parts)

def extract_text_from_pptx(prs):
    """Extract all text content from the presentation"""
    text_parts = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f'--- Slide {i} ---')
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text_parts.append(shape.text.strip())
            elif hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
    
    return '\n'.join(text_parts)

def escape_html(text):
    """Escape HTML special characters"""
    if not text:
        return ''
    
    text = text.replace('&', '&amp;')
    text = text.replace('<', '&lt;')
    text = text.replace('>', '&gt;')
    text = text.replace('"', '&quot;')
    text = text.replace("'", '&#x27;')
    return text

def main():
    """Main function to process PPTX file"""
    if len(sys.argv) != 2:
        print(json.dumps({
            'success': False,
            'error': 'Usage: python pptx_processor.py <file_path>'
        }))
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    if not os.path.exists(file_path):
        print(json.dumps({
            'success': False,
            'error': f'File not found: {file_path}'
        }))
        sys.exit(1)
    
    result = process_pptx(file_path)
    print(json.dumps(result))

if __name__ == '__main__':
    main()
