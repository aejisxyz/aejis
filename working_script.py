#!/usr/bin/env python3
import os
import sys
import json
import base64
import struct
import hashlib
import subprocess
from io import BytesIO
from datetime import datetime

# Install required packages
try:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'Pillow', '--quiet'])
except:
    pass  # Continue even if installation fails


def get_file_type(file_path):
    """Detect file type using magic bytes and extension"""
    try:
        import magic
        mime = magic.from_file(file_path, mime=True)
        return mime
    except:
        # Fallback to extension-based detection
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {{
            '.txt': 'text/plain', '.log': 'text/plain', '.md': 'text/markdown',
            '.json': 'application/json', '.xml': 'application/xml', '.csv': 'text/csv',
            '.py': 'text/x-python', '.js': 'application/javascript', '.html': 'text/html',
            '.css': 'text/css', '.sh': 'text/x-shellscript', '.bat': 'application/x-msdos-program',
            '.ps1': 'application/x-powershell', '.sql': 'application/sql',
            '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
            '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp',
            '.svg': 'image/svg+xml', '.tiff': 'image/tiff', '.ico': 'image/x-icon',
            '.pdf': 'application/pdf', '.doc': 'application/msword', '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.ppt': 'application/vnd.ms-powerpoint', '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xls': 'application/vnd.ms-excel', '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.zip': 'application/zip', '.rar': 'application/x-rar-compressed', '.7z': 'application/x-7z-compressed',
            '.tar': 'application/x-tar', '.gz': 'application/gzip', '.bz2': 'application/x-bzip2',
            '.mp4': 'video/mp4', '.avi': 'video/x-msvideo', '.mov': 'video/quicktime',
            '.mkv': 'video/x-matroska', '.wmv': 'video/x-ms-wmv', '.flv': 'video/x-flv',
            '.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.flac': 'audio/flac',
            '.aac': 'audio/aac', '.ogg': 'audio/ogg', '.m4a': 'audio/mp4',
            '.exe': 'application/x-msdownload', '.dll': 'application/x-msdownload',
            '.msi': 'application/x-msi', '.deb': 'application/x-debian-package',
            '.rpm': 'application/x-rpm', '.apk': 'application/vnd.android.package-archive',
            '.ttf': 'font/ttf', '.otf': 'font/otf', '.woff': 'font/woff',
            '.woff2': 'font/woff2', '.eot': 'application/vnd.ms-fontobject',
            '.db': 'application/x-sqlite3', '.sqlite': 'application/x-sqlite3',
            '.mdb': 'application/x-msaccess', '.accdb': 'application/vnd.msaccess'
        }}
        return type_map.get(ext, 'application/octet-stream')

def process_image(file_path, result):
    """Process image files with advanced metadata"""
    try:
        from PIL import Image, ExifTags
        import exifread
        
        # Get image dimensions and properties first
        with Image.open(file_path) as img:
            width, height = img.size
            mode = img.mode
            format_name = img.format or "Unknown"
        
        # Process image for thumbnail and analysis
        with Image.open(file_path) as img:
            
            # Create high-quality thumbnail
            img.thumbnail((800, 800), Image.Resampling.LANCZOS)
            buffer = BytesIO()
            img.save(buffer, format='PNG', optimize=True)
            thumbnail_b64 = base64.b64encode(buffer.getvalue()).decode()
            
            # Extract EXIF data
            exif_data = {{}}
            try:
                with open(file_path, 'rb') as f:
                    tags = exifread.process_file(f, details=False)
                    for tag in tags:
                        exif_data[str(tag)] = str(tags[tag])
            except:
                pass
            
            # Color analysis
            colors = img.getcolors(maxcolors=256*256*256)
            dominant_colors = sorted(colors, key=lambda x: x[0], reverse=True)[:5] if colors else []
            
            result["preview_content"] = f"Image Analysis\nDimensions: {{width}}x{{height}} pixels\nColor Mode: {{mode}}\nFormat: {{format_name}}\n\nEXIF Data: {{len(exif_data)}} tags found"
            result["preview_type"] = "image"
            result["metadata"] = {{
                "dimensions": f"{{width}}x{{height}}",
                "mode": mode,
                "format": format_name,
                "exif_data": exif_data,
                "dominant_colors": dominant_colors[:3],
                "file_size": result["file_size"]
            }}
            result["thumbnail"] = f"data:image/png;base64,{{thumbnail_b64}}"
            
    except ImportError:
        # Fallback when PIL is not available
        result["preview_content"] = f"Image File ({{result['file_size']}} bytes)\nFormat: {{get_file_type(file_path)}}\nBasic analysis only - PIL not available in container"
        result["preview_type"] = "image"
        result["metadata"] = {{
            "file_size": result["file_size"],
            "mime_type": get_file_type(file_path),
            "pil_available": False
        }}
        result["thumbnail"] = None
            
    except Exception as e:
        result["preview_content"] = f"Image processing error: {{str(e)}}"
        result["preview_type"] = "error"
        
def process_video(file_path, result):
    """Process video files with metadata extraction"""
    try:
        from mutagen import File as MutagenFile
        
        video_file = MutagenFile(file_path)
        duration = getattr(video_file.info, 'length', 0)
        bitrate = getattr(video_file.info, 'bitrate', 0)
        codec = getattr(video_file.info, 'codec', 'Unknown')
        
        # Initialize video variables
        width, height, fps, frame_count = 0, 0, 0, 0
        duration = 0
        
        # Try to get video dimensions
        try:
            import cv2
            cap = cv2.VideoCapture(file_path)
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = frame_count / fps if fps > 0 else 0
            cap.release()
        except:
            pass
        
        # Create thumbnail (first frame)
        thumbnail_b64 = None
        try:
            import cv2
            cap = cv2.CAP_PROP_POS_FRAMES
            cap = cv2.VideoCapture(file_path)
            ret, frame = cap.read()
            if ret:
                from PIL import Image
                img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                img.thumbnail((400, 300), Image.Resampling.LANCZOS)
                buffer = BytesIO()
                img.save(buffer, format='PNG')
                thumbnail_b64 = base64.b64encode(buffer.getvalue()).decode()
            cap.release()
        except:
            pass
                    
        result["preview_content"] = f"Video Analysis\nDuration: {{duration:.2f}}s\nResolution: {{width}}x{{height}}\nFPS: {{fps:.2f}}\nCodec: {{codec}}\nBitrate: {{bitrate}} bps"
        result["preview_type"] = "video"
        result["metadata"] = {{
            "duration": duration,
            "dimensions": f"{{width}}x{{height}}",
            "fps": fps,
            "frame_count": frame_count,
            "codec": codec,
            "bitrate": bitrate,
            "file_size": result["file_size"]
        }}
        if thumbnail_b64:
            result["thumbnail"] = f"data:image/png;base64,{{thumbnail_b64}}"
                    
    except Exception as e:
        result["preview_content"] = f"Video processing error: {{str(e)}}"
        result["preview_type"] = "error"
        
def process_audio(file_path, result):
    """Process audio files with metadata extraction"""
    try:
        from mutagen import File as MutagenFile
        
        audio_file = MutagenFile(file_path)
        duration = getattr(audio_file.info, 'length', 0)
        bitrate = getattr(audio_file.info, 'bitrate', 0)
        sample_rate = getattr(audio_file.info, 'sample_rate', 0)
        channels = getattr(audio_file.info, 'channels', 0)
        codec = getattr(audio_file.info, 'codec', 'Unknown')
        
        # Extract tags
        tags = {{}}
        if hasattr(audio_file, 'tags') and audio_file.tags:
            for key, value in audio_file.tags.items():
                tags[str(key)] = str(value)
        
        result["preview_content"] = f"Audio Analysis\nDuration: {{duration:.2f}}s\nBitrate: {{bitrate}} bps\nSample Rate: {{sample_rate}} Hz\nChannels: {{channels}}\nCodec: {{codec}}\n\nTags: {{len(tags)}} metadata entries"
        result["preview_type"] = "audio"
        result["metadata"] = {{
            "duration": duration,
            "bitrate": bitrate,
            "sample_rate": sample_rate,
            "channels": channels,
            "codec": codec,
            "tags": tags,
            "file_size": result["file_size"]
        }}
        
    except Exception as e:
        result["preview_content"] = f"Audio processing error: {{str(e)}}"
        result["preview_type"] = "error"

def process_pdf(file_path, result):
    """Process PDF files with secure Docker processing"""
    try:
        import base64
        
        # Read the raw PDF file for viewing
        with open(file_path, 'rb') as file:
            pdf_data = file.read()
            pdf_base64 = base64.b64encode(pdf_data).decode('utf-8')
        
        # Extract metadata
        text_content = ""
        metadata = {{}}
        page_count = 0
        
        try:
            import PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)
                    
                # Extract metadata
                if pdf_reader.metadata:
                    metadata = {{
                        'title': str(pdf_reader.metadata.get('/Title', '')),
                        'author': str(pdf_reader.metadata.get('/Author', '')),
                        'creator': str(pdf_reader.metadata.get('/Creator', '')),
                        'producer': str(pdf_reader.metadata.get('/Producer', '')),
                        'page_count': page_count
                    }}
                
                # Extract text from first 3 pages for metadata
                max_pages = min(3, page_count)
                for page_num in range(max_pages):
                    try:
                        page = pdf_reader.pages[page_num]
                        text_content += f"\n--- Page {{page_num + 1}} ---\n"
                        text_content += page.extract_text()
                    except:
                        text_content += f"\n--- Page {{page_num + 1}} (text extraction failed) ---\n"
        except Exception as e:
            text_content = f"Text extraction failed: {{str(e)}}"
            
        result.update({{
            "preview_type": "pdf",
            "preview_content": text_content[:5000] or "No text content found",
            "pdf_base64": pdf_base64,
            "base64_preview": pdf_base64,
            "pdf_metadata": {{
                    "page_count": page_count,
                "metadata": metadata,
                "extracted_pages": min(3, page_count) if page_count > 0 else 0
            }}
        }})
    except Exception as e:
        result.update({{
            "preview_type": "pdf",
            "preview_content": f"PDF processing error: {{str(e)}}"
        }})

def process_docx(file_path, result):
    """Process DOCX files with secure Docker processing"""
    print("DEBUG: process_docx function called")
    try:
        import base64
        print("DEBUG: base64 imported successfully")
        
        # Read the raw DOCX file for viewing
        with open(file_path, 'rb') as file:
            docx_data = file.read()
            docx_base64 = base64.b64encode(docx_data).decode('utf-8')
        
        # Extract content and metadata
        text_content = ""
        metadata = {{}}
        paragraph_count = 0
        table_count = 0
        
        try:
            from docx import Document
            doc = Document(file_path)
            
            # Extract paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            paragraph_count = len(paragraphs)
            
            # Extract tables
            table_count = len(doc.tables)
            
            # Get first 5 paragraphs for preview
            text_content = '\n'.join(paragraphs[:5])
            
            # Extract document properties
            core_props = doc.core_properties
            metadata = {{
                'title': str(core_props.title or ''),
                'author': str(core_props.author or ''),
                'subject': str(core_props.subject or ''),
                'created': str(core_props.created or ''),
                'modified': str(core_props.modified or ''),
                'paragraphs': paragraph_count,
                'tables': table_count
            }}
            
        except Exception as e:
            text_content = f"Content extraction failed: {{str(e)}}"
            metadata = {{'error': {{str(e)}}}}
            
        result.update({{
            "preview_type": "docx",
            "preview_content": text_content or "No text content found",
            "docx_base64": docx_base64,
            "base64_preview": docx_base64,
            "docx_metadata": {{
                "paragraphs": paragraph_count,
                "tables": table_count,
                "metadata": metadata
            }}
        }})
    except Exception as e:
        result.update({{
            "preview_type": "docx",
            "preview_content": f"DOCX processing error: {{str(e)}}"
        }})

def process_document(file_path, result):
    """Process document files (DOC, PPT, XLS, etc.)"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        
        # Skip PDF and DOCX processing here since they have dedicated functions
        if ext == '.pdf':
            result.update({{
                "preview_type": "pdf",
                "preview_content": "PDF processing handled by dedicated function"
            }})
            return
            
        if ext == '.docx':
            result.update({{
                "preview_type": "docx",
                "preview_content": "DOCX processing handled by dedicated function"
            }})
            return
            
        if ext in ['.doc']:
            try:
                from docx import Document
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                tables = len(doc.tables)
                
                content_preview = ' '.join(paragraphs[:5])
                result["preview_content"] = f"Word Document\nParagraphs: {{len(paragraphs)}}\nTables: {{tables}}\n\nContent Preview:\n{{content_preview}}..."
                result["preview_type"] = "document"
                result["metadata"] = {{
                    "paragraphs": len(paragraphs),
                    "tables": tables,
                    "file_size": result["file_size"]
                }}
            except:
                result["preview_content"] = f"Word Document ({{result['file_size']}} bytes)\nContent extraction not available"
                result["preview_type"] = "document"
        
        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides = len(prs.slides)
                
                slide_texts = []
                for slide in prs.slides[:3]:  # First 3 slides
                    text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
                    slide_texts.append(text.strip())
                
                slide_previews = '\n'.join(slide_texts)
                result["preview_content"] = f"PowerPoint Presentation\nSlides: {{slides}}\n\nSlide Previews:\n{{slide_previews}}"
                result["preview_type"] = "presentation"
                result["metadata"] = {{
                    "slides": slides,
                    "file_size": result["file_size"]
                }}
            except:
                file_size = result['file_size']
                result["preview_content"] = f"PowerPoint Presentation ({{file_size}} bytes)\nContent extraction not available"
                result["preview_type"] = "presentation"
        
        elif ext in ['.xlsx', '.xls']:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, read_only=True)
                sheets = wb.sheetnames
                
                # Get first sheet data
                ws = wb[sheets[0]]
                rows = list(ws.iter_rows(values_only=True))[:10]  # First 10 rows
                
                sheet_names = ', '.join(sheets)
                result["preview_content"] = f"Excel Spreadsheet\nSheets: {{len(sheets)}}\nSheet Names: {{sheet_names}}\n\nFirst Sheet Data:\n{{str(rows[:5])}}"
                result["preview_type"] = "spreadsheet"
                result["metadata"] = {{
                    "sheets": sheets,
                    "sheet_count": len(sheets),
                    "file_size": result["file_size"]
                }}
            except:
                result["preview_content"] = f"Excel Spreadsheet ({{result['file_size']}} bytes)\nContent extraction not available"
                result["preview_type"] = "spreadsheet"
                    
    except Exception as e:
        result["preview_content"] = f"Document processing error: {{str(e)}}"
        result["preview_type"] = "error"
        
def process_archive(file_path, result):
    """Process archive files with detailed content analysis"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
                
        if ext == '.zip':
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as zf:
                file_list = zf.namelist()
                total_size = sum(info.file_size for info in zf.infolist())
                compressed_size = sum(info.compress_size for info in zf.infolist())
                compression_ratio = (1 - compressed_size / total_size) * 100 if total_size > 0 else 0
                
                # Analyze file types in archive
                file_types = {{}}
                for filename in file_list:
                    ext = os.path.splitext(filename)[1].lower()
                    file_types[ext] = file_types.get(ext, 0) + 1
                
                result["preview_content"] = f"ZIP Archive\nFiles: {{len(file_list)}}\nTotal Size: {{total_size}} bytes\nCompressed: {{compressed_size}} bytes\nCompression: {{compression_ratio:.1f}}%\n\nFile Types: {{dict(list(file_types.items())[:5])}}"
                result["preview_type"] = "archive"
                result["metadata"] = {{
                    "file_count": len(file_list),
                    "total_size": total_size,
                    "compressed_size": compressed_size,
                    "compression_ratio": compression_ratio,
                    "file_types": file_types,
                    "files": file_list[:20]  # First 20 files
                }}
        
        elif ext == '.rar':
            try:
                import rarfile
                with rarfile.RarFile(file_path, 'r') as rf:
                    file_list = rf.namelist()
                    total_size = sum(info.file_size for info in rf.infolist())
                    
                    result["preview_content"] = f"RAR Archive\nFiles: {{len(file_list)}}\nTotal Size: {{total_size}} bytes\n\nFiles: {{file_list[:10]}}"
                    result["preview_type"] = "archive"
                    result["metadata"] = {{
                        "file_count": len(file_list),
                        "total_size": total_size,
                        "files": file_list[:20]
                    }}
            except:
                result["preview_content"] = f"RAR Archive ({{result['file_size']}} bytes)\nRAR processing not available"
                result["preview_type"] = "archive"
        
        elif ext == '.7z':
            try:
                import py7zr
                with py7zr.SevenZipFile(file_path, 'r') as zf:
                    file_list = zf.getnames()
                    total_size = sum(info.uncompressed for info in zf.list() if hasattr(info, 'uncompressed'))
                    
                    result["preview_content"] = f"7Z Archive\nFiles: {{len(file_list)}}\nTotal Size: {{total_size}} bytes\n\nFiles: {{file_list[:10]}}"
                    result["preview_type"] = "archive"
                    result["metadata"] = {{
                        "file_count": len(file_list),
                        "total_size": total_size,
                        "files": file_list[:20]
                    }}
            except:
                result["preview_content"] = f"7Z Archive ({{result['file_size']}} bytes)\n7Z processing not available"
                result["preview_type"] = "archive"
        
    except Exception as e:
        result["preview_content"] = f"Archive processing error: {{str(e)}}"
        result["preview_type"] = "error"
        
def process_executable(file_path, result):
    """Process executable files with security analysis"""
    try:
        with open(file_path, 'rb') as f:
            header = f.read(1024)
        
        # Analyze executable headers
        if header[:2] == b'MZ':
            pe_offset = struct.unpack('<I', header[60:64])[0]
            if pe_offset < len(header):
                pe_header = header[pe_offset:pe_offset+4]
                if pe_header == b'PE\x00':
                    result["preview_content"] = f"Windows Executable (PE)\nFile Size: {{result['file_size']}} bytes\nArchitecture: PE32/PE32+\n\nWARNING: This is an executable file. Only run if you trust the source."
                    result["preview_type"] = "executable"
                    result["metadata"] = {{
                        "platform": "Windows",
                        "type": "PE",
                        "file_size": result["file_size"],
                        "risk_level": "High"
                    }}
            else:
                result["preview_content"] = f"Executable File\nFile Size: {{result['file_size']}} bytes\nType: Unknown executable format"
                result["preview_type"] = "executable"
        elif header[:4] == b'\x7fELF':
            result["preview_content"] = f"Linux Executable (ELF)\nFile Size: {{result['file_size']}} bytes\nArchitecture: ELF\n\nWARNING: This is an executable file. Only run if you trust the source."
            result["preview_type"] = "executable"
            result["metadata"] = {{
                "platform": "Linux",
                "type": "ELF",
                "file_size": result["file_size"],
                "risk_level": "High"
            }}
        else:
            result["preview_content"] = f"Binary Executable\nFile Size: {{result['file_size']}} bytes\nType: Unknown binary format\n\nWARNING: This appears to be an executable file."
            result["preview_type"] = "executable"
            result["metadata"] = {{
                "platform": "Unknown",
                "type": "Binary",
                "file_size": result["file_size"],
                "risk_level": "High"
            }}
        
    except Exception as e:
        result["preview_content"] = f"Executable processing error: {{str(e)}}"
        result["preview_type"] = "error"

def process_font(file_path, result):
    """Process font files with metadata extraction"""
    try:
        from fontTools import TTFont
        font = TTFont(file_path)
        
        # Extract font metadata
        name_table = font['name']
        font_name = ""
        font_family = ""
        for record in name_table.names:
            if record.nameID == 1:  # Font Family
                font_family = record.string.decode('utf-16-be') if record.platformID == 3 else record.string.decode('utf-8')
            elif record.nameID == 4:  # Full Font Name
                font_name = record.string.decode('utf-16-be') if record.platformID == 3 else record.string.decode('utf-8')
        
        # Get font metrics
        head_table = font['head']
        units_per_em = head_table.unitsPerEm
        created = datetime.fromtimestamp(head_table.created)
        modified = datetime.fromtimestamp(head_table.modified)
        
        result["preview_content"] = f"Font File\nName: {{font_name}}\nFamily: {{font_family}}\nUnits per EM: {{units_per_em}}\nCreated: {{created.strftime('%Y-%m-%d %H:%M:%S')}}\nModified: {{modified.strftime('%Y-%m-%d %H:%M:%S')}}"
        result["preview_type"] = "font"
        result["metadata"] = {{
            "font_name": font_name,
            "font_family": font_family,
            "units_per_em": units_per_em,
            "created": created.isoformat(),
            "modified": modified.isoformat(),
            "file_size": result["file_size"]
        }}
        
    except Exception as e:
        result["preview_content"] = f"Font processing error: {{str(e)}}"
        result["preview_type"] = "error"

def process_database(file_path, result):
    """Process database files"""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.db', '.sqlite', '.sqlite3']:
            import sqlite3
            conn = sqlite3.connect(file_path)
            cursor = conn.cursor()
            
            # Get table information
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            
            # Get table counts
            table_info = {{}}
            for table in tables:
                table_name = table[0]
                cursor.execute(f"SELECT COUNT(*) FROM {{table_name}}")
                count = cursor.fetchone()[0]
                table_info[table_name] = count
            
            conn.close()
            
            result["preview_content"] = f"SQLite Database\nTables: {{len(tables)}}\nTable Names: {{[t[0] for t in tables]}}\n\nTable Row Counts: {{table_info}}"
            result["preview_type"] = "database"
            result["metadata"] = {{
                "tables": [t[0] for t in tables],
                "table_count": len(tables),
                "table_info": table_info,
                "file_size": result["file_size"]
            }}
        
    except Exception as e:
        result["preview_content"] = f"Database processing error: {{str(e)}}"
        result["preview_type"] = "error"

def safe_process_file():
    """Safely process file in isolated container with universal support"""
    try:
        file_path = "/tmp/file.docx"
        
        # Debug: List files in /tmp directory (for troubleshooting)
        # print(f"DEBUG: Files in /tmp: {{os.listdir('/tmp')}}")
        # print(f"DEBUG: Looking for file: {{file_path}}")
        # print(f"DEBUG: File exists: {{os.path.exists(file_path)}}")
        
        if not os.path.exists(file_path):
            print(json.dumps({{"error": "File not found", "success": False}}))
            return
        
        file_size = os.path.getsize(file_path)
        
        # File size check (max 200MB for preview)
        if file_size > 200 * 1024 * 1024:
            print(json.dumps({{
                "error": f"File too large: {{file_size}} bytes",
                "file_size": file_size,
                "success": False
            }}))
            return
        
        # Set UTF-8 encoding for stdout
        import sys
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
        
        # Install required packages (with error handling to ensure script continues)
        def install_packages():
            """Install required packages for file processing"""
            import subprocess
            import sys
            
            # Essential packages for file processing
            critical_packages = ["Pillow", "PyPDF2", "python-docx"]
            optional_packages = ["python-magic", "mutagen", "openpyxl", 
                                "python-pptx", "python-docx2txt", "rarfile", 
                                "py7zr", "python-magic-bin", "chardet", "exifread", "pydub"]
            
            # Install critical packages first (silently)
            for pkg in critical_packages:
                try:
                    # Try installing with pre-built wheels and no compilation
                    result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--only-binary=all', '--no-compile', '--quiet', pkg], 
                                          capture_output=True, check=True, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                except subprocess.CalledProcessError:
                    try:
                        result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--system', '--quiet', pkg], 
                                              capture_output=True, check=True, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                    except subprocess.CalledProcessError:
                        try:
                            result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--quiet', pkg], 
                                                  capture_output=True, check=True, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                        except:
                            pass  # Silent fail
                except Exception:
                    pass  # Silent fail
            
            # Install optional packages (silently)
            for pkg in optional_packages:
                try:
                    result = subprocess.run([sys.executable, '-m', 'pip', 'install', '--quiet', pkg], 
                                          capture_output=True, timeout=30, text=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                except:
                    pass  # Silent fail
        
        # Install packages with error handling (silently)
        try:
            install_packages()
        except Exception as e:
            pass  # Silent fail, continue with basic processing
        
        # Initialize result with basic info
        result = {{
            "file_size": file_size,
            "file_ext": ".docx",
            "success": True,
            "preview_content": None,
            "preview_type": "unknown",
            "metadata": {{}},
            "thumbnail": None,
            "mime_type": get_file_type(file_path)
        }}
        
        # Determine file type and process accordingly
        mime_type = result["mime_type"]
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Handle DOCX files even if MIME type detection fails
        if file_ext == '.docx' or mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
            print("DEBUG: Routing to process_docx function")
            process_docx(file_path, result)
            print("DEBUG: process_docx completed successfully")
        elif mime_type.startswith('image/'):
            print("DEBUG: Routing to process_image function")
            process_image(file_path, result)
        elif mime_type.startswith('video/'):
            process_video(file_path, result)
        elif mime_type.startswith('audio/'):
            process_audio(file_path, result)
        elif mime_type == 'application/pdf' or file_ext == '.pdf':
            process_pdf(file_path, result)
        elif mime_type in ['application/msword', 'application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                          'application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            process_document(file_path, result)
        elif mime_type in ['application/zip', 'application/x-rar-compressed', 'application/x-7z-compressed',
                          'application/x-tar', 'application/gzip', 'application/x-bzip2']:
            process_archive(file_path, result)
        elif mime_type in ['application/x-msdownload', 'application/x-msi', 'application/x-debian-package',
                          'application/x-rpm', 'application/vnd.android.package-archive']:
            process_executable(file_path, result)
        elif mime_type.startswith('font/') or mime_type in ['application/font-woff', 'application/font-woff2']:
            process_font(file_path, result)
        elif mime_type in ['application/x-sqlite3', 'application/vnd.msaccess']:
            process_database(file_path, result)
        elif mime_type.startswith('text/'):
            # Text files with multiple encoding attempts
            try:
                content = None
                encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read(100000)  # Max 100KB preview
                            if len(content) >= 100000:
                                content += "\n\n[... content truncated for security ...]"
                            break
                    except UnicodeDecodeError:
                        continue
                
                if content is None:
                    # If all encodings fail, read as binary and decode with errors='replace'
                    with open(file_path, 'rb') as f:
                        raw_content = f.read(100000)
                        content = raw_content.decode('utf-8', errors='replace')
                        if len(content) >= 100000:
                            content += "\n\n[... content truncated for security ...]"
                
                result["preview_content"] = f"Text File\nSize: {{file_size}} bytes\n\nContent:\n{{content}}"
                result["preview_type"] = "text"
                result["metadata"] = {{
                    "line_count": len(content.splitlines()),
                    "char_count": len(content),
                    "file_size": file_size
                }}
            except Exception as e:
                result["preview_content"] = f"Text processing error: {{str(e)}}"
                result["preview_type"] = "error"
        else:
            print("DEBUG: FALLBACK - File fell through to binary processing!")
            # Binary/unknown files - enhanced analysis
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(512).hex()  # First 512 bytes as hex
                
                # Calculate file entropy
                f.seek(0)
                data = f.read(min(8192, file_size))  # First 8KB
                
                if data:
                    byte_counts = [0] * 256
                    for byte in data:
                        byte_counts[byte] += 1
                    
                    entropy = 0
                    for count in byte_counts:
                        if count > 0:
                            p = count / len(data)
                            entropy -= p * math.log2(p)
                    
                    result["preview_content"] = f"Binary File Analysis\nSize: {{file_size}} bytes\nEntropy: {{entropy:.2f}}/8.0\nMIME Type: {{mime_type}}\n\nHex Header (first 512 bytes):\n{{header[:200]}}..."
                    result["preview_type"] = "binary"
                    result["metadata"] = {{
                        "header": header,
                        "entropy": entropy,
                        "mime_type": mime_type,
                        "file_size": file_size,
                        "is_likely_encrypted": entropy > 7.0,
                        "is_likely_compressed": entropy > 6.0
                    }}
                else:
                    result["preview_content"] = f"Empty Binary File\nSize: {{file_size}} bytes"
                    result["preview_type"] = "binary"
                
            except Exception as e:
                result["preview_content"] = f"Binary processing error: {{str(e)}}"
                result["preview_type"] = "error"
        
        # Ensure JSON output is properly encoded
        try:
            json_output = json.dumps(result, ensure_ascii=False, indent=2)
            print(json_output)
        except UnicodeEncodeError:
            # Fallback: ensure ASCII encoding
            json_output = json.dumps(result, ensure_ascii=True, indent=2)
            print(json_output)
        
    except Exception as e:
        try:
            error_result = {{"error": str(e), "success": False}}
            json_output = json.dumps(error_result, ensure_ascii=False, indent=2)
            print(json_output)
        except UnicodeEncodeError:
            error_result = {"error": "Processing error occurred", "success": False}
            json_output = json.dumps(error_result, ensure_ascii=True, indent=2)
            print(json_output)

if __name__ == "__main__":
    safe_process_file()
